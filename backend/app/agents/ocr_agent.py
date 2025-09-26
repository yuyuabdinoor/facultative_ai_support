"""
OCR Processing Agent for Facultative Reinsurance System

This module provides OCR processing capabilities for various document types:
- PDF text extraction (text-based and scanned)
- Email (.msg) parsing with attachment handling
- Excel/CSV file processing
- PowerPoint presentation processing
- Word document processing
- Scanned document OCR using PaddleOCR
- Text region detection and extraction

Handles reinsurance-specific documents:
- Property survey reports
- Risk placement slips/contract wording
- RI Slips
- Inspection reports
- Additional supporting documents
"""

from __future__ import annotations

import logging
import os
import os
import tempfile
import zipfile
from pathlib import Path
from typing import Dict, List, Optional, Union, Any
import io
import base64

# Document processing imports
import pandas as pd
import numpy as np
from PIL import Image
try:
    import fitz  # PyMuPDF for better PDF handling
    PYMUPDF_AVAILABLE = True
except ImportError:
    fitz = None
    PYMUPDF_AVAILABLE = False

import pdfplumber

# PaddleOCR for OCR processing
try:
    from paddleocr import PaddleOCR
    PADDLEOCR_AVAILABLE = True
except ImportError:
    PaddleOCR = None
    PADDLEOCR_AVAILABLE = False

# Document format handlers
import extract_msg
from openpyxl import load_workbook

# Word document processing
try:
    from docx import Document as DocxDocument
    DOCX_AVAILABLE = True
except ImportError:
    DocxDocument = None
    DOCX_AVAILABLE = False

# PowerPoint processing
try:
    from pptx import Presentation
    PPTX_AVAILABLE = True
except ImportError:
    Presentation = None
    PPTX_AVAILABLE = False

# Pydantic models for structured responses
from pydantic import BaseModel, Field
from datetime import datetime

# Configure logging
logger = logging.getLogger(__name__)


class TextRegion(BaseModel):
    """Represents a detected text region in a document"""
    text: str
    confidence: float
    bbox: List[float] = Field(description="Bounding box coordinates [x1, y1, x2, y2]")
    page_number: Optional[int] = None


class OCRResult(BaseModel):
    """Result of OCR processing"""
    text: str
    confidence: float
    regions: List[TextRegion] = []
    metadata: Dict[str, Any] = {}
    processing_time: float
    success: bool = True
    error_message: Optional[str] = None


class ExcelData(BaseModel):
    """Excel/CSV file processing result"""
    sheets: Dict[str, List[Dict[str, Any]]] = {}
    metadata: Dict[str, Any] = {}
    total_rows: int = 0
    total_sheets: int = 0


class WordDocumentData(BaseModel):
    """Word document processing result"""
    text: str
    paragraphs: List[str] = []
    tables: List[List[List[str]]] = []
    metadata: Dict[str, Any] = {}


class PowerPointData(BaseModel):
    """PowerPoint presentation processing result"""
    slides: List[Dict[str, Any]] = []
    text: str
    metadata: Dict[str, Any] = {}
    total_slides: int = 0


class AttachmentData(BaseModel):
    """Email attachment processing result"""
    filename: str
    content_type: str
    size: int
    processed_content: Optional[Union[OCRResult, ExcelData, WordDocumentData, PowerPointData]] = None
    extraction_success: bool = True
    error_message: Optional[str] = None


class EmailContent(BaseModel):
    """Parsed email content structure"""
    subject: str
    sender: str
    recipients: List[str] = []
    body: str
    attachments: List[AttachmentData] = []
    date: Optional[datetime] = None
    metadata: Dict[str, Any] = {}
    document_type: Optional[str] = None  # e.g., "property_survey", "ri_slip", etc.


class OCRProcessingAgent:
    """
    OCR Processing Agent that handles various document types for reinsurance applications
    
    Supports:
    - Property survey reports
    - Risk placement slips/contract wording
    - RI Slips
    - Inspection reports
    - Additional supporting documents
    """
    
    def __init__(self):
        """Initialize the OCR agent with required models and configurations"""
        self.ocr_model = None
        # Ensure PaddleOCR cache directory is a string
        self.hf_cache_dir = os.environ.get("HF_HOME", "/app/.cache/huggingface")
        self._initialize_ocr_model()
        
        # Document type classification keywords
        self.document_type_keywords = {
            'property_survey': ['property survey', 'survey report', 'property inspection', 'building survey'],
            'ri_slip': ['ri slip', 'reinsurance slip', 'facultative slip', 'treaty slip'],
            'risk_placement': ['placement slip', 'contract wording', 'policy wording', 'terms and conditions'],
            'inspection_report': ['inspection report', 'loss inspection', 'claim inspection', 'site inspection'],
            'supporting_document': ['certificate', 'valuation', 'financial statement', 'loss history']
        }
        
    def _initialize_ocr_model(self):
        """Initialize PaddleOCR model"""
        try:
            if PADDLEOCR_AVAILABLE:
                # Initialize PaddleOCR with English language support
                self.ocr_model = PaddleOCR(use_angle_cls=True, lang='en', show_log=False)
                logger.info("PaddleOCR model initialized successfully")
            else:
                self.ocr_model = None
                logger.warning("PaddleOCR not available, OCR functionality will be limited")
        except Exception as e:
            logger.error(f"Failed to initialize OCR model: {str(e)}")
            self.ocr_model = None
    
    def process_pdf(self, file_path: str) -> OCRResult:
        """
        Extract text from PDF files using pdfplumber for text-based PDFs
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            OCRResult with extracted text and metadata
        """
        start_time = datetime.now()
        
        try:
            text_content = []
            regions = []
            total_confidence = 0
            page_count = 0
            
            with pdfplumber.open(file_path) as pdf:
                for page_num, page in enumerate(pdf.pages):
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(page_text)
                        
                        # Create text regions for each line
                        lines = page_text.split('\n')
                        for i, line in enumerate(lines):
                            if line.strip():
                                regions.append(TextRegion(
                                    text=line.strip(),
                                    confidence=0.95,  # High confidence for text-based PDFs
                                    bbox=[0, i * 20, page.width, (i + 1) * 20],
                                    page_number=page_num + 1
                                ))
                        
                        total_confidence += 0.95
                        page_count += 1
            
            full_text = '\n'.join(text_content)
            avg_confidence = total_confidence / page_count if page_count > 0 else 0
            
            processing_time = (datetime.now() - start_time).total_seconds()
            
            return OCRResult(
                text=full_text,
                confidence=avg_confidence,
                regions=regions,
                metadata={
                    'file_type': 'pdf',
                    'page_count': page_count,
                    'extraction_method': 'pdfplumber'
                },
                processing_time=processing_time,
                success=True
            )
            
        except Exception as e:
            logger.error(f"Error processing PDF {file_path}: {str(e)}")
            processing_time = (datetime.now() - start_time).total_seconds()
            
            return OCRResult(
                text="",
                confidence=0.0,
                processing_time=processing_time,
                success=False,
                error_message=str(e)
            )
    
    def process_scanned_pdf(self, file_path: str) -> OCRResult:
        """
        Process scanned PDF using PaddleOCR
        
        Args:
            file_path: Path to the scanned PDF file
            
        Returns:
            OCRResult with OCR-extracted text and regions
        """
        start_time = datetime.now()
        
        try:
            if not self.ocr_model or not PADDLEOCR_AVAILABLE:
                raise Exception("OCR model not initialized or PaddleOCR not available")
            
            text_content = []
            regions = []
            total_confidence = 0
            word_count = 0
            
            # Convert PDF pages to images using PyMuPDF if available, otherwise use fitz
            if PYMUPDF_AVAILABLE:
                pdf_document = fitz.open(file_path)
                
                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    # Convert page to image
                    mat = fitz.Matrix(2.0, 2.0)  # 2x zoom for better OCR
                    pix = page.get_pixmap(matrix=mat)
                    img_data = pix.tobytes("png")
                    
                    # Convert to PIL Image
                    image = Image.open(io.BytesIO(img_data))
                    
                    # Perform OCR on the image
                    ocr_result = self.ocr_model.ocr(np.array(image), cls=True)
                    
                    page_text = []
                    
                    if ocr_result and ocr_result[0]:
                        for line in ocr_result[0]:
                            if line and len(line) >= 2:
                                bbox_coords = line[0]
                                text_info = line[1]
                                
                                if text_info and len(text_info) >= 2:
                                    text = text_info[0]
                                    confidence = text_info[1]
                                    
                                    page_text.append(text)
                                    total_confidence += confidence
                                    word_count += 1
                                    
                                    # Create text region
                                    regions.append(TextRegion(
                                        text=text,
                                        confidence=confidence,
                                        bbox=[
                                            min([p[0] for p in bbox_coords]),
                                            min([p[1] for p in bbox_coords]),
                                            max([p[0] for p in bbox_coords]),
                                            max([p[1] for p in bbox_coords])
                                        ],
                                        page_number=page_num + 1
                                    ))
                    
                    if page_text:
                        text_content.append(' '.join(page_text))
                
                pdf_document.close()
                
            else:
                # Fallback: try to process as image directly
                ocr_result = self.ocr_model.ocr(file_path, cls=True)
                
                if ocr_result and ocr_result[0]:
                    page_text = []
                    for line in ocr_result[0]:
                        if line and len(line) >= 2:
                            text_info = line[1]
                            if text_info and len(text_info) >= 2:
                                text = text_info[0]
                                confidence = text_info[1]
                                page_text.append(text)
                                total_confidence += confidence
                                word_count += 1
                    
                    if page_text:
                        text_content.append(' '.join(page_text))
            
            full_text = '\n\n'.join(text_content)
            avg_confidence = total_confidence / word_count if word_count > 0 else 0
            
            processing_time = (datetime.now() - start_time).total_seconds()
            
            return OCRResult(
                text=full_text,
                confidence=avg_confidence,
                regions=regions,
                metadata={
                    'file_type': 'scanned_pdf',
                    'page_count': len(text_content),
                    'word_count': word_count,
                    'extraction_method': 'paddleocr'
                },
                processing_time=processing_time,
                success=True
            )
            
        except Exception as e:
            logger.error(f"Error processing scanned PDF {file_path}: {str(e)}")
            processing_time = (datetime.now() - start_time).total_seconds()
            
            return OCRResult(
                text="",
                confidence=0.0,
                processing_time=processing_time,
                success=False,
                error_message=str(e)
            )
    
    def _classify_document_type(self, text: str, filename: str = "") -> str:
        """
        Classify document type based on content and filename
        
        Args:
            text: Document text content
            filename: Document filename
            
        Returns:
            Document type classification
        """
        text_lower = text.lower()
        filename_lower = filename.lower()
        
        for doc_type, keywords in self.document_type_keywords.items():
            for keyword in keywords:
                if keyword in text_lower or keyword in filename_lower:
                    return doc_type
        
        return 'supporting_document'  # Default classification
    
    def _process_attachment(self, attachment, temp_dir: str) -> AttachmentData:
        """
        Process a single email attachment
        
        Args:
            attachment: Email attachment object
            temp_dir: Temporary directory for saving attachment
            
        Returns:
            AttachmentData with processed content
        """
        filename = getattr(attachment, 'longFilename', '') or getattr(attachment, 'shortFilename', '')
        size = getattr(attachment, 'size', 0)
        
        try:
            # Save attachment to temporary file
            temp_file_path = os.path.join(temp_dir, filename)
            attachment.save(temp_file_path)
            
            # Determine file type and process accordingly
            file_extension = Path(filename).suffix.lower()
            processed_content = None
            
            if file_extension == '.pdf':
                # Try text extraction first, then OCR if needed
                if self.detect_document_type(temp_file_path) == 'text_pdf':
                    processed_content = self.process_pdf(temp_file_path)
                else:
                    processed_content = self.process_scanned_pdf(temp_file_path)
            elif file_extension == '.docx':
                processed_content = self.process_word_document(temp_file_path)
            elif file_extension in ['.xlsx', '.xls']:
                processed_content = self.process_excel(temp_file_path)
            elif file_extension == '.csv':
                processed_content = self.process_csv(temp_file_path)
            elif file_extension == '.pptx':
                processed_content = self.process_powerpoint(temp_file_path)
            elif file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
                processed_content = self.process_image_ocr(temp_file_path)
            
            # Clean up temporary file
            if os.path.exists(temp_file_path):
                os.remove(temp_file_path)
            
            return AttachmentData(
                filename=filename,
                content_type=file_extension,
                size=size,
                processed_content=processed_content,
                extraction_success=processed_content is not None,
                error_message=None
            )
            
        except Exception as e:
            logger.error(f"Error processing attachment {filename}: {str(e)}")
            return AttachmentData(
                filename=filename,
                content_type=Path(filename).suffix.lower(),
                size=size,
                processed_content=None,
                extraction_success=False,
                error_message=str(e)
            )
    
    def process_email(self, file_path: str) -> EmailContent:
        """
        Parse .msg email files with comprehensive attachment processing
        
        Args:
            file_path: Path to the .msg file
            
        Returns:
            EmailContent with parsed email data and processed attachments
        """
        try:
            msg = extract_msg.Message(file_path)
            
            # Extract basic email information
            subject = msg.subject or ""
            sender = msg.sender or ""
            body = msg.body or ""
            date = msg.date
            
            # Extract recipients
            recipients = []
            if msg.to:
                recipients.extend([r.strip() for r in msg.to.split(';') if r.strip()])
            if msg.cc:
                recipients.extend([r.strip() for r in msg.cc.split(';') if r.strip()])
            
            # Process attachments
            attachments = []
            if msg.attachments:
                # Create temporary directory for attachment processing
                with tempfile.TemporaryDirectory() as temp_dir:
                    for attachment in msg.attachments:
                        attachment_data = self._process_attachment(attachment, temp_dir)
                        attachments.append(attachment_data)
            
            # Classify document type based on subject and body
            document_type = self._classify_document_type(f"{subject} {body}")
            
            # Metadata
            metadata = {
                'message_class': getattr(msg, 'messageClass', ''),
                'importance': getattr(msg, 'importance', ''),
                'sensitivity': getattr(msg, 'sensitivity', ''),
                'attachment_count': len(attachments),
                'processed_attachments': sum(1 for att in attachments if att.extraction_success)
            }
            
            msg.close()
            
            return EmailContent(
                subject=subject,
                sender=sender,
                recipients=recipients,
                body=body,
                attachments=attachments,
                date=date,
                metadata=metadata,
                document_type=document_type
            )
            
        except Exception as e:
            logger.error(f"Error processing email {file_path}: {str(e)}")
            return EmailContent(
                subject="",
                sender="",
                body="",
                attachments=[],
                metadata={'error': str(e)}
            )
    
    def process_word_document(self, file_path: str) -> WordDocumentData:
        """
        Process Word documents (.docx) using python-docx
        
        Args:
            file_path: Path to the Word document
            
        Returns:
            WordDocumentData with extracted content
        """
        try:
            if not DOCX_AVAILABLE:
                raise Exception("python-docx not available for Word document processing")
            
            doc = DocxDocument(file_path)
            
            # Extract paragraphs
            paragraphs = []
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    paragraphs.append(paragraph.text.strip())
            
            # Extract tables
            tables = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = []
                    for cell in row.cells:
                        row_data.append(cell.text.strip())
                    table_data.append(row_data)
                tables.append(table_data)
            
            # Combine all text
            full_text = '\n'.join(paragraphs)
            
            # Metadata
            metadata = {
                'file_format': '.docx',
                'paragraph_count': len(paragraphs),
                'table_count': len(tables),
                'processing_method': 'python_docx'
            }
            
            return WordDocumentData(
                text=full_text,
                paragraphs=paragraphs,
                tables=tables,
                metadata=metadata
            )
            
        except Exception as e:
            logger.error(f"Error processing Word document {file_path}: {str(e)}")
            return WordDocumentData(
                text="",
                metadata={'error': str(e)}
            )
    
    def process_powerpoint(self, file_path: str) -> PowerPointData:
        """
        Process PowerPoint presentations (.pptx) using python-pptx
        
        Args:
            file_path: Path to the PowerPoint file
            
        Returns:
            PowerPointData with extracted content
        """
        try:
            if not PPTX_AVAILABLE:
                raise Exception("python-pptx not available for PowerPoint processing")
            
            prs = Presentation(file_path)
            
            slides_data = []
            all_text = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = []
                slide_info = {
                    'slide_number': slide_num + 1,
                    'text': [],
                    'shapes': []
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                        slide_info['text'].append(shape.text.strip())
                    
                    # Record shape information
                    slide_info['shapes'].append({
                        'type': str(shape.shape_type),
                        'has_text': hasattr(shape, "text")
                    })
                
                slides_data.append(slide_info)
                if slide_text:
                    all_text.extend(slide_text)
            
            full_text = '\n'.join(all_text)
            
            # Metadata
            metadata = {
                'file_format': '.pptx',
                'processing_method': 'python_pptx'
            }
            
            return PowerPointData(
                slides=slides_data,
                text=full_text,
                metadata=metadata,
                total_slides=len(slides_data)
            )
            
        except Exception as e:
            logger.error(f"Error processing PowerPoint {file_path}: {str(e)}")
            return PowerPointData(
                slides=[],
                text="",
                metadata={'error': str(e)},
                total_slides=0
            )
    
    def process_csv(self, file_path: str) -> ExcelData:
        """
        Process CSV files using pandas
        
        Args:
            file_path: Path to the CSV file
            
        Returns:
            ExcelData with processed CSV data
        """
        try:
            # Read CSV with pandas
            df = pd.read_csv(file_path)
            
            # Convert to list of dictionaries
            csv_data = df.to_dict('records')
            
            # Metadata
            metadata = {
                'file_format': '.csv',
                'processing_method': 'pandas_csv'
            }
            
            return ExcelData(
                sheets={'Sheet1': csv_data},  # CSV has only one "sheet"
                metadata=metadata,
                total_rows=len(csv_data),
                total_sheets=1
            )
            
        except Exception as e:
            logger.error(f"Error processing CSV file {file_path}: {str(e)}")
            return ExcelData(
                metadata={'error': str(e)}
            )
    
    def process_excel(self, file_path: str) -> ExcelData:
        """
        Process Excel files using pandas and openpyxl
        
        Args:
            file_path: Path to the Excel file
            
        Returns:
            ExcelData with processed spreadsheet data
        """
        try:
            # Load workbook to get sheet names and metadata
            workbook = load_workbook(file_path, read_only=True)
            sheet_names = workbook.sheetnames
            
            sheets_data = {}
            total_rows = 0
            
            # Process each sheet
            for sheet_name in sheet_names:
                try:
                    # Read sheet data with pandas
                    df = pd.read_excel(file_path, sheet_name=sheet_name)
                    
                    # Convert to list of dictionaries
                    sheet_data = df.to_dict('records')
                    sheets_data[sheet_name] = sheet_data
                    total_rows += len(sheet_data)
                    
                except Exception as sheet_error:
                    logger.warning(f"Error processing sheet {sheet_name}: {str(sheet_error)}")
                    sheets_data[sheet_name] = []
            
            workbook.close()
            
            # Metadata
            metadata = {
                'file_format': Path(file_path).suffix.lower(),
                'sheet_names': sheet_names,
                'processing_method': 'pandas_openpyxl'
            }
            
            return ExcelData(
                sheets=sheets_data,
                metadata=metadata,
                total_rows=total_rows,
                total_sheets=len(sheet_names)
            )
            
        except Exception as e:
            logger.error(f"Error processing Excel file {file_path}: {str(e)}")
            return ExcelData(
                metadata={'error': str(e)}
            )
    
    def process_image_ocr(self, image_path: str) -> OCRResult:
        """
        Process image files using PaddleOCR
        
        Args:
            image_path: Path to the image file
            
        Returns:
            OCRResult with extracted text and regions
        """
        start_time = datetime.now()
        
        try:
            if not self.ocr_model or not PADDLEOCR_AVAILABLE:
                raise Exception("OCR model not initialized or PaddleOCR not available")
            
            # Perform OCR on the image
            ocr_result = self.ocr_model.ocr(image_path, cls=True)
            
            text_content = []
            regions = []
            total_confidence = 0
            word_count = 0
            
            if ocr_result and ocr_result[0]:
                for line in ocr_result[0]:
                    if line and len(line) >= 2:
                        bbox_coords = line[0]
                        text_info = line[1]
                        
                        if text_info and len(text_info) >= 2:
                            text = text_info[0]
                            confidence = text_info[1]
                            
                            text_content.append(text)
                            total_confidence += confidence
                            word_count += 1
                            
                            # Create text region
                            regions.append(TextRegion(
                                text=text,
                                confidence=confidence,
                                bbox=[
                                    min([p[0] for p in bbox_coords]),
                                    min([p[1] for p in bbox_coords]),
                                    max([p[0] for p in bbox_coords]),
                                    max([p[1] for p in bbox_coords])
                                ]
                            ))
            
            full_text = ' '.join(text_content)
            avg_confidence = total_confidence / word_count if word_count > 0 else 0
            
            processing_time = (datetime.now() - start_time).total_seconds()
            
            return OCRResult(
                text=full_text,
                confidence=avg_confidence,
                regions=regions,
                metadata={
                    'file_type': 'image',
                    'word_count': word_count,
                    'extraction_method': 'paddleocr'
                },
                processing_time=processing_time,
                success=True
            )
            
        except Exception as e:
            logger.error(f"Error processing image {image_path}: {str(e)}")
            processing_time = (datetime.now() - start_time).total_seconds()
            
            return OCRResult(
                text="",
                confidence=0.0,
                processing_time=processing_time,
                success=False,
                error_message=str(e)
            )
    
    def extract_text_regions(self, image: Union[np.ndarray, str]) -> List[TextRegion]:
        """
        Extract text regions from an image using PaddleOCR
        
        Args:
            image: Image as numpy array or path to image file
            
        Returns:
            List of TextRegion objects with detected text
        """
        try:
            if not self.ocr_model or not PADDLEOCR_AVAILABLE:
                raise Exception("OCR model not initialized or PaddleOCR not available")
            
            # Handle different input types
            if isinstance(image, str):
                # Image file path
                ocr_result = self.ocr_model.ocr(image, cls=True)
            elif isinstance(image, np.ndarray):
                # Numpy array
                ocr_result = self.ocr_model.ocr(image, cls=True)
            else:
                raise ValueError("Image must be a file path or numpy array")
            
            regions = []
            
            if ocr_result and ocr_result[0]:
                for line in ocr_result[0]:
                    if line and len(line) >= 2:
                        bbox_coords = line[0]
                        text_info = line[1]
                        
                        if text_info and len(text_info) >= 2:
                            text = text_info[0]
                            confidence = text_info[1]
                            
                            regions.append(TextRegion(
                                text=text,
                                confidence=confidence,
                                bbox=[
                                    min([p[0] for p in bbox_coords]),
                                    min([p[1] for p in bbox_coords]),
                                    max([p[0] for p in bbox_coords]),
                                    max([p[1] for p in bbox_coords])
                                ]
                            ))
            
            return regions
            
        except Exception as e:
            logger.error(f"Error extracting text regions: {str(e)}")
            return []
    
    def detect_document_type(self, file_path: str) -> str:
        """
        Detect if a PDF is text-based or scanned
        
        Args:
            file_path: Path to the PDF file
            
        Returns:
            'text_pdf' or 'scanned_pdf'
        """
        try:
            with pdfplumber.open(file_path) as pdf:
                # Check first few pages for text content
                text_found = False
                pages_to_check = min(3, len(pdf.pages))
                
                for i in range(pages_to_check):
                    page_text = pdf.pages[i].extract_text()
                    if page_text and len(page_text.strip()) > 50:
                        text_found = True
                        break
                
                return 'text_pdf' if text_found else 'scanned_pdf'
                
        except Exception as e:
            logger.error(f"Error detecting document type for {file_path}: {str(e)}")
            return 'scanned_pdf'  # Default to OCR processing
    
    def process_document(self, file_path: str, document_type: Optional[str] = None) -> Union[OCRResult, EmailContent, ExcelData, WordDocumentData, PowerPointData]:
        """
        Process any supported document type for reinsurance applications
        
        Args:
            file_path: Path to the document
            document_type: Optional document type hint
            
        Returns:
            Appropriate result object based on document type
        """
        file_extension = Path(file_path).suffix.lower()
        
        # Auto-detect document type if not provided
        if not document_type:
            if file_extension == '.pdf':
                document_type = self.detect_document_type(file_path)
            elif file_extension == '.msg':
                document_type = 'email'
            elif file_extension in ['.xlsx', '.xls']:
                document_type = 'excel'
            elif file_extension == '.csv':
                document_type = 'csv'
            elif file_extension == '.docx':
                document_type = 'word'
            elif file_extension == '.pptx':
                document_type = 'powerpoint'
            elif file_extension in ['.png', '.jpg', '.jpeg', '.tiff', '.bmp']:
                document_type = 'image'
            else:
                document_type = 'scanned_pdf'  # Default to OCR
        
        # Process based on type
        if document_type == 'text_pdf':
            return self.process_pdf(file_path)
        elif document_type == 'scanned_pdf':
            return self.process_scanned_pdf(file_path)
        elif document_type == 'email':
            return self.process_email(file_path)
        elif document_type == 'excel':
            return self.process_excel(file_path)
        elif document_type == 'csv':
            return self.process_csv(file_path)
        elif document_type == 'word':
            return self.process_word_document(file_path)
        elif document_type == 'powerpoint':
            return self.process_powerpoint(file_path)
        elif document_type == 'image':
            return self.process_image_ocr(file_path)
        else:
            raise ValueError(f"Unsupported document type: {document_type}")


# Global OCR agent instance
ocr_agent = OCRProcessingAgent()