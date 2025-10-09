import time
import datetime
from typing import Dict, Any, List, Optional, Union, Tuple
import ollama
import extract_msg
import json
import numpy as np
from pathlib import Path
import fitz
import cv2
from paddleocr import PaddleOCR
import traceback
import pandas as pd
from docx import Document
from pptx import Presentation
import openpyxl
import tempfile
import os
from PIL import Image
import re
from tqdm import tqdm
from config import OLLAMA_CONFIG, OCR_CONFIG, PROCESSING_CONFIG

class OfficeDocumentProcessor:
    """Process Office documents (docx, pptx, xlsx, csv) with table extraction"""

    @staticmethod
    def extract_docx(file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extract text and tables from Word documents with table detection

        Args:
            file_path: Path to the .docx file

        Returns:
            Dictionary containing paragraphs, tables, and metadata
        """
        doc: Document = Document(file_path)
        content: Dict[str, Any] = {
            'paragraphs': [],
            'tables': [],
            'metadata': {
                'total_paragraphs': 0,
                'total_tables': 0,
                'table_validation': []
            }
        }

        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                content['paragraphs'].append(para.text)

        # Enhanced table extraction
        for table_idx, table in enumerate(doc.tables):
            table_data: Dict[str, Any] = {
                'table_number': table_idx + 1,
                'rows': len(table.rows),
                'columns': len(table.columns),
                'data': [],
                'dataframe': None,
                'validation': {
                    'has_header': False,
                    'empty_cells': 0,
                    'merged_cells': 0,
                    'data_quality_score': 0.0
                }
            }

            # Extract table data with validation
            total_cells: int = 0
            empty_cells: int = 0
            for row_idx, row in enumerate(table.rows):
                row_data: List[str] = []
                for cell in row.cells:
                    cell_text: str = cell.text.strip()
                    total_cells += 1
                    if not cell_text:
                        empty_cells += 1
                    row_data.append(cell_text)
                table_data['data'].append(row_data)

            # Calculate data quality metrics
            table_data['validation']['empty_cells'] = empty_cells
            table_data['validation']['data_quality_score'] = (
                (total_cells - empty_cells) / total_cells if total_cells > 0 else 0
            )

            try:
                if len(table_data['data']) > 1:
                    # Check if first row looks like headers
                    first_row: List[str] = table_data['data'][0]
                    header_score: float = sum(
                        1 for cell in first_row
                        if cell and not cell.replace('.', '').replace(',', '').isdigit()
                    ) / len(first_row)

                    df: pd.DataFrame
                    if header_score > 0.6:  # More than 60% non-numeric
                        table_data['validation']['has_header'] = True
                        df = pd.DataFrame(table_data['data'][1:], columns=first_row)
                    else:
                        df = pd.DataFrame(table_data['data'])
                else:
                    df = pd.DataFrame(table_data['data'])

                # Clean DataFrame
                df = df.dropna(how='all').dropna(axis=1, how='all')

                if not df.empty:
                    table_data['dataframe'] = df.to_dict('records')
                    table_data['dataframe_shape'] = df.shape
                    table_data['columns'] = df.columns.tolist()
                else:
                    table_data['dataframe'] = []
                    table_data['dataframe_shape'] = (0, 0)
                    table_data['columns'] = []

            except Exception as e:
                print(f"Could not convert table {table_idx + 1} to DataFrame: {e}")
                table_data['dataframe'] = []
                table_data['dataframe_shape'] = (0, 0)
                table_data['columns'] = []

            content['tables'].append(table_data)
            content['metadata']['table_validation'].append(table_data['validation'])

        content['metadata']['total_paragraphs'] = len(content['paragraphs'])
        content['metadata']['total_tables'] = len(content['tables'])

        return content

    @staticmethod
    def extract_pptx(file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extract text and tables from PowerPoint presentations

        Args:
            file_path: Path to the .pptx file

        Returns:
            Dictionary containing slides with text, tables, and metadata
        """
        prs: Presentation = Presentation(file_path)
        content: Dict[str, Any] = {
            'slides': [],
            'metadata': {
                'total_slides': len(prs.slides),
                'total_tables': 0
            }
        }

        for slide_idx, slide in enumerate(prs.slides):
            slide_content: Dict[str, Any] = {
                'slide_number': slide_idx + 1,
                'title': '',
                'text': [],
                'tables': []
            }

            # Extract title
            if slide.shapes.title:
                slide_content['title'] = slide.shapes.title.text

            # Extract text and tables
            for shape in slide.shapes:
                # Text extraction
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content['text'].append(shape.text)

                # Table extraction
                if shape.has_table:
                    table = shape.table
                    table_data: Dict[str, Any] = {
                        'table_number': len(slide_content['tables']) + 1,
                        'rows': len(table.rows),
                        'columns': len(table.columns),
                        'data': [],
                        'dataframe': None
                    }

                    # Extract table data
                    for row in table.rows:
                        row_data: List[str] = [cell.text.strip() for cell in row.cells]
                        table_data['data'].append(row_data)

                    # Convert to DataFrame
                    try:
                        df: pd.DataFrame
                        if len(table_data['data']) > 1:
                            df = pd.DataFrame(table_data['data'][1:], columns=table_data['data'][0])
                            table_data['dataframe'] = df.to_dict('records')
                            table_data['dataframe_shape'] = df.shape
                        else:
                            df = pd.DataFrame(table_data['data'])
                            table_data['dataframe'] = df.to_dict('records')
                            table_data['dataframe_shape'] = df.shape
                    except Exception as e:
                        print(f"Could not convert table to DataFrame: {e}")

                    slide_content['tables'].append(table_data)
                    content['metadata']['total_tables'] += 1

            content['slides'].append(slide_content)

        return content

    @staticmethod
    def extract_xlsx(file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Excel extraction with sheet analysis

        Args:
            file_path: Path to the .xlsx file

        Returns:
            Dictionary containing sheets with data and metadata
        """
        wb: openpyxl.Workbook = openpyxl.load_workbook(file_path, data_only=True)
        content: Dict[str, Any] = {
            'sheets': [],
            'metadata': {
                'total_sheets': len(wb.sheetnames),
                'sheet_names': wb.sheetnames,
                'sheet_analysis': []
            }
        }

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_data: Dict[str, Any] = {
                'sheet_name': sheet_name,
                'used_range': f"{ws.dimensions}",
                'data': [],
                'dataframe': None,
                'analysis': {
                    'has_data': False,
                    'data_density': 0.0,
                    'likely_headers': False
                }
            }

            # Extract all rows
            all_rows: List[List[Any]] = []
            for row in ws.iter_rows(values_only=True):
                all_rows.append(list(row))

            sheet_data['data'] = all_rows

            # Analyze sheet content
            if all_rows:
                # Check if sheet has meaningful data
                non_empty_cells: int = sum(
                    1 for row in all_rows for cell in row
                    if cell is not None and str(cell).strip()
                )
                total_cells: int = sum(len(row) for row in all_rows)
                sheet_data['analysis']['data_density'] = (
                    non_empty_cells / total_cells if total_cells > 0 else 0
                )
                sheet_data['analysis']['has_data'] = sheet_data['analysis']['data_density'] > 0.1

                try:
                    if len(all_rows) > 1:
                        # Better header detection
                        first_row: List[Any] = all_rows[0]
                        header_indicators: int = sum(
                            1 for cell in first_row
                            if cell and isinstance(cell, str) and
                            not str(cell).replace('.', '').replace(',', '').replace('-', '').isdigit()
                        )

                        df: pd.DataFrame
                        if header_indicators > len(first_row) * 0.5:
                            sheet_data['analysis']['likely_headers'] = True
                            df = pd.DataFrame(all_rows[1:], columns=first_row)
                        else:
                            df = pd.DataFrame(all_rows)
                    else:
                        df = pd.DataFrame(all_rows)

                    # Clean DataFrame
                    df = df.dropna(how='all').dropna(axis=1, how='all')

                    if not df.empty:
                        sheet_data['dataframe'] = df.to_dict('records')
                        sheet_data['dataframe_shape'] = df.shape
                        sheet_data['columns'] = df.columns.tolist()
                    else:
                        sheet_data['dataframe'] = []
                        sheet_data['dataframe_shape'] = (0, 0)
                        sheet_data['columns'] = []

                except Exception as e:
                    print(f"Could not convert sheet '{sheet_name}' to DataFrame: {e}")
                    sheet_data['dataframe'] = []
                    sheet_data['dataframe_shape'] = (0, 0)
                    sheet_data['columns'] = []

            content['sheets'].append(sheet_data)
            content['metadata']['sheet_analysis'].append(sheet_data['analysis'])

        return content

    @staticmethod
    def extract_csv(file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extract data from CSV files with encoding detection

        Args:
            file_path: Path to the .csv file

        Returns:
            Dictionary containing CSV data and metadata
        """
        content: Dict[str, Any] = {
            'filename': Path(file_path).name,
            'data': [],
            'dataframe': None,
            'metadata': {
                'encoding': 'unknown',
                'delimiter': 'unknown',
                'data_quality': 0.0
            }
        }

        # Try multiple encodings
        encodings: List[str] = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']

        for encoding in encodings:
            try:
                df: pd.DataFrame = pd.read_csv(file_path, encoding=encoding, on_bad_lines='skip')
                content['metadata']['encoding'] = encoding

                # Detect delimiter
                with open(file_path, 'r', encoding=encoding) as f:
                    sample: str = f.read(1024)
                    delimiters: List[str] = [',', ';', '\t', '|']
                    delimiter_counts: Dict[str, int] = {
                        delim: sample.count(delim) for delim in delimiters
                    }
                    content['metadata']['delimiter'] = max(delimiter_counts, key=delimiter_counts.get)

                # Calculate data quality
                total_cells: int = df.size
                non_null_cells: int = df.count().sum()
                content['metadata']['data_quality'] = (
                    non_null_cells / total_cells if total_cells > 0 else 0
                )

                content['dataframe'] = df.to_dict('records')
                content['dataframe_shape'] = df.shape
                content['columns'] = df.columns.tolist()
                content['metadata']['rows'] = len(df)
                content['metadata']['columns'] = len(df.columns)

                # Store raw data
                content['data'] = [df.columns.tolist()] + df.values.tolist()
                break

            except Exception as e:
                continue

        return content

class EmailProcessor:
    """Process .msg files with existing attachments in folder structure"""

    def __init__(self, subfolder_path: Union[str, Path]) -> None:
        """
        Initialize EmailProcessor

        Args:
            subfolder_path: Path to the folder containing the .msg file
        """
        self.subfolder_path: Path = Path(subfolder_path)
        self.msg_file: Optional[Path] = None
        self.msg_data: Optional[Dict[str, Any]] = None

    def find_msg_file(self) -> Optional[Path]:
        """
        Find the .msg file in the subfolder

        Returns:
            Path to .msg file if found, None otherwise
        """
        msg_files: List[Path] = list(self.subfolder_path.glob('*.msg'))
        if msg_files:
            self.msg_file = msg_files[0]
            return self.msg_file
        return None

    def _safe_str(self, value: Any) -> Optional[str]:
        """
        Safely convert any value to string, handling bytes

        Args:
            value: Value to convert

        Returns:
            String representation or None
        """
        if value is None:
            return None
        if isinstance(value, bytes):
            try:
                return value.decode('utf-8', errors='ignore')
            except:
                return str(value)
        return str(value)

    def extract_email_data(self) -> Optional[Dict[str, Any]]:
        """
        Extract structured data from .msg file

        Returns:
            Dictionary containing email metadata and content, or None on error
        """
        if not self.msg_file:
            self.find_msg_file()

        if not self.msg_file:
            return None

        try:
            msg = extract_msg.Message(str(self.msg_file))

            self.msg_data = {
                'metadata': {
                    'subject': self._safe_str(msg.subject) if msg.subject else '',
                    'sender': self._safe_str(msg.sender) if msg.sender else '',
                    'sender_email': self._safe_str(getattr(msg, 'senderEmail', None)),
                    'date': self._safe_str(msg.date) if msg.date else None,
                    'received_time': self._safe_str(getattr(msg, 'receivedTime', None)) if hasattr(msg,
                                                                                                   'receivedTime') else None,
                    'message_id': self._safe_str(msg.messageId) if msg.messageId else '',
                    'importance': self._safe_str(msg.importance) if msg.importance else '',
                },
                'recipients': {
                    'to': self._format_recipients(msg.to),
                    'cc': self._format_recipients(msg.cc),
                    'bcc': self._format_recipients(msg.bcc)
                },
                'body': {
                    'plain_text': self._safe_str(msg.body) if msg.body else '',
                    'html': self._safe_str(msg.htmlBody) if msg.htmlBody else '',
                },
                'attachments_in_msg': self._get_msg_attachments(msg),
                'folder_path': str(self.subfolder_path),
                'msg_filename': self.msg_file.name
            }

            msg.close()
            return self.msg_data

        except Exception as e:
            print(f"Error extracting {self.msg_file}: {e}")
            traceback.print_exc()
            return None

    def _format_recipients(self, recipients: Any) -> List[Dict[str, str]]:
        """
        Format recipient list - handles both string and object formats

        Args:
            recipients: Recipients in various formats

        Returns:
            List of dictionaries with name and email
        """
        if not recipients:
            return []

        formatted: List[Dict[str, str]] = []

        if isinstance(recipients, str):
            return [{'name': self._safe_str(recipients), 'email': self._safe_str(recipients)}]

        if isinstance(recipients, list):
            for r in recipients:
                if isinstance(r, str):
                    formatted.append({'name': self._safe_str(r), 'email': self._safe_str(r)})
                elif hasattr(r, 'name') and hasattr(r, 'email'):
                    formatted.append({
                        'name': self._safe_str(r.name) if r.name else '',
                        'email': self._safe_str(r.email) if r.email else ''
                    })
                else:
                    formatted.append({'name': self._safe_str(r), 'email': self._safe_str(r)})

        return formatted

    def _get_msg_attachments(self, msg: Any) -> List[Dict[str, Any]]:
        """
        Get attachment info from .msg file - WITHOUT binary data

        Args:
            msg: Message object from extract_msg

        Returns:
            List of attachment metadata dictionaries
        """
        attachments: List[Dict[str, Any]] = []
        try:
            for att in msg.attachments:
                filename: str = att.longFilename or att.shortFilename or 'unknown'
                att_info: Dict[str, Any] = {
                    'filename': self._safe_str(filename),
                    'size': len(att.data) if att.data else 0,
                }
                attachments.append(att_info)
        except Exception as e:
            print(f"  Warning: Could not read attachments from msg: {e}")

        return attachments

    def list_folder_files(self) -> List[Dict[str, Any]]:
        """
        List all files in the subfolder (saved attachments)

        Returns:
            List of file metadata dictionaries
        """
        all_files: List[Dict[str, Any]] = []
        for file in self.subfolder_path.iterdir():
            if file.is_file() and not file.name.endswith('.msg'):
                all_files.append({
                    'filename': file.name,
                    'path': str(file),
                    'size': file.stat().st_size
                })
        return all_files

    def get_complete_data(self) -> Optional[Dict[str, Any]]:
        """
        Get email data + list of existing attachment files

        Returns:
            Complete email data dictionary or None on error
        """
        email_data: Optional[Dict[str, Any]] = self.extract_email_data()
        if not email_data:
            return None

        email_data['saved_attachments'] = self.list_folder_files()
        return email_data

    def process_attachments_with_ocr(
            self,
            email_data: Dict[str, Any],
            ocr_engine: 'Detect_and_Recognize',
            save_visuals: bool = True
    ) -> List[Dict[str, Any]]:
        """
        Process attachments with OCR, office documents, and direct text extraction.

        Args:
            email_data: Dictionary containing email metadata and attachments
            ocr_engine: Instance of Detect_and_Recognize class for OCR operations
            save_visuals: If True, saves annotated images and JSON results to:
                         <email_folder>/text_detected_and_recognized/<TIMESTAMP>/

        Returns:
            List of dictionaries containing extracted text and metadata for each attachment
        """
        attachment_texts: List[Dict[str, Any]] = []
        office_processor = OfficeDocumentProcessor()

        # If saving visuals, create timestamped output root in the email folder
        out_root: Optional[Path] = None
        if save_visuals:
            ts: str = datetime.datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
            out_root = Path(self.subfolder_path) / "text_detected_and_recognized" / ts
            out_root.mkdir(parents=True, exist_ok=True)
            tqdm.write(f"[OCR] Saving visuals to: {out_root}")

        attachments: List[Dict[str, Any]] = email_data.get('saved_attachments', [])

        # Main attachment processing loop with progress bar
        for attachment in tqdm(attachments, desc="Processing attachments", unit="file", leave=False):
            file_path: str = attachment['path']
            ext: str = Path(file_path).suffix.lower()
            base_name: str = Path(file_path).stem

            if ext in ['.json']:
                continue

            start_time: float = time.time()

            # OFFICE DOCUMENTS - Process first before OCR
            if ext in ['.docx']:
                try:
                    content: Dict[str, Any] = office_processor.extract_docx(file_path)

                    # Extract all text for LLM context
                    text_parts: List[str] = content['paragraphs'].copy()
                    for table in content['tables']:
                        table_text: str = f"\n[TABLE {table['table_number']} - {table['rows']}x{table['columns']}]\n"
                        for row in table['data']:
                            table_text += " | ".join(str(cell) for cell in row) + "\n"
                        text_parts.append(table_text)

                    full_text: str = "\n".join(text_parts)
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': full_text,
                        'method': 'docx_extraction',
                        'structured_content': content,
                        'time': time.time() - start_time
                    })

                    # Save structured data
                    if save_visuals and out_root is not None:
                        json_path: Path = out_root / f"{base_name}_docx_extracted.json"
                        with open(json_path, 'w', encoding='utf-8') as jf:
                            json.dump(content, jf, indent=2, ensure_ascii=False)

                except Exception as e:
                    tqdm.write(f"[DOCX] Error processing {file_path}: {e}")

            elif ext in ['.pptx']:
                try:
                    content: Dict[str, Any] = office_processor.extract_pptx(file_path)

                    # Extract all text for LLM context
                    text_parts: List[str] = []
                    for slide in content['slides']:
                        slide_text: str = f"\n[SLIDE {slide['slide_number']}: {slide['title']}]\n"
                        slide_text += "\n".join(slide['text'])

                        for table in slide['tables']:
                            table_text: str = f"\n[TABLE {table['table_number']} - {table['rows']}x{table['columns']}]\n"
                            for row in table['data']:
                                table_text += " | ".join(str(cell) for cell in row) + "\n"
                            slide_text += table_text

                        text_parts.append(slide_text)

                    full_text: str = "\n".join(text_parts)
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': full_text,
                        'method': 'pptx_extraction',
                        'structured_content': content,
                        'time': time.time() - start_time
                    })

                    # Save structured data
                    if save_visuals and out_root is not None:
                        json_path: Path = out_root / f"{base_name}_pptx_extracted.json"
                        with open(json_path, 'w', encoding='utf-8') as jf:
                            json.dump(content, jf, indent=2, ensure_ascii=False)

                except Exception as e:
                    tqdm.write(f"[PPTX] Error processing {file_path}: {e}")

            elif ext in ['.xlsx']:
                try:
                    content: Dict[str, Any] = office_processor.extract_xlsx(file_path)

                    # Extract all sheets for LLM context
                    text_parts: List[str] = []
                    for sheet in content['sheets']:
                        sheet_text: str = f"\n[SHEET: {sheet['sheet_name']} - {sheet['used_range']}]\n"
                        if sheet['dataframe']:
                            # Create table representation
                            for row in sheet['data'][:100]:  # Limit to first 100 rows
                                sheet_text += " | ".join(str(cell) if cell is not None else "" for cell in row) + "\n"
                        text_parts.append(sheet_text)

                    full_text: str = "\n".join(text_parts)
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': full_text,
                        'method': 'xlsx_extraction',
                        'structured_content': content,
                        'time': time.time() - start_time
                    })

                    # Save structured data
                    if save_visuals and out_root is not None:
                        json_path: Path = out_root / f"{base_name}_xlsx_extracted.json"
                        with open(json_path, 'w', encoding='utf-8') as jf:
                            json.dump(content, jf, indent=2, ensure_ascii=False)

                except Exception as e:
                    tqdm.write(f"[XLSX] Error processing {file_path}: {e}")

            elif ext in ['.csv']:
                try:
                    content: Dict[str, Any] = office_processor.extract_csv(file_path)

                    # Create text representation
                    text_parts: List[str] = [
                        f"[CSV FILE - {content['metadata'].get('rows', 0)} rows x {content['metadata'].get('columns', 0)} columns]\n"
                    ]
                    for row in content['data'][:100]:  # Limit to first 100 rows
                        text_parts.append(" | ".join(str(cell) if cell is not None else "" for cell in row))

                    full_text: str = "\n".join(text_parts)
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': full_text,
                        'method': 'csv_extraction',
                        'structured_content': content,
                        'time': time.time() - start_time
                    })

                    # Save structured data
                    if save_visuals and out_root is not None:
                        json_path: Path = out_root / f"{base_name}_csv_extracted.json"
                        with open(json_path, 'w', encoding='utf-8') as jf:
                            json.dump(content, jf, indent=2, ensure_ascii=False)

                except Exception as e:
                    tqdm.write(f"[CSV] Error processing {file_path}: {e}")

            # PDF - OCR Processing with nested progress bar
            elif ext in ['.pdf']:
                try:
                    pdf_start: float = time.time()
                    images: List[Any]
                    originals: List[Any]
                    images, originals = ocr_engine.pdf_to_images(file_path)

                    # Nested progress bar for PDF pages
                    for page_num, (processed, original) in tqdm(
                            enumerate(zip(images, originals)),
                            desc=f"  OCR pages: {Path(file_path).name}",
                            total=len(images),
                            unit="pg",
                            leave=False
                    ):
                        page_start: float = time.time()

                        # Save temp image for PaddleOCR
                        with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                            img_pil: Image.Image = Image.fromarray(original)
                            img_pil.save(tmp.name)
                            tmp_path: str = tmp.name

                        # OCR with timing
                        ocr_start: float = time.time()
                        results: List[Any] = ocr_engine.ocr_with_detection_and_recognition(tmp_path)
                        ocr_time: float = time.time() - ocr_start

                        # SAVE FIRST - before any filtering
                        if save_visuals and out_root is not None and results:
                            save_start: float = time.time()
                            try:
                                # PaddleOCR returns a list - get first result
                                if isinstance(results, list) and len(results) > 0:
                                    result_obj: Any = results[0]

                                    # Create page-specific subdirectory for this page's outputs
                                    page_dir: Path = out_root / f"{base_name}_page{page_num + 1}"
                                    page_dir.mkdir(exist_ok=True)

                                    # save_to_img expects a DIRECTORY, not a filename
                                    result_obj.save_to_img(str(page_dir))

                                    # Save JSON with specific filename
                                    json_path: Path = page_dir / "ocr_result.json"
                                    result_obj.save_to_json(str(json_path))

                                    save_time: float = time.time() - save_start
                                    tqdm.write(
                                        f"    [PDF] Page {page_num + 1} - OCR: {ocr_time:.2f}s, Save: {save_time:.2f}s")
                            except Exception as e:
                                tqdm.write(f"    [PDF] Warning: Could not save page {page_num + 1}: {e}")
                                traceback.print_exc()

                        text_parts: List[str] = []
                        for res in results:
                            # Access the JSON dict directly from the result object
                            result_data: Dict[str, Any]
                            if isinstance(res, dict):
                                result_data = res
                            else:
                                result_data = getattr(res, "json", None) or res

                            rec_texts: List[str] = result_data.get('rec_texts', []) or []
                            rec_scores: List[float] = result_data.get('rec_scores', []) or []

                            if rec_scores and len(rec_scores) == len(rec_texts):
                                for t, s in zip(rec_texts, rec_scores):
                                    try:
                                        if float(s) >= 0.7 and str(t).strip():
                                            text_parts.append(str(t).strip())
                                    except Exception:
                                        continue
                            else:
                                for t in rec_texts:
                                    if isinstance(t, str) and len(t.strip()) >= 2:
                                        text_parts.append(t.strip())

                        full_text: str = ' '.join(text_parts).strip()

                        # Only skip ADDING to attachment_texts if no text
                        # But we already saved the visual outputs above
                        if full_text:
                            attachment_texts.append({
                                'file': attachment['filename'],
                                'page': page_num + 1,
                                'text': full_text,
                                'method': 'ocr',
                                'time': time.time() - page_start,
                                'ocr_time': ocr_time
                            })

                        # Cleanup temp file
                        os.unlink(tmp_path)

                    tqdm.write(f"  [PDF] Total: {time.time() - pdf_start:.2f}s")

                except Exception as e:
                    tqdm.write(f"[OCR] Error processing PDF {file_path}: {e}")
                    traceback.print_exc()

            # IMAGE FILES - FIXED (matches your current code structure)
            elif ext in ['.png', '.jpg', '.jpeg', '.bmp', '.tiff', '.tif']:
                try:
                    img_start: float = time.time()

                    # OCR with timing
                    ocr_start: float = time.time()
                    results: List[Any] = ocr_engine.ocr_with_detection_and_recognition(file_path)
                    ocr_time: float = time.time() - ocr_start

                    # SAVE FIRST - before any text filtering
                    if save_visuals and out_root is not None and results:
                        save_start: float = time.time()
                        try:
                            # Handle results structure - should be a list
                            result_obj: Any
                            if isinstance(results, list) and len(results) > 0:
                                result_obj = results[0]
                            else:
                                result_obj = results

                            # Create file-specific subdirectory for this image's outputs
                            img_dir: Path = out_root / f"{base_name}_image"
                            img_dir.mkdir(exist_ok=True)

                            # save_to_img expects a DIRECTORY
                            result_obj.save_to_img(str(img_dir))

                            # Save JSON with specific filename
                            json_path: Path = img_dir / "ocr_result.json"
                            result_obj.save_to_json(str(json_path))

                            save_time: float = time.time() - save_start
                            tqdm.write(f"  [IMG] {base_name} - OCR: {ocr_time:.2f}s, Save: {save_time:.2f}s")
                        except Exception as e:
                            tqdm.write(f"  [IMG] Warning: Could not save visuals: {e}")
                            traceback.print_exc()

                    # NOW extract text from result objects
                    text_parts: List[str] = []
                    for res in results:
                        # Access the JSON dict directly from the result object
                        result_data: Dict[str, Any]
                        if isinstance(res, dict):
                            result_data = res
                        else:
                            result_data = getattr(res, "json", None) or res

                        rec_texts: List[str] = result_data.get('rec_texts', []) or []
                        rec_scores: List[float] = result_data.get('rec_scores', []) or []

                        if rec_scores and len(rec_scores) == len(rec_texts):
                            for t, s in zip(rec_texts, rec_scores):
                                try:
                                    if float(s) >= 0.7 and str(t).strip():
                                        text_parts.append(str(t).strip())
                                except Exception:
                                    continue
                        else:
                            for t in rec_texts:
                                if isinstance(t, str) and len(t.strip()) >= 2:
                                    text_parts.append(t.strip())

                    full_text: str = ' '.join(text_parts).strip()

                    # Only add to attachment_texts if there's filtered text
                    # But we already saved visuals above regardless
                    if full_text:
                        attachment_texts.append({
                            'file': attachment['filename'],
                            'text': full_text,
                            'method': 'ocr',
                            'time': time.time() - img_start,
                            'ocr_time': ocr_time
                        })

                except Exception as e:
                    tqdm.write(f"[OCR] Error processing image {file_path}: {e}")
                    traceback.print_exc()

            # TEXT FILES - Direct extraction
            elif ext in ['.txt', '.log']:
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text: str = f.read()
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': text,
                        'method': 'direct',
                        'time': time.time() - start_time
                    })

                    # optionally save a JSON copy of the raw text to visuals folder
                    if save_visuals and out_root is not None:
                        json_name: str = f"{base_name}_raw_text.json"
                        json_path: Path = out_root / json_name
                        try:
                            with open(json_path, 'w', encoding='utf-8') as jf:
                                json.dump({"file": attachment['filename'], "text": text}, jf, indent=2,
                                          ensure_ascii=False)
                        except Exception as e:
                            tqdm.write(f"[OCR] Warning: failed to save raw text json {json_path}: {e}")
                except Exception as e:
                    tqdm.write(f"[TXT] Error processing {file_path}: {e}")

            # UNKNOWN - Try OCR as fallback
            else:
                try:
                    ocr_start: float = time.time()
                    results: List[Any] = ocr_engine.ocr_with_detection_and_recognition(file_path)
                    ocr_time: float = time.time() - ocr_start

                    # Extract text from results
                    text_parts: List[str] = []
                    for res in results:
                        result_data: Dict[str, Any]
                        if isinstance(res, dict):
                            result_data = res
                        else:
                            result_data = getattr(res, "json", None) or res

                        rec_texts: List[str] = result_data.get('rec_texts', []) or []
                        rec_scores: List[float] = result_data.get('rec_scores', []) or []

                        if rec_scores and len(rec_scores) == len(rec_texts):
                            for t, s in zip(rec_texts, rec_scores):
                                try:
                                    if float(s) >= 0.7 and str(t).strip():
                                        text_parts.append(str(t).strip())
                                except Exception:
                                    continue

                    text: str = ' '.join(text_parts).strip()

                    if text:
                        attachment_texts.append({
                            'file': attachment['filename'],
                            'text': text,
                            'method': 'ocr_fallback',
                            'time': time.time() - start_time,
                            'ocr_time': ocr_time
                        })

                    if save_visuals and out_root is not None and results:
                        try:
                            result_obj: Any = results[0] if isinstance(results, list) and len(results) > 0 else results
                            unknown_dir: Path = out_root / f"{base_name}_unknown"
                            unknown_dir.mkdir(exist_ok=True)
                            result_obj.save_to_img(str(unknown_dir))
                            json_path: Path = unknown_dir / "ocr_result.json"
                            result_obj.save_to_json(str(json_path))
                        except Exception as e:
                            tqdm.write(f"[OCR] Warning: failed to save unknown file type visuals: {e}")

                except Exception as e:
                    tqdm.write(f"[OCR] Skipping unknown file type {file_path}: {e}")

        return attachment_texts


class LocalLLMProcessor:
    """Process LLM requests using local model pulled through Ollama"""

    def __init__(self, model_name: Optional[str] = None) -> None:
        self.model_name: str = model_name or OLLAMA_CONFIG["model_name"]

    def test_connection(self) -> bool:
        """Test if local Ollama is running and model is available"""
        try:
            with tqdm(total=3, desc="Testing Ollama connection", leave=False) as pbar:
                pbar.set_description("Fetching model list")
                models = ollama.list()
                pbar.update(1)

                # ListResponse object has a 'models' attribute
                model_list = models.models if hasattr(models, 'models') else []

                # Extract model names
                pbar.set_description("Extracting model names")
                model_names: List[str] = []
                for m in model_list:
                    if hasattr(m, 'model'):
                        model_names.append(m.model)
                    elif hasattr(m, 'name'):
                        model_names.append(m.name)
                pbar.update(1)

                # Check if target model exists
                pbar.set_description("Verifying model availability")
                if any(self.model_name in name for name in model_names):
                    print(f"âœ“ Local model {self.model_name} is available")
                    pbar.update(1)
                    return True
                else:
                    print(f"âœ— Model {self.model_name} not found in model list.")
                    print(f"   Available models: {model_names}")

                    # Try to actually use the model anyway - maybe it's there but list failed
                    try:
                        print(f"   Attempting to use model directly...")
                        test_result = ollama.generate(
                            model=self.model_name,
                            prompt="Respond with just the word 'OK' and nothing else.",
                            options={'num_predict': 10}
                        )

                        # Verify we got a valid response
                        response_text: str = test_result.get("response", "").strip()

                        if response_text:
                            print(f"âœ“ Model {self.model_name} is actually available!")
                            print(
                                f"   Test response received: '{response_text[:50]}{'...' if len(response_text) > 50 else ''}'")
                            pbar.update(1)
                            return True
                        else:
                            print(f"âœ— Model returned empty response")
                            print(f"   To pull the model, run: ollama pull {self.model_name}")
                            pbar.update(1)
                            return False

                    except Exception as test_error:
                        print(f"âœ— Model test failed: {test_error}")
                        print(f"   To pull the model, run: ollama pull {self.model_name}")
                        pbar.update(1)
                        return False

        except Exception as e:
            print(f"âœ— Cannot connect to local Ollama: {e}")
            traceback.print_exc()
            return False

    def generate_response(
            self,
            prompt: str,
            max_tokens: int = 4000,
            temperature: float = 0.1
    ) -> Dict[str, Any]:
        """
        Generate response from local Ollama model using ollama library

        Args:
            prompt: The input prompt
            max_tokens: Maximum tokens to generate
            temperature: Sampling temperature (0.0 to 1.0)

        Returns:
            Dictionary with response data and metadata
        """
        start_time: float = time.time()

        try:
            with tqdm(total=2, desc="Generating response", leave=False) as pbar:
                pbar.set_description("Sending request to model")
                result = ollama.generate(
                    model=self.model_name,
                    prompt=prompt,
                    options={
                        "temperature": temperature,
                        "num_predict": max_tokens,
                        "top_p": 0.9,
                        "top_k": 40
                    }
                )
                pbar.update(1)

                pbar.set_description("Processing response")
                generation_time: float = time.time() - start_time
                response_text: str = result.get("response", "")
                pbar.update(1)

            return {
                "success": True,
                "response": response_text,
                "model": result.get("model", self.model_name),
                "generation_time": generation_time,
                "tokens_generated": len(response_text.split()),
                "metadata": {
                    "prompt_length": len(prompt),
                    "temperature": temperature,
                    "max_tokens": max_tokens
                }
            }

        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "generation_time": time.time() - start_time
            }

    def extract_structured_data(self, llm_prompt: str) -> Dict[str, Any]:
        """
        Extract structured reinsurance data using the local LLM

        Args:
            llm_prompt: The formatted prompt with email and attachment data

        Returns:
            Dictionary with extracted data and processing metadata
        """
        print("ðŸ¤– Sending request to local Ollama...")

        # Generate response
        result: Dict[str, Any] = self.generate_response(
            llm_prompt,
            max_tokens=OLLAMA_CONFIG["max_tokens"],
            temperature=OLLAMA_CONFIG["temperature"]
        )

        if not result["success"]:
            return {
                "success": False,
                "error": result["error"],
                "raw_response": None,
                "extracted_data": None
            }

        # Try to parse JSON response
        raw_response: str = result["response"].strip()

        # Clean up response - remove any markdown formatting
        with tqdm(total=3, desc="Processing LLM output", leave=False) as pbar:
            pbar.set_description("Cleaning response")
            if raw_response.startswith("```json"):
                raw_response = raw_response[7:]
            if raw_response.endswith("```"):
                raw_response = raw_response[:-3]
            raw_response = raw_response.strip()
            pbar.update(1)

            try:
                pbar.set_description("Parsing JSON")
                extracted_data: Dict[str, Any] = json.loads(raw_response)
                pbar.update(2)

                return {
                    "success": True,
                    "raw_response": raw_response,
                    "extracted_data": extracted_data,
                    "generation_time": result["generation_time"],
                    "tokens_generated": result["tokens_generated"],
                    "model_used": result["model"]
                }
            except json.JSONDecodeError as e:
                pbar.set_description("JSON parse failed, trying extraction")
                pbar.update(1)

                print(f"âš ï¸  JSON parsing failed: {e}")
                print(f"Raw response: {raw_response[:500]}...")

                # Try to extract JSON from the response
                json_match = re.search(r'\{.*\}', raw_response, re.DOTALL)
                if json_match:
                    try:
                        extracted_data = json.loads(json_match.group())
                        pbar.update(1)
                        return {
                            "success": True,
                            "raw_response": raw_response,
                            "extracted_data": extracted_data,
                            "generation_time": result["generation_time"],
                            "tokens_generated": result["tokens_generated"],
                            "model_used": result["model"],
                            "warning": "JSON extracted from response with regex"
                        }
                    except json.JSONDecodeError:
                        pass

                pbar.update(1)
                return {
                    "success": False,
                    "error": f"JSON parsing failed: {e}",
                    "raw_response": raw_response,
                    "extracted_data": None,
                    "generation_time": result["generation_time"]
                }

def extract_candidates_from_text(text: str) -> Dict[str, Any]:
    """
    Return a small dict of candidate values useful as LLM hints.

    Keys possibly returned:
      - insured: str
      - cedant: str
      - total_sum_insured: str (raw)
      - total_sum_insured_float: Optional[float]
      - currency: Optional[str] (symbol or 3-letter code)
      - period_of_insurance: str
      - country: str
      - retention_of_cedant: str
      - share_offered: str
    """
    cand: Dict[str, Any] = {}
    if not text:
        return cand

    # ----- INSURED -----
    m = re.search(
        r'(?:Name of the Original Insured|Named Insured|Insured|Assured)[:\s\-]*([A-Z0-9&\-\.,\(\)\'\"/ ]{3,200})',
        text, re.I
    )
    if m:
        cand['insured'] = m.group(1).strip()

    # ----- CEDANT -----
    m = re.search(
        r'(?:Cedant|Cedant Name|Ceding Company|Ceding Insurer|Ceding Party|Insurer)[:\s\-]*([A-Z0-9&\-\.,\(\)\'\"/ ]{3,200})',
        text, re.I
    )
    if m:
        cand['cedant'] = m.group(1).strip()

    # Monetary piece regex: currency may appear before OR after the number
    # currency token = either common symbol OR 3-letter code (ISO-like)
    cur_sym = r'[\$\â‚¬\Â£\Â¥\â‚¹\â‚©\â‚ª\à¸¿Â¢]'
    cur_code = r'[A-Z]{3}'
    money_inner_pat = (
        rf'(?:(?P<cur_before>{cur_code}|{cur_sym})\s*)?'
        r'(?P<num>[-+]?\d[\d,\s]*\.?\d*)'
        rf'(?:\s*(?P<cur_after>{cur_code}|{cur_sym}))?'
    )

    # ----- TSI / Total Sum Insured -----
    tsi_label_pat = (
        r'(?:Total Sum Insured|Total Sum Insured \(TSI\)|Total \(QAR\)|SUM INSURED|TSI|'
        r'Total Limit|Aggregate Limit|Limit of Indemnity)[^\n\r]{0,200}'
        r'([^\n\r]{1,160})'
    )
    m = re.search(tsi_label_pat, text, re.I)
    tsi_span = None
    if m:
        raw = m.group(1).strip()
        cand['total_sum_insured'] = raw
        tsi_span = m.span(1)
        # extract number + currency (before or after) from the raw neighborhood
        im = re.search(money_inner_pat, raw, re.I)
        if im:
            num_raw = im.group('num')
            parsed_num = None
            # handle comma-as-decimal heuristics
            if num_raw.count(',') == 1 and num_raw.count('.') > 0 and num_raw.rfind('.') < num_raw.rfind(','):
                norm = num_raw.replace('.', '').replace(',', '.')
            else:
                norm = num_raw.replace(',', '').replace(' ', '')
            try:
                parsed_num = float(norm)
            except Exception:
                parsed_num = None
            cand['total_sum_insured_float'] = parsed_num
            # prefer a 3-letter code, otherwise symbol; prefer before over after if both present
            cur = None
            cb = im.group('cur_before') or ''
            ca = im.group('cur_after') or ''
            # normalize to uppercase for codes
            if re.fullmatch(cur_code, cb, re.I):
                cur = cb.upper()
            elif re.fullmatch(cur_code, ca, re.I):
                cur = ca.upper()
            elif cb:
                cur = cb
            elif ca:
                cur = ca
            if cur:
                cand['currency'] = cur
        else:
            # fallback: try to find any standalone number in raw
            nm = re.search(r'[-+]?\d[\d,\s]*\.?\d*', raw)
            if nm:
                s = nm.group(0)
                if s.count(',') == 1 and s.count('.') > 0 and s.rfind('.') < s.rfind(','):
                    s_norm = s.replace('.', '').replace(',', '.')
                else:
                    s_norm = s.replace(',', '').replace(' ', '')
                try:
                    cand['total_sum_insured_float'] = float(s_norm)
                except Exception:
                    cand['total_sum_insured_float'] = None
            # try to find any 3-letter code in raw as currency
            cc = re.search(r'\b([A-Z]{3})\b', raw, re.I)
            if cc:
                cand['currency'] = cc.group(1).upper()

    # If TSI wasn't found by label, scan for money-like tokens across entire document (currency before/after)
    if 'total_sum_insured' not in cand:
        candidates = []
        for mm in re.finditer(money_inner_pat, text, re.I):
            num_part = mm.group('num')
            if not num_part:
                continue
            # normalize number
            if num_part.count(',') == 1 and num_part.count('.') > 0 and num_part.rfind('.') < num_part.rfind(','):
                n = num_part.replace('.', '').replace(',', '.')
            else:
                n = num_part.replace(',', '').replace(' ', '')
            try:
                val = float(n)
            except Exception:
                continue
            raw_tok = mm.group(0).strip()
            candidates.append((val, raw_tok, mm.span()))
        if candidates:
            largest = max(candidates, key=lambda t: t[0])
            cand['total_sum_insured'] = largest[1]
            cand['total_sum_insured_float'] = float(largest[0])
            tsi_span = largest[2]
            # try to capture currency for this chosen token
            mm_text = largest[1]
            im = re.search(money_inner_pat, mm_text, re.I)
            if im:
                cb = im.group('cur_before') or ''
                ca = im.group('cur_after') or ''
                if re.fullmatch(cur_code, cb, re.I):
                    cand['currency'] = cb.upper()
                elif re.fullmatch(cur_code, ca, re.I):
                    cand['currency'] = ca.upper()
                elif cb:
                    cand['currency'] = cb
                elif ca:
                    cand['currency'] = ca

    # If still no currency, try neighborhood around tsi_span or whole document
    if 'currency' not in cand:
        search_window = 80
        if tsi_span:
            start = max(0, tsi_span[0] - search_window)
            end = min(len(text), tsi_span[1] + search_window)
            neighborhood = text[start:end]
            cc = re.search(r'\b([A-Z]{3})\b', neighborhood, re.I)
            sym = re.search(cur_sym, neighborhood)
            if cc:
                cand['currency'] = cc.group(1).upper()
            elif sym:
                cand['currency'] = sym.group(0)
        if 'currency' not in cand:
            cc = re.search(r'\b([A-Z]{3})\b', text, re.I)
            sym = re.search(cur_sym, text)
            if cc:
                cand['currency'] = cc.group(1).upper()
            elif sym:
                cand['currency'] = sym.group(0)

    # ----- PERIOD -----
    m = re.search(
        r'(?:Period of Insurance|Period|Insurance Period)[:\s\-]*(From\s+[^\n\r]{3,60}?\s+to\s+[^\n\r]{3,60}|'
        r'\d{1,2}[\/\-]\d{1,2}[\/\-]\d{2,4}\s*to\s*\d{1,2}[\/\-]\d{1,2}[\/\-]\d{2,4}|'
        r'[A-Za-z]{3,9}\s+\d{1,2},?\s*\d{2,4}\s*to\s*[A-Za-z]{3,9}\s+\d{1,2},?\s*\d{2,4})',
        text, re.I
    )
    if m:
        cand['period_of_insurance'] = m.group(1).strip()

    # ----- COUNTRY / RISK LOCATION -----
    m = re.search(r'(?:Risk Location|Country|Territorial Limit|State of)[:\s\-]*([A-Za-z \-\(\)\/]{2,80})',
                  text, re.I)
    if m:
        cand['country'] = m.group(1).strip()

    # ----- RETENTION -----
    m = re.search(
        r'(?:Cedantâ€™s retention|Cedant retention|Retention of Cedant|Cedantâ€™s retention in %|Cedant retention in %|Retention)[:\s\-]{0,5}([0-9]{1,3}\s*%?)',
        text, re.I
    )
    if m:
        cand['retention_of_cedant'] = m.group(1).strip()

    # ----- SHARE OFFERED -----
    m = re.search(r'(?:Share Offered|Offered Share|Share)[:\s\-]*([0-9]{1,3}\s*%?)',
                  text, re.I)
    if m:
        cand['share_offered'] = m.group(1).strip()

    return cand


MASTER_PROMPT: str = """INSTRUCTIONS FOR STRUCTURED EXTRACTION (READ CAREFULLY)

You are given the text of a single email and any OCR-extracted text from its attachments (only regions with confidence â‰¥ 0.7). Use this context to extract reinsurance information and output it as a JSON object.

Fields to extract

insured: Name of the insured party (string).

cedant: Name of the ceding insurer or company (string).

broker: Name of the broker or intermediary (string).

occupation_of_insured: Occupation or role of the insured (string).

main_activities: Main business activities (string).

perils_covered: Covered perils or risks (string).

geographical_limit: Geographical territory or limits (string).

situation_of_risk: Location or situation of the risk (string).

total_sum_insured: Total insured sum (numeric part only, as string).

total_sum_insured_float: Total insured sum as a number (float).

currency: Currency of the insured sum (string).

period_of_insurance: Insurance coverage period (dates or description, string).

excess_deductible: Excess or deductible (string; use "TBD" if unknown).

retention_of_cedant: Cedant's retention (string; use "TBD" if unknown).

possible_maximum_loss: Possible maximum loss (string; use "TBD" if unknown).

cat_exposure: Catastrophe exposure (string; use "TBD" if unknown).

claims_experience: Claims experience details (string; use "TBD" if unknown).

reinsurance_deductions: Reinsurance deductions (string; use "TBD" if unknown).

share_offered: Share offered (string; use "TBD" if unknown).

inward_acceptances: Inward acceptances (string; use "TBD" if unknown).

risk_surveyor_report: Risk surveyor's report (string; use "TBD" if unknown).

premium_rates: Premium rates (string; use "TBD" if unknown).

premium: Premium amount (string; use "TBD" if unknown).

climate_change_risk_factors: Climate change risk factors (string; use "TBD" if unknown).

esg_risk_assessment: ESG risk assessment (string; use "TBD" if unknown).

country: Country of the risk or insured (string).

Extraction Instructions

Output Format: Provide exactly one valid JSON object with all the fields above. Do not include any additional keys or text. Do not output markdown or code fencesâ€”only raw JSON.

Missing Values: If a field's value is not present or is unclear in the text, set it to "TBD" (the string).

Value Types: All values should be strings except for total_sum_insured_float, which must be a numeric value (no quotes). For total_sum_insured, strip any formatting so it's a plain number as a string. Use a standard currency code or symbol for currency.

Semantic Mapping: The text may use different terms or synonyms. Use context to map them to the correct fields. For example:

Words like "Insured", "Assured", or the policyholder's name â†’ insured.

"Cedant", "Ceding insurer", or similar â†’ cedant.

"Broker" or "Intermediary" â†’ broker.

Descriptions of the insured's job or role â†’ occupation_of_insured.

Descriptions of business or operations â†’ main_activities.

Phrases like "covered perils", "risks" â†’ perils_covered.

Territory or country names â†’ geographical_limit or country as appropriate.

Locations or site descriptions â†’ situation_of_risk.

Amounts with currency (e.g. "USD 100000") â†’ total_sum_insured (e.g. "100000"), currency (e.g. "USD"), total_sum_insured_float (e.g. 100000.0).

Date ranges or terms like "from"/"to" â†’ period_of_insurance.

Words like "Excess", "Deductible" â†’ excess_deductible.

"Retention" â†’ retention_of_cedant.

"Maximum loss", "PML" â†’ possible_maximum_loss.

"Cat exposure" or "Catastrophe" â†’ cat_exposure.

"Claims experience" â†’ claims_experience.

"Reinsurance deductions" â†’ reinsurance_deductions.

"Share offered" â†’ share_offered.

"Inward acceptances" â†’ inward_acceptances.

"Surveyor report" â†’ risk_surveyor_report or survey report or risk report.

"Premium rate" â†’ premium_rates.

"Premium" (amount) â†’ premium.

"Climate risk factors" â†’ climate_change_risk_factors.

"ESG" or "sustainability" terms â†’ esg_risk_assessment.

No Extra Text: The model should only output the JSON object. It must not output any explanatory text or additional formatting. Ensure the JSON is syntactically correct and exactly matches the fields above.

---TEXT BLOCK STARTS BELOW---
{LLM_INPUT}
---TEXT BLOCK ENDS ABOVE---
"""




def _truncate_keep_ends(text: str, max_len: int = PROCESSING_CONFIG['max_text_length']) -> str:
    """If text > max_len, keep start and end with a truncation marker in the middle."""
    if not text or len(text) <= max_len:
        return text
    head: str = text[: int(max_len * 0.7)]
    tail: str = text[-int(max_len * 0.3):]
    return head + "\n\n...[TRUNCATED]...\n\n" + tail

def create_llm_context(
        email_data: Dict[str, Any],
        attachment_texts: List[Dict[str, Any]]
) -> Tuple[Dict[str, Any], str]:
    """
    Build LLM context and fill the MASTER_PROMPT with the combined text.
    - Skips attachments with empty text.
    - Returns (context_dict, llm_prompt_string).

    Args:
        email_data: Dictionary containing email metadata and body
        attachment_texts: List of dictionaries with attachment text and metadata

    Returns:
        Tuple of (context dictionary, formatted LLM prompt string)
    """
    # keep only attachments that have non-empty text (already filtered by OCR threshold earlier)
    kept_attachments: List[Dict[str, Any]] = []
    for att in attachment_texts:
        txt: Any = att.get('text', '')
        if isinstance(txt, str) and txt.strip():
            kept_attachments.append(att)

    # Build context (same shape as before but with kept attachments)
    context: Dict[str, Any] = {
        'email': {
            'subject': email_data['metadata'].get('subject', ''),
            'from': email_data['metadata'].get('sender', ''),
            'date': email_data['metadata'].get('date', ''),
            'body': email_data['body'].get('plain_text', '') if 'body' in email_data else ''
        },
        'attachments': kept_attachments,
        'total_attachments': len(kept_attachments),
        'processing_times': {
            'email': 0,
            'attachments': sum(a.get('time', 0) for a in kept_attachments)
        }
    }

    # Build combined human-readable text
    parts: List[str] = []
    parts.append(f"Email Subject: {context['email']['subject']}")
    parts.append(f"From: {context['email']['from']}")
    parts.append(f"Date: {context['email']['date']}")
    parts.append("\nEmail Body:\n" + (context['email']['body'].strip() or ""))

    if kept_attachments:
        parts.append("\nAttachments included for extraction:\n")
        for i, att in enumerate(kept_attachments, start=1):
            header: str = f"--- Attachment {i}: {att['file']} ---"
            page_info: str = f" (page {att.get('page')})" if att.get('page') else ""
            body_text: str = att.get('text', '').strip()

            # if extremely long, keep the head and tail
            max_att_len: int = PROCESSING_CONFIG['max_attachment_length']
            if len(body_text) > max_att_len:
                head_portion: str = body_text[:int(max_att_len * 0.7)]
                tail_portion: str = body_text[-int(max_att_len * 0.3):]
                body_text = f"{head_portion}\n\n...[ATTACHMENT TRUNCATED]...\n\n{tail_portion}"

            parts.append(header + page_info + "\n" + body_text)

    combined_text: str = "\n\n".join(parts).strip()

    # Add candidate hints from combined_text and email body to help mapping
    try:
        candidates: Dict[str, Any] = extract_candidates_from_text(
            combined_text + "\n" + context['email']['body']
        )
    except Exception:
        candidates = {}

    if candidates:
        hint_lines: List[str] = ["\nCandidate Hints (auto-extracted):"]
        for k, v in candidates.items():
            hint_lines.append(f"{k}: {v}")
        combined_text += "\n\n" + "\n".join(hint_lines)

    # Truncate overall combined_text to avoid overlong input
    combined_text = _truncate_keep_ends(combined_text, max_len=11000)

    # Fill the MASTER_PROMPT with the combined text
    llm_prompt: str = MASTER_PROMPT.replace("{LLM_INPUT}", combined_text)

    # Return both the structured context and the full prompt to send to the LLM
    return context, llm_prompt

class Detect_and_Recognize:
    """OCR processor using PaddleOCR for text detection and recognition"""

    def __init__(self) -> None:
        """
        Initialize PaddleOCR with custom parameters

        Args:
            device: 'cpu' or 'gpu' or 'gpu:0,1' for multiple GPUs
            enable_mkldnn: Enable MKL-DNN acceleration on CPU
            cpu_threads: Number of CPU threads
            det_limit_side_len: Max/min side length for detection
            det_limit_type: 'max' or 'min' - how to apply side length limit
            det_thresh: Detection threshold for text pixels
            det_box_thresh: Threshold for text region boxes
            det_unclip_ratio: Box expansion ratio
            rec_batch_size: Batch size for recognition
            use_doc_orientation: Enable document orientation classification
            use_doc_unwarping: Enable document unwarping
            use_textline_orientation: Enable text line orientation
        """
        self.ocr: PaddleOCR = PaddleOCR(
            # doc_orientation_classify_model_name='PP-LCNet_x1_0_doc_ori',
            # doc_orientation_classify_model_dir='PP-LCNet_x1_0_doc_ori',
            # textline_orientation_model_name='PP-LCNet_x1_0_textline_ori',
            # textline_orientation_model_dir='PP-LCNet_x1_0_textline_ori',
            text_detection_model_name='PP-OCRv5_mobile_det',
            text_recognition_model_name='PP-OCRv5_mobile_rec',
            text_detection_model_dir="PP-OCRv5_mobile_det_infer",
            text_recognition_model_dir="PP-OCRv5_mobile_rec_infer",
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False,
            device=OCR_CONFIG['device'],
            cpu_threads=OCR_CONFIG['cpu_threads'],
            enable_mkldnn=OCR_CONFIG['enable_mkldnn'],
            text_det_limit_side_len=OCR_CONFIG['det_limit_side_len'],
            text_det_limit_type='max',
            text_det_thresh=OCR_CONFIG['text_det_thresh'],
            text_det_box_thresh=OCR_CONFIG['text_det_box_thresh'],
            text_det_unclip_ratio=1.5,
            text_recognition_batch_size=OCR_CONFIG['text_recognition_batch_size'],
            # ocr_version="PP-OCRv5",
        )

    def ocr_with_detection_and_recognition(
            self,
            input_path: Union[str, np.ndarray, List[np.ndarray]]
    ) -> Any:
        """
        Perform OCR with both detection and recognition

        Args:
            input_path: Path to image, PDF, directory, or numpy array(s)

        Returns:
            PaddleOCR result object(s) with detection boxes and recognized text
            Format: List of page results, each containing:
                    [[bbox, (text, confidence)], ...]
        """
        return self.ocr.predict(input_path)

    def pdf_to_images(
            self,
            pdf_path: str,
            dpi: int = 300
    ) -> Tuple[List[np.ndarray], List[np.ndarray]]:
        """
        Convert PDF pages to images for OCR processing
        Memory-efficient version with configurable DPI

        Args:
            pdf_path: Path to PDF file
            dpi: Resolution for rendering (200-300 is good balance for CPU)

        Returns:
            tuple: (processed_images, original_images)
                - processed_images: List of numpy arrays ready for OCR
                - original_images: List of original numpy arrays
        """
        doc: fitz.Document = fitz.open(pdf_path)
        processed_images: List[np.ndarray] = []
        original_images: List[np.ndarray] = []

        # ADD TQDM FOR PDF PAGE CONVERSION
        for page_num in tqdm(
                range(len(doc)),
                desc=f"Converting PDF pages ({len(doc)} total)",
                unit="page",
                leave=False
        ):
            page: fitz.Page = doc.load_page(page_num)
            mat: fitz.Matrix = fitz.Matrix(dpi / 72, dpi / 72)
            pix: fitz.Pixmap = page.get_pixmap(matrix=mat)

            img: np.ndarray = np.frombuffer(
                pix.samples,
                dtype=np.uint8
            ).reshape(pix.height, pix.width, pix.n)

            # Convert RGBA to RGB if necessary
            if img.shape[2] == 4:
                img = cv2.cvtColor(img, cv2.COLOR_RGBA2RGB)

            original_images.append(img.copy())
            processed_images.append(img)

            # Clean up pixmap to free memory
            del pix

        doc.close()
        return processed_images, original_images


def check_ollama_installed() -> bool:
    """
    Check if Ollama service is responding and list available models

    Returns:
        bool: True if Ollama is running and accessible, False otherwise
    """
    try:
        with tqdm(total=2, desc="Checking Ollama service", leave=False) as pbar:
            pbar.set_description("Connecting to Ollama")
            models = ollama.list()
            pbar.update(1)

            # ListResponse object has a 'models' attribute
            model_list: List[Any] = models.models if hasattr(models, 'models') else []

            # Extract model names
            pbar.set_description("Listing models")
            model_names: List[str] = []
            for m in model_list:
                if hasattr(m, 'model'):
                    model_names.append(m.model)
                elif hasattr(m, 'name'):
                    model_names.append(m.name)
            pbar.update(1)

        print(f"âœ“ Ollama service is running")
        if model_names:
            print(f"   Available models ({len(model_names)}):")
            for idx, model_name in enumerate(model_names, 1):
                print(f"      {idx}. {model_name}")
        else:
            print(f"   âš ï¸  No models found. Pull a model with: ollama pull <model-name>")

        return True

    except Exception as e:
        print(f"âœ— Cannot connect to Ollama service: {e}")
        print("   Make sure Ollama is running:")
        print("      â€¢ Check if the service is started")
        print("      â€¢ Try running: ollama serve")
        return False

def save_processing_results(
        folder_path: Path,
        context: Dict[str, Any],
        llm_prompt: str,
        llm_result: Dict[str, Any]
) -> None:
    """
    Save all processing results to JSON files with proper timing

    Args:
        folder_path: Directory path where results will be saved
        context: LLM context dictionary containing email and attachment data
        llm_prompt: The formatted master prompt sent to the LLM
        llm_result: Dictionary containing LLM response and metadata

    Returns:
        None
    """
    timestamp: str = datetime.datetime.now().isoformat()

    save_tasks: list = ['context', 'prompt', 'response/error', 'clean_data']
    with tqdm(total=len(save_tasks), desc="Saving results", unit="file", leave=False) as pbar:

        # Save LLM context
        pbar.set_description("Saving LLM context")
        context_path: Path = folder_path / 'llm_context.json'
        context_data: Dict[str, Any] = {
            "context": context,
            "timestamp": timestamp,
            "metadata": {
                "email_subject": context.get('email', {}).get('subject', ''),
                "total_attachments": context.get('total_attachments', 0),
                "processing_times": context.get('processing_times', {})
            }
        }
        with open(context_path, 'w', encoding='utf-8') as f:
            json.dump(context_data, f, indent=2, ensure_ascii=False)
        print(f"âœ“ Saved LLM context: {context_path.name}")
        pbar.update(1)

        # Save master prompt with LLM context
        pbar.set_description("Saving master prompt")
        prompt_path: Path = folder_path / 'master_prompt.json'
        prompt_length: int = len(llm_prompt)
        prompt_data: Dict[str, Any] = {
            "master_prompt": llm_prompt,
            "llm_context": context,
            "timestamp": timestamp,
            "prompt_length": prompt_length,
            "context_summary": {
                "email_subject": context.get('email', {}).get('subject', ''),
                "total_attachments": context.get('total_attachments', 0),
                "processing_times": context.get('processing_times', {})
            }
        }
        with open(prompt_path, 'w', encoding='utf-8') as f:
            json.dump(prompt_data, f, indent=2, ensure_ascii=False)
        print(f"âœ“ Saved master prompt with LLM context: {prompt_path.name}")
        pbar.update(1)

        # Save LLM response with timing
        if llm_result.get("success", False):
            pbar.set_description("Saving LLM response")
            result_path: Path = folder_path / 'llm_response.json'

            # Get metadata safely
            metadata: Dict[str, Any] = llm_result.get("metadata", {})

            result_data: Dict[str, Any] = {
                "success": True,
                "extracted_data": llm_result.get("extracted_data", {}),
                "raw_response": llm_result.get("raw_response", ""),
                "timing": {
                    "generation_time_seconds": llm_result.get("generation_time", 0.0),
                    "tokens_generated": llm_result.get("tokens_generated", 0),
                    "timestamp": timestamp
                },
                "model_info": {
                    "model_used": llm_result.get("model_used", "unknown"),
                    "prompt_length": prompt_length,
                    "temperature": metadata.get("temperature", OLLAMA_CONFIG.get("temperature", 0.1)),
                    "max_tokens": metadata.get("max_tokens", OLLAMA_CONFIG.get("max_tokens", 2000))
                }
            }

            # Add warning if present
            warning: Optional[str] = llm_result.get("warning")
            if warning:
                result_data["warning"] = warning

            with open(result_path, 'w', encoding='utf-8') as f:
                json.dump(result_data, f, indent=2, ensure_ascii=False)
            print(f"âœ“ Saved LLM response with timing: {result_path.name}")
            pbar.update(1)

            # Also save a clean version with just the extracted data
            pbar.set_description("Saving clean extraction")
            clean_path: Path = folder_path / 'extracted_reinsurance_data.json'
            extracted_data: Dict[str, Any] = llm_result.get("extracted_data", {})
            with open(clean_path, 'w', encoding='utf-8') as f:
                json.dump(extracted_data, f, indent=2, ensure_ascii=False)
            print(f"âœ“ Saved clean extraction: {clean_path.name}")
            pbar.update(1)

        else:
            pbar.set_description("Saving error response")
            error_path: Path = folder_path / 'llm_error.json'
            error_data: Dict[str, Any] = {
                "success": False,
                "error": llm_result.get("error", "Unknown error"),
                "raw_response": llm_result.get("raw_response", ""),
                "timing": {
                    "generation_time_seconds": llm_result.get("generation_time", 0.0),
                    "timestamp": timestamp
                }
            }
            with open(error_path, 'w', encoding='utf-8') as f:
                json.dump(error_data, f, indent=2, ensure_ascii=False)
            print(f"âœ— Saved error with timing: {error_path.name}")
            pbar.update(2)  # Skip the clean data step


def validate_processing_results(folder_path: Path) -> Dict[str, Any]:
    """
    Returns:
      {
        'validation': { 'context': True, ... },
        'missing_files': ['master_prompt.json'],
        'invalid_json': ['llm_response.json']
      }
    """
    validation: Dict[str, bool] = {}
    missing: List[str] = []
    invalid: List[str] = []
    expected_files: Dict[str, str] = {
        'context': 'llm_context.json',
        'prompt': 'master_prompt.json',
        'response': 'llm_response.json',
        'error': 'llm_error.json',
        'extraction': 'extracted_reinsurance_data.json'
    }

    for file_type, filename in expected_files.items():
        file_path = folder_path / filename
        if not file_path.exists():
            validation[file_type] = False
            missing.append(filename)
            continue
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                json.load(f)
            validation[file_type] = True
        except json.JSONDecodeError:
            validation[file_type] = False
            invalid.append(filename)

    return {'validation': validation, 'missing_files': missing, 'invalid_json': invalid}

def load_processing_results(folder_path: Path, strict: bool = False) -> Optional[Dict[str, Any]]:
    """
    Load results; if strict=True raise exceptions on missing/invalid files,
    otherwise return a dict with whatever could be read plus a 'diagnostics' key.
    """
    results: Dict[str, Any] = {}
    diag = {'missing': [], 'invalid': []}

    def try_load(name: str, filename: str):
        p = folder_path / filename
        if not p.exists():
            diag['missing'].append(filename)
            return None
        try:
            with open(p, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception:
            diag['invalid'].append(filename)
            return None

    results['context'] = try_load('context', 'llm_context.json')
    results['prompt'] = try_load('prompt', 'master_prompt.json')
    # prefer response, fallback to error
    r = try_load('response', 'llm_response.json')
    if r is None:
        r = try_load('error', 'llm_error.json')
    results['llm_result'] = r
    results['extracted_data'] = try_load('extraction', 'extracted_reinsurance_data.json')
    results['diagnostics'] = diag

    if strict and (diag['missing'] or diag['invalid']):
        raise RuntimeError(f"Load failed: {diag}")

    return results

def process_with_llm_integration(root_folder: str, model_name: str = None,
                                 skip_processed: bool = True,
                                 strict_load_on_error: bool = False):
    """
    Complete pipeline: Email processing + OCR + LLM extraction

    Args:
        root_folder: Path to folder containing email subfolders
        model_name: Local Ollama model name to use
        skip_processed: If True, skip subfolders that already have valid saved results
        strict_load_on_error: If True, call load_processing_results(..., strict=True) when validation fails to raise
                              and surface problems; otherwise keep diagnostics in results.
    """
    check_ollama_installed()
    print("ðŸš€ Initializing components...")

    llm_processor = LocalLLMProcessor(model_name)
    if not llm_processor.test_connection():
        print("âŒ Cannot proceed without local Ollama connection")
        return

    print("Initializing OCR engine...")
    ocr_engine = Detect_and_Recognize()

    root_path = Path(root_folder)
    subfolders = [f for f in root_path.iterdir() if f.is_dir()]

    print(f"ðŸ“ Found {len(subfolders)} email folders to process")

    for subfolder in tqdm(sorted(subfolders), desc="Processing emails", unit="email"):
        print(f"\n{'=' * 60}")
        print(f"ðŸ“§ Processing: {subfolder.name}")
        print(f"{'=' * 60}")

        try:
            # --- Early skip: if folder already has valid saved results, skip (useful for resume)
            if skip_processed:
                vres = validate_processing_results(subfolder)
                validation = vres.get('validation', {})
                # consider 'complete' when either 'response' OR 'error' exists AND 'context' and 'prompt' exist
                if validation.get('context') and validation.get('prompt') and (validation.get('response') or validation.get('error')):
                    print("â„¹ï¸  Valid saved results present â€” skipping (use skip_processed=False to force).")
                    continue

            print("ðŸ“¨ Extracting email data...")
            processor = EmailProcessor(subfolder)
            email_data = processor.get_complete_data()

            if not email_data:
                print("âŒ No email data found")
                continue

            print(f"âœ“ Subject: {email_data['metadata']['subject']}")
            print(f"âœ“ From: {email_data['metadata']['sender']}")
            print(f"âœ“ Attachments: {len(email_data['saved_attachments'])}")

            print("ðŸ“„ Processing attachments...")
            attachment_texts = processor.process_attachments_with_ocr(
                email_data, ocr_engine, save_visuals=PROCESSING_CONFIG['save_visuals']
            )

            print("ðŸ§  Creating LLM context...")
            context, llm_prompt = create_llm_context(email_data, attachment_texts)

            print("ðŸ¤– Running LLM extraction...")
            x = time.time()
            llm_result = llm_processor.extract_structured_data(llm_prompt)
            print(f'Response generation took: {time.time() - x:.2f} seconds')

            print("ðŸ’¾ Saving results...")
            save_processing_results(subfolder, context, llm_prompt, llm_result)

            # --- Immediately validate saved files and surface diagnostics
            vres = validate_processing_results(subfolder)
            if not (vres['validation'].get('context') and vres['validation'].get('prompt') and (vres['validation'].get('response') or vres['validation'].get('error'))):
                print("âš ï¸  Post-save validation FAILED for folder:", subfolder.name)
                print("    missing_files:", vres.get('missing_files'))
                print("    invalid_json:", vres.get('invalid_json'))

                # optionally auto-load diagnostics (strict or not)
                try:
                    load_results = load_processing_results(subfolder, strict=strict_load_on_error)
                    print("    Diagnostics loaded; see load_results['diagnostics']")
                except Exception as ex:
                    print("    load_processing_results raised:", ex)
            else:
                # saved and validated successfully
                if llm_result.get("success"):
                    print(f"âœ… Success! Generated {llm_result.get('tokens_generated', '?')} tokens in {llm_result.get('generation_time', 0):.2f}s")
                    extracted = llm_result.get("extracted_data", {})
                    print(f"ðŸ“Š Extracted: {extracted.get('insured', 'N/A')} | {extracted.get('cedant', 'N/A')} | {extracted.get('total_sum_insured', 'N/A')}")
                else:
                    print(f"âŒ LLM extraction failed: {llm_result.get('error')}")

        except Exception as e:
            print(f"âŒ Error processing {subfolder.name}: {e}")
            traceback.print_exc()

    print(f"\nðŸŽ‰ Processing complete!")

if __name__ == "__main__":
    root_folder = "wednesday evening test data"
    model_name = OLLAMA_CONFIG['model_name']
    process_with_llm_integration(root_folder, model_name)
