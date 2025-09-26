import re
import logging
from typing import Dict, List, Optional, Tuple, Any
from dataclasses import dataclass
from pathlib import Path
import pandas as pd
from datetime import datetime
import json

# Email processing
try:
    import extract_msg
    import email
    from email.mime.multipart import MIMEMultipart
    from email.mime.text import MIMEText

    HAS_EMAIL_LIBS = True
except ImportError:
    HAS_EMAIL_LIBS = False
    print("Email libraries not available. Install with: pip install extract-msg")

# Document processing
try:
    import fitz  # PyMuPDF for PDF processing
    import openpyxl
    from docx import Document
    from pptx import Presentation
    HAS_DOC_LIBS = True
except ImportError:
    HAS_DOC_LIBS = False
    print("Document libraries not available")

# OCR (only loaded when needed)
PADDLEOCR_AVAILABLE = False
try:
    from paddleocr import PaddleOCR
    import cv2
    from PIL import Image
    PADDLEOCR_AVAILABLE = True
except ImportError:
    pass

logger = logging.getLogger(__name__)


class OptimizedEmailProcessor:
    """Fast email processing with header and body analysis."""

    def __init__(self):
        """Initialize email processor."""
        if not HAS_EMAIL_LIBS:
            raise ImportError("Email processing libraries not available")

    def process_msg_file(self, file_path: str) -> Tuple[str, Dict[str, str]]:
        """
        Process MSG file and extract email content with metadata.

        Returns:
            Tuple of (full_text, email_metadata)
        """
        try:
            msg = extract_msg.Message(file_path)

            metadata = {
                'subject': self._fix_encoding(getattr(msg, 'subject', '')),
                'sender': self._fix_encoding(getattr(msg, 'sender', '')),
                'to': self._fix_encoding(getattr(msg, 'to', '')),
                'cc': self._fix_encoding(getattr(msg, 'cc', '')),
                'date': str(getattr(msg, 'date', '')),
            }

            # Extract body
            body = getattr(msg, 'body', '')

            # Combine all text
            full_text = f"""
                Subject: {metadata['subject']}
                From: {metadata['sender']}
                To: {metadata['to']}
                Date: {metadata['date']}
                
                Body:
                {body}
            """

            # Process attachments if they exist
            attachments_text = ""
            if hasattr(msg, 'attachments') and msg.attachments:
                attachments_text += "\nAttachments found:\n"
                for attachment in msg.attachments:
                    if hasattr(attachment, 'longFilename'):
                        attachments_text += f"- {attachment.longFilename}\n"

            return full_text + attachments_text, metadata

        except Exception as e:
            logger.error(f"Error processing MSG file {file_path}: {e}")
            return "", {}

    def process_msg_with_attachments(self, file_path: str, doc_processor) -> Tuple[str, Dict[str, str]]:
        """
        Process MSG file with attachments and combine all content for unified extraction.

        Returns:
            Tuple of (combined_text, email_metadata)
        """
        try:
            msg = extract_msg.Message(file_path)

            # Extract email metadata
            metadata = {
                'subject': self._fix_encoding(getattr(msg, 'subject', '')),
                'sender': self._fix_encoding(getattr(msg, 'sender', '')),
                'to': getattr(msg, 'to', ''),
                'cc': getattr(msg, 'cc', ''),
                'date': str(getattr(msg, 'date', '')),
            }

            # Start with email body
            body = getattr(msg, 'body', '')
            combined_text = f"""
                EMAIL CONTENT:
                Subject: {metadata['subject']}
                From: {metadata['sender']}
                Date: {metadata['date']}

                Body:
                {body}
            """

            # Process and combine attachments
            attachment_texts = []
            if hasattr(msg, 'attachments') and msg.attachments:
                import tempfile

                for i, attachment in enumerate(msg.attachments):
                    try:
                        if hasattr(attachment, 'longFilename') and hasattr(attachment, 'data'):
                            filename = attachment.longFilename or f"attachment_{i}"

                            # Create temp file for processing
                            with tempfile.NamedTemporaryFile(suffix=Path(filename).suffix, delete=False) as temp_file:
                                temp_file.write(attachment.data)
                                temp_path = temp_file.name

                            # Process attachment based on type
                            attachment_text = self._process_single_attachment(temp_path, filename, doc_processor)

                            if attachment_text and attachment_text.strip():
                                attachment_texts.append(f"\n\nATTACHMENT: {filename}\n{'-' * 50}\n{attachment_text}")

                            # Clean up
                            Path(temp_path).unlink(missing_ok=True)

                    except Exception as e:
                        logger.warning(f"Failed to process attachment {i} in {file_path}: {e}")
                        continue

            # Combine email and all attachments
            if attachment_texts:
                combined_text += "\n\nATTACHMENTS CONTENT:\n" + "\n".join(attachment_texts)

            return combined_text, metadata

        except Exception as e:
            logger.error(f"Error processing MSG with attachments {file_path}: {e}")
            return "", {}

    def _process_single_attachment(self, temp_path: str, filename: str, doc_processor) -> str:
        """Process a single attachment and return its text content."""
        file_ext = Path(filename).suffix.lower()

        try:
            if file_ext == '.pdf':
                text, _ = doc_processor.process_pdf_smart(temp_path)
                return text
            elif file_ext == '.docx':
                return doc_processor.process_docx(temp_path)
            elif file_ext in ['.xlsx', '.xls']:
                return doc_processor.process_excel_smart(temp_path)
            elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
                return self._process_attachment_image_ocr(temp_path, doc_processor)
            else:
                logger.info(f"Unsupported attachment type: {file_ext}")
                return ""
        except Exception as e:
            logger.error(f"Error processing attachment {filename}: {e}")
            return ""

    def _process_attachment_image_ocr(self, file_path: str, doc_processor) -> str:
        """Process image attachments with OCR."""
        if not PADDLEOCR_AVAILABLE:
            return ""

        try:
            ocr = doc_processor._get_ocr()
            if ocr is None:
                return ""

            result = ocr.ocr(file_path)
            if result and result[0]:
                return '\n'.join([line[1][0] for line in result[0] if line[1][0].strip()])
            return ""
        except Exception as e:
            logger.error(f"OCR failed for {file_path}: {e}")
            return ""

    def _fix_encoding(self, text: str) -> str:
        """Fix common encoding issues."""
        if not text:
            return text

        encoding_fixes = {
            'â€"': '–', 'â€™': "'", 'â€œ': '"', 'â€': '"', 'â€¢': '•',
        }

        for wrong, correct in encoding_fixes.items():
            text = text.replace(wrong, correct)
        return text

    def _process_attachment_image(self, file_path: str, doc_processor) -> str:
        """Process image attachments using OCR."""
        try:
            ocr = doc_processor._get_ocr()
            if ocr is None:
                return ""

            result = ocr.ocr(file_path)
            if result and result[0]:
                return '\n'.join([line[1][0] for line in result[0] if line[1][0].strip()])
            return ""
        except Exception as e:
            logger.error(f"OCR processing failed for attachment {file_path}: {e}")
            return ""

    def quick_check_email_relevance(self, text: str, metadata: Dict[str, str]) -> bool:
        """email relevance check with comprehensive keywords."""
        keywords = [
            # Core reinsurance terms
            'facultative', 'reinsurance', 'reinsuarance', 'reinsurer', 'facultative reinsurance',
            'cedant', 'ceding', 'ceding company', 'broker', 'reinsurance broker',
            'insured', 'assured', 'policyholder', 'risk holder',

            # Financial terms
            'tsi', 'sum insured', 'total sum insured', 'insured value', 'coverage amount',
            'premium', 'rate', 'pricing', 'quote', 'quotation',
            'liability', 'exposure', 'limit', 'policy limit', 'coverage limit',
            'excess', 'deductible', 'retention', 'share', 'percentage',

            # Risk terms
            'peril', 'hazard', 'risk', 'coverage', 'protection', 'claim', 'loss',
            'pml', 'possible maximum loss', 'maximum probable loss',
            'catastrophe', 'cat exposure', 'natural disaster',

            # Document types
            'placement', 'submission', 'slip', 'insurance slip', 'reinsuarance slip',
            'proposal', 'offering', 'quote slip', 'placing slip',
            'facultative submission', 'reinsurance submission',
            'r.i', 'ri', 'insurance proposal',

            # Business terms
            'underwriter', 'underwriting', 'acceptance', 'decline',
            'terms', 'conditions', 'warranty', 'exclusion',
            'renewal', 'extension', 'amendment',

            # Geography/Location
            'worldwide', 'territory', 'geographical', 'location', 'situated',

            # Industry specific
            'fire', 'flood', 'earthquake', 'storm', 'explosion', 'terrorism',
            'machinery breakdown', 'business interruption', 'marine', 'aviation',
            'property', 'casualty', 'motor', 'energy', 'construction'
        ]

        text_lower = text.lower()
        subject_lower = metadata.get('subject', '').lower()
        sender_lower = metadata.get('sender', '').lower()

        # Check subject (highest weight)
        subject_matches = sum(1 for keyword in keywords if keyword in subject_lower)

        # Check sender for broker/insurance companies
        sender_indicators = ['broker', 'insurance', 'reinsurance', 'underwriter', 'risk']
        sender_matches = sum(1 for indicator in sender_indicators if indicator in sender_lower)

        # Check body text
        text_matches = sum(1 for keyword in keywords if keyword in text_lower)

        # More sophisticated scoring
        relevance_score = (subject_matches * 3) + (sender_matches * 2) + text_matches

        # Email is relevant if score >= 5 or subject has 2+ matches
        return relevance_score >= 5 or subject_matches >= 2


class SmartDocumentProcessor:
    """Intelligent document processor with selective OCR."""

    def __init__(self):
        """Initialize document processor."""
        self.ocr = None  # Lazy load OCR only when needed

    def _get_ocr(self):
        """Lazy load OCR when needed."""
        if self.ocr is None and PADDLEOCR_AVAILABLE:
            self.ocr = PaddleOCR(
                text_detection_model_name='PP-OCRv4_mobile_det',
                text_recognition_model_name='PP-OCRv4_mobile_rec',
                text_detection_model_dir="PP-OCRv4_mobile_det_infer",
                text_recognition_model_dir="PP-OCRv4_mobile_rec_infer",
                use_doc_orientation_classify=False,
                use_doc_unwarping=False,
                use_textline_orientation=False,
                lang='en',
                ocr_version="PP-OCRv4",
                device='cpu',
            )
            # self.ocr = PaddleOCR(
            #     text_detection_model_name='PP-OCRv5_mobile_det',
            #     text_recognition_model_name='PP-OCRv5_server_rec',
            #     text_detection_model_dir="PP-OCRv5_mobile_det_infer",
            #     text_recognition_model_dir="PP-OCRv5_server_rec_infer",
            #     use_doc_orientation_classify=False,
            #     use_doc_unwarping=False,
            #     use_textline_orientation=False,
            #     lang='en',
            #     ocr_version="PP-OCRv5",
            #     device='cpu',
            # )
        return self.ocr

    def process_pdf_smart(self, file_path: str, max_pages: int = 20) -> Tuple[str, bool]:
        """
        Smart PDF processing with fallback to OCR for scanned documents.

        Returns:
            Tuple of (extracted_text, used_ocr)
        """
        try:
            doc = fitz.open(file_path)
            text = ""
            pages_to_process = min(len(doc), max_pages)

            # Try text extraction first
            for page_num in range(pages_to_process):
                page = doc[page_num]
                page_text = page.get_text().strip()
                text += page_text + "\n"

            doc.close()

            # Check if extracted text is meaningful
            if len(text.strip()) > 50 and self._has_meaningful_text(text):
                return text, False  # Successfully extracted without OCR

            # If text extraction failed, try OCR on first few pages
            logger.info(f"Text extraction failed for {file_path}, attempting OCR...")
            return self._ocr_pdf_fallback(file_path, min(3, pages_to_process)), True

        except Exception as e:
            logger.error(f"PDF processing failed for {file_path}: {e}")
            return "", False

    def _has_meaningful_text(self, text: str) -> bool:
        """Check if extracted text appears to be meaningful (not just symbols/garbage)."""
        if len(text.strip()) < 50:
            return False

        # Count alphabetic characters vs total characters
        alpha_count = sum(1 for c in text if c.isalpha())
        total_chars = len(text.replace(' ', '').replace('\n', ''))

        if total_chars == 0:
            return False

        alpha_ratio = alpha_count / total_chars
        return alpha_ratio > 0.7  # At least 70% alphabetic characters

    def _ocr_pdf_fallback(self, file_path: str, max_pages: int = 3) -> str:
        """Use OCR as fallback for scanned PDFs."""
        ocr = self._get_ocr()
        if ocr is None:
            logger.warning("OCR not available for scanned PDF processing")
            return ""

        try:
            doc = fitz.open(file_path)
            text = ""

            for page_num in range(min(len(doc), max_pages)):
                page = doc[page_num]
                pix = page.get_pixmap()
                img_data = pix.tobytes("png")

                # Save temporary image
                temp_img_path = f"temp_page_{page_num}.png"
                with open(temp_img_path, "wb") as f:
                    f.write(img_data)

                # OCR the image
                result = ocr.ocr(temp_img_path, cls=True)

                if result:
                    page_text = '\n'.join([line[1][0] for line in result[0] if line[1][0].strip()])
                    text += page_text + "\n"

                # Clean up temp file
                Path(temp_img_path).unlink(missing_ok=True)

            doc.close()
            return text

        except Exception as e:
            logger.error(f"OCR fallback failed for {file_path}: {e}")
            return ""

    def process_docx(self, file_path: str) -> str:
        """Extract text from DOCX files."""
        try:
            doc = Document(file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])

            # Also extract from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += "\n" + cell.text

            return text
        except Exception as e:
            logger.error(f"DOCX processing failed for {file_path}: {e}")
            return ""

    def process_excel_smart(self, file_path: str, max_sheets: int = 5) -> str:
        """Smart Excel processing focusing on relevant sheets."""
        try:
            # Get sheet names first
            xl_file = pd.ExcelFile(file_path)
            sheet_names = xl_file.sheet_names

            # Prioritize sheets that might contain facultative data
            priority_keywords = ['facultative', 'submission', 'quote', 'placement', 'risk', 'tsi', 'sum insured', 'total sum insured']
            prioritized_sheets = []
            other_sheets = []

            for sheet in sheet_names:
                sheet_lower = sheet.lower()
                if any(keyword in sheet_lower for keyword in priority_keywords):
                    prioritized_sheets.append(sheet)
                else:
                    other_sheets.append(sheet)

            # Process prioritized sheets first
            sheets_to_process = (prioritized_sheets + other_sheets)[:max_sheets]

            text = ""
            for sheet_name in sheets_to_process:
                try:
                    df = pd.read_excel(file_path, sheet_name=sheet_name, nrows=1000)
                    text += f"\n--- Sheet: {sheet_name} ---\n"
                    text += df.to_string(index=False, na_rep='') + "\n"
                except Exception as e:
                    logger.warning(f"Error reading sheet {sheet_name}: {e}")
                    continue

            return text
        except Exception as e:
            logger.error(f"Excel processing failed for {file_path}: {e}")
            return ""


class EnhancedInformationExtractor:
    def __init__(self):
        """Initialize with enhanced extraction patterns."""
        self.patterns = {

            # 'insured': [
            #     # More specific patterns to avoid grabbing policy text
            #     r'(?:name\s*of\s*insured|insured\s*name|company\s*name)[\s:]+([A-Za-z][A-Za-z0-9\s\.,&\-()]*?(?:\s+(?:S\.A\.E\.|Ltd\.|Limited|Inc\.|Corp\.|Company|Co\.|Pvt\.|Private|LLC|S\.A\.|GmbH))+)',
            #     r'(?:insured|assured|policyholder)[\s:]+([A-Z][A-Za-z0-9\s\.,&\-()]*?(?:\s+(?:S\.A\.E\.|Ltd\.|Limited|Inc\.|Corp\.|Company|Co\.|Pvt\.|Private|LLC|S\.A\.|GmbH))+)',
            #     r'(?:client|customer|entity)[\s:]+([A-Z][A-Za-z0-9\s\.,&\-()]*?(?:\s+(?:S\.A\.E\.|Ltd\.|Limited|Inc\.|Corp\.|Company|Co\.|Pvt\.|Private|LLC|S\.A\.|GmbH))+)',
            #     # Fallback for cases without clear labels
            #     r'\b([A-Z][A-Za-z0-9\s&\-()]{10,100}(?:\s+(?:S\.A\.E\.|Ltd\.|Limited|Inc\.|Corp\.|Company|Co\.|Pvt\.|Private|LLC|S\.A\.|GmbH))+)\b',
            # ],

            'insured': [
                # Direct patterns for company names
                r'(?:name\s*of\s*insured|insured\s*name|insured|assured|policyholder)[\s:]+([A-Za-z][A-Za-z0-9\s\.,&\-()]*?(?:Ltd\.|Limited|Inc\.|Corp\.|Company|Co\.|Pvt\.|Private|S\.A\.E\.|Assurance).*?)(?=\s*(?:\n|accepts|excludes|warranties|private\s*limited))',
                r'(?:company|entity|client)[\s:]+([A-Za-z][A-Za-z0-9\s\.,&\-()]*?(?:Ltd\.|Limited|Inc\.|Corp\.|Company|Co\.|Pvt\.|Private|S\.A\.E\.|Assurance).*?)(?=\s*(?:\n|accepts|excludes))',
                # Catch standalone company names with indicators
                r'\b([A-Za-z][A-Za-z0-9\s&\-()]{5,80}(?:\s+(?:Ltd\.|Limited|Inc\.|Corp\.|Company|Co\.|Pvt\.|Private|S\.A\.E\.|Assurance)))\b',
            ],

            'cedant': [
                r'(?:cedant|ceding\s*company|ceding\s*insurer|direct\s*insurer)[\s:]+([A-Za-z][A-Za-z0-9\s\.,&\-()]*?(?:Ltd\.|Limited|Inc\.|Corp\.|Company|Co\.|Pvt\.|Private|S\.A\.E\.|Assurance).*?)(?=\s*(?:\n|accepts))',
                # Look for insurance company names
                r'\b([A-Za-z][A-Za-z0-9\s&\-()]{5,80}(?:\s+(?:Insurance|Assurance))(?:\s+(?:Ltd\.|Limited|Inc\.|Corp\.|Company|Co\.|Pvt\.|Private|S\.A\.E\.))?)\b',
            ],
            'broker': [
                r'(?:broker|brokerage|intermediary)[\s:]*["\']?([A-Za-z][A-Za-z0-9\s\.,&\-()]{4,150})["\']?',
                r'(?:placing\s*broker|reinsurance\s*broker)[\s:]*["\']?([A-Za-z][A-Za-z0-9\s\.,&\-()]{4,150})["\']?',
                r'(?:via|through|brokered\s*by)[\s:]*["\']?([A-Za-z][A-Za-z0-9\s\.,&\-()]{4,150})["\']?',
            ],
            'total_sum_insured': [
                # More precise TSI patterns with currency capture
                r'(?:total\s*sum\s*insured|t\.?s\.?i\.?|sum\s*insured)[\s:]*([A-Z]{3}[\s]*[0-9,.\s]+(?:million|billion|thousand|m|b|k)?)',
                r'(?:insured\s*value|coverage\s*amount|insured\s*amount)[\s:]*([A-Z]{3}[\s]*[0-9,.\s]+(?:million|billion|thousand|m|b|k)?)',
                r'(?:limit|policy\s*limit|coverage\s*limit)[\s:]*([A-Z]{3}[\s]*[0-9,.\s]+(?:million|billion|thousand|m|b|k)?)',
                # Pattern for amounts with currency symbols
                r'((?:USD|US\$|\$|EUR|€|GBP|£|KES|KSh|AED|SAR|ZAR)[\s]*[0-9,.\s]+(?:million|billion|thousand|m|b|k)?)',
                r'([0-9,.\s]+(?:million|billion|thousand|m|b|k)?[\s]*(?:USD|US\$|\$|EUR|€|GBP|£|KES|KSh|AED|SAR|ZAR))',
            ],
            # 'risk_surveyor_report': [
            #     r'(?:surveyor|survey|inspection|risk\s*survey)[\s]*(?:report)?[\s:]*([^,\n\r;]{20,500})',
            #     r'(?:technical\s*survey|engineering\s*report|site\s*inspection)[\s:]*([^,\n\r;]{20,500})',
            #     r'(?:recommendations|findings|observations)[\s:]*([^,\n\r;]{20,500})',
            #     r'(?:fire\s*protection|safety\s*measures|security\s*arrangements)[\s:]*([^,\n\r;]{20,500})',
            #     r'(?:housekeeping|maintenance|construction\s*quality)[\s:]*([^,\n\r;]{20,500})',
            # ],

            'risk_surveyor_report': [
                # Pattern to capture report names/titles
                r'(?:surveyor[\'s]*\s*report|survey\s*report|inspection\s*report|technical\s*survey|risk\s*survey)[\s:]*["\']?([A-Za-z][A-Za-z0-9\s\-_\.\/()]{5,150})["\']?',
                r'(?:report\s*by|surveyed\s*by|inspected\s*by)[\s:]*["\']?([A-Za-z][A-Za-z0-9\s&\-_\.\/()]{5,100})["\']?',
                # Capture report file names
                r'(?:survey|inspection|technical)[\s_\-]*(?:report|assessment)[\s]*[:\-]?[\s]*([A-Za-z0-9][A-Za-z0-9\s\-_\.\/()]{5,80}\.(?:pdf|doc|docx|txt))',
                # Capture surveyor company names
                r'(?:risk\s*management|engineering\s*consultant|surveyor|inspection\s*company)[\s:]*["\']?([A-Za-z][A-Za-z0-9\s&\-_\.\/()]{10,100})["\']?',
                # Generic report reference patterns
                r'(?:as\s*per|according\s*to|reference)[\s]*(?:the\s*)?(?:survey|inspection|technical\s*report)[\s:]*["\']?([A-Za-z0-9][A-Za-z0-9\s\-_\.\/()]{5,80})["\']?',
                # Date-based report references
                r'(?:survey|inspection|report)[\s]*(?:dated|on|of)[\s]*(\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4})',
            ],
            'inward_acceptances': [
                r'(?:inward\s*acceptance|inward\s*business|treaty\s*acceptance)[\s:]*([^,\n\r;]{10,200})',
                r'(?:accepted\s*from|received\s*from|ceded\s*by)[\s:]*([^,\n\r;]{10,200})',
                r'(?:participating\s*companies|co-insurers|consortium)[\s:]*([^,\n\r;]{10,200})',
            ],

            'share_offered': [
                r'(?:share\s*offered|percentage\s*offered|participation\s*sought)[\s:]*([0-9]+(?:\.[0-9]+)?%?)',
                r'(?:seeking|looking\s*for|require)[\s:]*([0-9]+(?:\.[0-9]+)?%?)[\s]*(?:share|participation)',
                r'(?:quota\s*share|surplus|percentage)[\s:]*([0-9]+(?:\.[0-9]+)?%?)',
                r'([0-9]+(?:\.[0-9]+)?%?)[\s]*(?:of|share|participation|quota)',
                r'(?:line\s*capacity|our\s*share|participation)[\s:]*([0-9]+(?:\.[0-9]+)?%?)',
            ],
            'claims_experience': [
                # Enhanced patterns for comprehensive claims data
                # Year-wise claims with amounts
                r'(20[0-9]{2}[\s]*[-:][\s]*(?:USD|EUR|GBP|INR|KES|AED|SAR|CAD|AUD|\$|â‚¬|Â£)?\s*[0-9,.\s]*(?:million|thousand|m|k)?[^,\n\r;\.]{0,100})',

                # Detailed loss ratios and frequencies
                r'((?:loss|combined|claims?)\s*ratio[\s:]*[0-9]+(?:\.[0-9]+)?%?[^,\n\r;\.]{0,150})',
                r'((?:frequency|severity)\s*(?:of\s*)?(?:claims?|losses?)[\s:]*[^,\n\r;\.]{10,200})',

                # Claims-free periods with specificity
                r'((?:[0-9]+\s*years?|last\s*[0-9]+\s*years?|since\s*[0-9]{4})\s*(?:claims?\s*free|no\s*claims?|nil\s*losses?|claims?\s*free\s*period)[^,\n\r;\.]{0,100})',

                # Major/minor claims classification
                r'((?:major|minor|significant|small)\s*(?:claims?|losses?)[\s:]*[^,\n\r;\.]{10,200})',

                # Well below/above average references
                r'((?:well\s*below|above|below)\s*(?:average|industry|market)[\s]*(?:claims?|losses?|experience)[^,\n\r;\.]{0,150})',

                # Specific incident descriptions
                r'((?:incident|claim|loss)\s*(?:related\s*to|due\s*to|from)[\s]*[^,\n\r;\.]{15,300})',

                # Recovery and settlement information
                r'((?:recovered|settlement|paid|outstanding)[\s:]*(?:USD|EUR|GBP|INR|KES|AED|SAR|CAD|AUD|\$|â‚¬|Â£)?\s*[0-9,.\s]*(?:million|thousand|m|k)?[^,\n\r;\.]{0,100})',

                # Comprehensive 3-year summaries
                r'((?:last|past|previous)\s*(?:3|three)\s*years?[\s:]*[^\.]{30,500})',

                # Attachment/excess information in claims context
                r'((?:excess|deductible|attachment)\s*(?:of|point)[\s:]*(?:USD|EUR|GBP|INR|KES|AED|SAR|CAD|AUD|\$|â‚¬|Â£)?\s*[0-9,.\s]*(?:million|thousand|m|k)?[^,\n\r;\.]{0,150})',
            ],
            'period_of_insurance': [
                r'(?:period\s*of\s*insurance|insurance\s*period|policy\s*period|coverage\s*period)[\s:]*([^,\n\r;]{10,100})',
                r'(?:effective\s*from|commencing\s*from|valid\s*from)[\s:]*(\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4})[\s]*(?:to|until|through|expiring)[\s]*(\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4})',
                r'(\d{1,2}\s*(?:months?|years?))[\s]*(?:from|commencing|starting|effective)',
                r'(?:12\s*months|24\s*months|36\s*months|one\s*year|two\s*years)[\s]*(?:from|commencing|period)',
                r'(?:annual|yearly|12\s*month)[\s]*(?:policy|coverage|period)',
            ],
            'excess_deductible': [
                r'(?:excess|deductible)[\s:]*([A-Z]{3}[\s]*[0-9,.\s]+(?:million|thousand|m|k)?)',
                r'(?:first\s*loss|retained\s*by\s*insured)[\s:]*([A-Z]{3}[\s]*[0-9,.\s]+(?:million|thousand|m|k)?)',
                r'(?:subject\s*to\s*excess\s*of)[\s:]*([A-Z]{3}[\s]*[0-9,.\s]+(?:million|thousand|m|k)?)',
            ],
            'premium_rates': [
                r'(?:rate|premium\s*rate)[\s:]*([0-9]+(?:\.[0-9]+)?%?(?:\s*per\s*(?:mille|mill))?)',
                r'(?:pricing\s*at|priced\s*at)[\s:]*([0-9]+(?:\.[0-9]+)?%?(?:\s*per\s*(?:mille|mill))?)',
                r'([0-9]+(?:\.[0-9]+)?)[\s]*(?:%|per\s*mill?e?)[\s]*(?:rate|pricing)',
            ],
            'premium': [
                r'(?:premium)[\s:]*([A-Z]{3}[\s]*[0-9,.\s]+(?:million|thousand|m|k)?)',
                r'(?:total\s*premium|gross\s*premium)[\s:]*([A-Z]{3}[\s]*[0-9,.\s]+(?:million|thousand|m|k)?)',
            ],
            'perils_covered': [
                r'(?:perils?\s*covered|risks?\s*covered|coverage\s*includes?)[\s:]*([^,\n\r;]{10,300})',
                r'(?:insured\s*against|protection\s*against)[\s:]+([^,\n\r;]{10,300})',
                r'(?:covers|covering|coverage\s*for)[\s:]*([^,\n\r;]{10,300})',
            ],
            'geographical_limit': [
                r'(?:geographical?\s*limit|territory|territorial\s*scope)[\s:]+([^,\n\r;]{4,100})',
                r'(?:coverage\s*territory|insured\s*location|worldwide|global)[\s:]*([^,\n\r;]{4,100})',
                r'(?:within|throughout|covering)[\s:]+([A-Za-z][^,\n\r;]{4,100})',
            ],

            'situation_of_risk': [
                r'(?:situation\s*of\s*risk|risk\s*location|premises\s*address)[\s:]+([^,\n\r;]{10,300})',
                r'(?:located\s*at|situated\s*at|address)[\s:]+([^,\n\r;]{10,300})',
                r'(?:property\s*located|facility\s*at)[\s:]+([^,\n\r;]{10,300})',
            ],
            'occupation_of_insured': [
                r'(?:occupation|business|industry|sector|trade)[\s:]+([A-Za-z][^,\n\r;]{4,100})',
                r'(?:type\s*of\s*business|business\s*activity|commercial\s*activity)[\s:]+([A-Za-z][^,\n\r;]{4,100})',
                r'(?:industry\s*classification|business\s*type)[\s:]+([A-Za-z][^,\n\r;]{4,100})',
            ],
            'main_activities': [
                r'(?:main\s*activities|primary\s*activities|core\s*business|operations)[\s:]+([^,\n\r;]{10,300})',
                r'(?:business\s*activities|commercial\s*activities|operational\s*activities)[\s:]+([^,\n\r;]{10,300})',
                r'(?:nature\s*of\s*business|scope\s*of\s*operations)[\s:]+([^,\n\r;]{10,300})',
            ],
            # 'currency': [
            #     r'(?:currency|denominated\s*in|in\s*currency\s*of|equivalent\s*in)[\s:]+([A-Z]{3})',
            #     r'(?:equivalent\s*in|converted\s*to)[\s:]+([A-Z]{3})',
            #     r'(?:all\s*amounts?\s*in|figures?\s*in|values?\s*in)[\s:]+([A-Z]{3})',
            #     # Major global currencies
            #     r'(?:USD|US\$|\$|United\s*States\s*Dollar)',
            #     r'(?:EUR|€|Euro)',
            #     r'(?:GBP|£|British\s*Pound|Pound\s*Sterling)',
            #     r'(?:JPY|¥|Japanese\s*Yen)',
            #     r'(?:CHF|Swiss\s*Franc)',
            #     r'(?:CAD|Canadian\s*Dollar)',
            #     r'(?:AUD|Australian\s*Dollar)',
            #
            #     # Middle East & Africa
            #     r'(?:KES|KSh|Ksh|Kenya\s*Shilling)',
            #     r'(?:AED|UAE\s*Dirham)',
            #     r'(?:SAR|Saudi\s*Riyal)',
            #     r'(?:QAR|Qatari\s*Riyal)',
            #     r'(?:BHD|Bahraini\s*Dinar)',
            #     r'(?:KWD|Kuwaiti\s*Dinar)',
            #     r'(?:OMR|Omani\s*Rial)',
            #     r'(?:ZAR|South\s*African\s*Rand)',
            #     r'(?:EGP|Egyptian\s*Pound)',
            #     r'(?:MAD|Moroccan\s*Dirham)',
            #     r'(?:TND|Tunisian\s*Dinar)',
            #     r'(?:NGN|Nigerian\s*Naira)',
            #     r'(?:GHS|Ghanaian\s*Cedi)',
            #     r'(?:UGX|Ugandan\s*Shilling)',
            #     r'(?:TZS|Tanzanian\s*Shilling)',
            #     r'(?:ETB|Ethiopian\s*Birr)',
            #
            #     # Asian currencies
            #     r'(?:CNY|RMB|Chinese\s*Yuan)',
            #     r'(?:INR|Indian\s*Rupee)',
            #     r'(?:SGD|Singapore\s*Dollar)',
            #     r'(?:HKD|Hong\s*Kong\s*Dollar)',
            #     r'(?:MYR|Malaysian\s*Ringgit)',
            #     r'(?:THB|Thai\s*Baht)',
            #     r'(?:IDR|Indonesian\s*Rupiah)',
            #     r'(?:PHP|Philippine\s*Peso)',
            #     r'(?:KRW|South\s*Korean\s*Won)',
            #
            #
            # ],
            'currency': [
                # Context-based detection
                r'(?:currency|denominated\s*in|amounts?\s*in|values?\s*in)[\s:]+([A-Z]{3})',
                r'(?:equivalent\s*in|converted\s*to)[\s:]+([A-Z]{3})',
                # Direct currency mentions
                r'\b(USD|EUR|GBP|INR|CAD|AUD|PHP|SAR|AED|KES)\b',
                r'(?:US\s*Dollar|Euro|British\s*Pound|Indian\s*Rupee)',
                r'(?:Canadian\s*Dollar|Australian\s*Dollar|Philippine\s*Peso)',
            ],
        }

    def parse_tsi_and_currency(self, tsi_string: str) -> Tuple[Optional[float], Optional[str]]:
        """Enhanced TSI parsing with better currency detection."""
        if not tsi_string:
            return None, None

        # Currency mapping
        currency_map = {
            '$': 'USD', 'US$': 'USD', 'USD': 'USD',
            '€': 'EUR', 'EUR': 'EUR', 'EURO': 'EUR',
            '£': 'GBP', 'GBP': 'GBP',
            'INR': 'INR', 'RUPEE': 'INR',
            'CAD': 'CAD', 'AUD': 'AUD', 'PHP': 'PHP',
            'SAR': 'SAR', 'AED': 'AED', 'KES': 'KES', 'KSh': 'KES'
        }

        # Find currency (prioritize 3-letter codes)
        currency = None
        currency_patterns = [
            r'\b(USD|EUR|GBP|INR|CAD|AUD|PHP|SAR|AED|KES)\b',  # ISO codes
            r'(US\$|\$|€|£)',  # Symbols
        ]

        for pattern in currency_patterns:
            match = re.search(pattern, tsi_string, re.IGNORECASE)
            if match:
                found = match.group(1).upper()
                currency = currency_map.get(found, found)
                break

        # Extract number
        # Remove currency and letters, keep numbers and separators
        clean_string = re.sub(r'[A-Za-z$€£]', '', tsi_string)

        # Find the main number
        number_patterns = [
            r'(\d{1,3}(?:,\d{3})*(?:\.\d{2})?)',  # 1,000,000.00
            r'(\d+(?:\.\d+)?)',  # Simple number
        ]

        base_amount = None
        for pattern in number_patterns:
            match = re.search(pattern, clean_string)
            if match:
                try:
                    number_str = match.group(1).replace(',', '')
                    base_amount = float(number_str)
                    if base_amount > 0:
                        break
                except ValueError:
                    continue

        if base_amount is None:
            return None, currency

        # Apply multipliers
        lower_text = tsi_string.lower()
        if 'billion' in lower_text or ' bn' in lower_text:
            base_amount *= 1_000_000_000
        elif 'million' in lower_text or ' mn' in lower_text or ' m ' in lower_text:
            base_amount *= 1_000_000
        elif 'thousand' in lower_text or ' k ' in lower_text:
            base_amount *= 1_000

        return base_amount, currency

    def extract_surveyor_report_enhanced(self, text: str) -> Optional[str]:
        """Enhanced extraction for risk surveyor reports focusing on report names and sources."""
        patterns = self.patterns['risk_surveyor_report']

        # Priority order for different types of surveyor information
        surveyor_info = []

        # First, try to find specific report names/files
        report_name_patterns = [
            r'(?:surveyor[\'s]*\s*report|survey\s*report|inspection\s*report)[\s:]*["\']?([A-Za-z][A-Za-z0-9\s\-_\.\/()]{10,150})["\']?',
            r'([A-Za-z0-9][A-Za-z0-9\s\-_\.\/()]{5,80}\.(?:pdf|doc|docx|txt))',  # File names
        ]

        for pattern in report_name_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                if isinstance(match, tuple):
                    match = match[0] if match[0] else match[1]

                cleaned_match = str(match).strip()
                if len(cleaned_match) > 5 and not any(exclude in cleaned_match.lower()
                                                      for exclude in
                                                      ['accepts no liability', 'subject to', 'warranty']):
                    surveyor_info.append(f"Report: {cleaned_match}")

        # Then try to find surveyor companies/names
        company_patterns = [
            r'(?:report\s*by|surveyed\s*by|inspected\s*by)[\s:]*["\']?([A-Za-z][A-Za-z0-9\s&\-_\.\/()]{5,100})["\']?',
            r'(?:risk\s*management|engineering\s*consultant|surveyor)[\s:]*["\']?([A-Za-z][A-Za-z0-9\s&\-_\.\/()]{10,100})["\']?',
        ]

        for pattern in company_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                if isinstance(match, tuple):
                    match = match[0] if match[0] else match[1]

                cleaned_match = str(match).strip()
                if len(cleaned_match) > 5 and not any(exclude in cleaned_match.lower()
                                                      for exclude in
                                                      ['accepts no liability', 'subject to', 'warranty']):
                    surveyor_info.append(f"Surveyor: {cleaned_match}")

        # Finally, look for report dates
        date_patterns = [
            r'(?:survey|inspection|report)[\s]*(?:dated|on|of)[\s]*(\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4})',
        ]

        for pattern in date_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            for match in matches:
                surveyor_info.append(f"Date: {match}")

        # Return the most relevant information
        if surveyor_info:
            # Prioritize report names over surveyor names over dates
            report_items = [item for item in surveyor_info if item.startswith('Report:')]
            surveyor_items = [item for item in surveyor_info if item.startswith('Surveyor:')]
            date_items = [item for item in surveyor_info if item.startswith('Date:')]

            result_parts = []
            if report_items:
                result_parts.append(report_items[0])  # Take the first report name
            if surveyor_items:
                result_parts.append(surveyor_items[0])  # Take the first surveyor
            if date_items:
                result_parts.append(date_items[0])  # Take the first date

            return '; '.join(result_parts) if result_parts else None

        return None

    def extract_claims_experience_enhanced(self, text: str) -> Optional[str]:
        """Enhanced extraction for comprehensive claims experience data."""
        patterns = self.patterns['claims_experience']

        claims_data = {
            'yearly_claims': [],
            'loss_ratios': [],
            'claims_free_periods': [],
            'major_minor_claims': [],
            'recoveries': [],
            'general_experience': []
        }

        # Extract year-wise claims
        year_pattern = r'(20[0-9]{2})[\s]*[-:][\s]*([^,\n\r;\.]{5,200})'
        year_matches = re.findall(year_pattern, text, re.IGNORECASE | re.MULTILINE)
        for year, description in year_matches:
            clean_desc = description.strip()
            if len(clean_desc) > 3:
                claims_data['yearly_claims'].append(f"{year}: {clean_desc}")

        # Extract loss ratios
        ratio_patterns = [
            r'((?:loss|combined|claims?)\s*ratio[\s:]*[0-9]+(?:\.[0-9]+)?%?[^,\n\r;\.]{0,100})',
            r'((?:frequency|severity)\s*(?:rate|ratio)[\s:]*[^,\n\r;\.]{10,150})'
        ]

        for pattern in ratio_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                clean_match = str(match).strip()
                if len(clean_match) > 10:
                    claims_data['loss_ratios'].append(clean_match)

        # Extract claims-free periods
        claims_free_patterns = [
            r'((?:[0-9]+\s*years?|last\s*[0-9]+\s*years?|since\s*[0-9]{4})\s*(?:claims?\s*free|no\s*claims?|nil\s*losses?)[^,\n\r;\.]{0,100})',
            r'(claims?\s*free\s*(?:for|since)[\s]*[^,\n\r;\.]{5,100})',
        ]

        for pattern in claims_free_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                clean_match = str(match).strip()
                if len(clean_match) > 10:
                    claims_data['claims_free_periods'].append(clean_match)

        # Extract major/minor claims information
        major_minor_pattern = r'((?:major|minor|significant|small|large)\s*(?:claims?|losses?)[\s:]*[^,\n\r;\.]{10,200})'
        major_minor_matches = re.findall(major_minor_pattern, text, re.IGNORECASE | re.MULTILINE)
        for match in major_minor_matches:
            clean_match = str(match).strip()
            if len(clean_match) > 15:
                claims_data['major_minor_claims'].append(clean_match)

        # Extract recovery information
        recovery_pattern = r'((?:recovered|settlement|recovery|reimbursement)[\s:]*(?:USD|EUR|GBP|INR|KES|AED|SAR|CAD|AUD|\$|â‚¬|Â£)?\s*[0-9,.\s]*(?:million|thousand|m|k)?[^,\n\r;\.]{0,150})'
        recovery_matches = re.findall(recovery_pattern, text, re.IGNORECASE | re.MULTILINE)
        for match in recovery_matches:
            clean_match = str(match).strip()
            if len(clean_match) > 10:
                claims_data['recoveries'].append(clean_match)

        # Extract general experience statements
        general_patterns = [
            r'((?:well\s*below|above|below)\s*(?:average|industry|market)[\s]*(?:claims?|losses?|experience)[^,\n\r;\.]{0,150})',
            r'((?:satisfactory|excellent|poor|good)\s*(?:claims?|loss)\s*(?:experience|record)[^,\n\r;\.]{0,150})',
            r'((?:last|past|previous)\s*(?:3|three)\s*years?[\s]*[^\.]{20,200})'
        ]

        for pattern in general_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
            for match in matches:
                clean_match = str(match).strip()
                if len(clean_match) > 15:
                    claims_data['general_experience'].append(clean_match)

        # Compile comprehensive claims experience
        result_parts = []

        # Add yearly claims (limit to most recent 5 years)
        if claims_data['yearly_claims']:
            yearly_sorted = sorted(claims_data['yearly_claims'], reverse=True)[:5]
            result_parts.append("Yearly Claims: " + "; ".join(yearly_sorted))

        # Add loss ratios
        if claims_data['loss_ratios']:
            result_parts.append("Loss Ratios: " + "; ".join(claims_data['loss_ratios'][:3]))

        # Add claims-free periods
        if claims_data['claims_free_periods']:
            result_parts.append("Claims-Free: " + "; ".join(claims_data['claims_free_periods'][:2]))

        # Add major/minor claims
        if claims_data['major_minor_claims']:
            result_parts.append("Claim Types: " + "; ".join(claims_data['major_minor_claims'][:3]))

        # Add recoveries
        if claims_data['recoveries']:
            result_parts.append("Recoveries: " + "; ".join(claims_data['recoveries'][:2]))

        # Add general experience
        if claims_data['general_experience']:
            result_parts.append("General: " + "; ".join(claims_data['general_experience'][:2]))

        if result_parts:
            return " | ".join(result_parts)

        return None

    def extract_structured_claims(self, text: str) -> Optional[str]:
        """Extract structured claims experience data."""
        claims_patterns = [
            # Year-based claims: 2021: USD 50,000, 2022: Nil, etc.
            r'(20[0-9]{2}[\s]*[-:][\s]*[^,\n\r;]{5,100})',
            # Loss ratios: Loss ratio 15%, Combined ratio 85%
            r'((?:loss|combined)\s*ratio[\s:]*[0-9]+(?:\.[0-9]+)?%?)',
            # Claims summary: 3 years claims free, No major losses
            r'((?:[0-9]+\s*years?|last\s*[0-9]+)\s*(?:claims?\s*free|no\s*claims?|nil\s*losses?))',
        ]

        claims_data = []
        for pattern in claims_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
            claims_data.extend(matches)

        if claims_data:
            return '; '.join(claims_data[:5])  # Limit to 5 most relevant entries
        return None

    def extract_share_percentage(self, text: str) -> Optional[str]:
        """Extract share percentage with validation."""
        share_patterns = [
            r'(?:share\s*offered|seeking|require)[\s:]*([0-9]+(?:\.[0-9]+)?)\s*%',
            r'([0-9]+(?:\.[0-9]+)?)%\s*(?:share|participation|line)',
            r'(?:line\s*capacity|participation)[\s:]*([0-9]+(?:\.[0-9]+)?)\s*%',
        ]

        for pattern in share_patterns:
            matches = re.findall(pattern, text, re.IGNORECASE)
            if matches:
                try:
                    percentage = float(matches[0])
                    if 0 < percentage <= 100:  # Valid percentage range
                        return f"{percentage}%"
                except ValueError:
                    continue
        return None

    # def extract_field(self, text: str, field_key: str) -> Optional[str]:
    #     """Enhanced field extraction with special handling for complex fields."""
    #     if field_key not in self.patterns:
    #         return None
    #
    #     # Special handling for specific fields
    #     if field_key == 'claims_experience':
    #         structured_claims = self.extract_structured_claims(text)
    #         if structured_claims:
    #             return structured_claims
    #
    #     if field_key == 'share_offered':
    #         share_percentage = self.extract_share_percentage(text)
    #         if share_percentage:
    #             return share_percentage
    #
    #     # Standard pattern matching for other fields
    #     patterns = self.patterns[field_key]
    #     text_lines = text.split('\n')
    #
    #     # Context keywords for better targeting
    #     context_keywords = {
    #         'risk_surveyor_report': ['survey', 'inspection', 'recommendations', 'technical', 'engineering'],
    #         'inward_acceptances': ['inward', 'treaty', 'accepted', 'participating', 'consortium'],
    #         'claims_experience': ['claims', 'loss', 'experience', 'history', 'ratio'],
    #         'insured': ['name of insured', 'insured:', 'assured:', 'policyholder:', 'client:'],
    #         'cedant': ['cedant:', 'ceding company', 'direct insurer', 'original insurer'],
    #         'broker': ['broker:', 'intermediary:', 'placed by', 'via'],
    #         'total_sum_insured': ['sum insured', 'tsi', 'insured value', 'coverage amount', 'limit'],
    #         'share_offered': ['share offered', '% offered', 'seeking', 'participation','quota', 'line capacity','share'],
    #         'premium_rates': ['rate', 'pricing', '% rate', 'per mille'],
    #         'period_of_insurance': ['period', 'effective', 'coverage period', 'policy period', 'duration']
    #     }
    #
    #     field_contexts = context_keywords.get(field_key, [])
    #
    #     # Try contextual extraction first
    #     for line in text_lines:
    #         line_lower = line.lower()
    #         has_context = any(ctx.lower() in line_lower for ctx in field_contexts)
    #
    #         if has_context:
    #             for pattern in patterns:
    #                 try:
    #                     matches = re.findall(pattern, line, re.IGNORECASE | re.MULTILINE)
    #                     if matches:
    #                         for match in matches:
    #                             cleaned_match = self._validate_and_clean_match(match, field_key)
    #                             if cleaned_match:
    #                                 return cleaned_match
    #                 except re.error as e:
    #                     logger.warning(f"Regex error: {e}")
    #                     continue
    #
    #     # Fallback to full text searchs
    #     for pattern in patterns:
    #         try:
    #             matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
    #             if matches:
    #                 for match in matches:
    #                     cleaned_match = self._validate_and_clean_match(match, field_key)
    #                     if cleaned_match:
    #                         return cleaned_match
    #         except re.error as e:
    #             logger.warning(f"Regex error: {e}")
    #             continue
    #
    #     return None

    def extract_field(self, text: str, field_key: str) -> Optional[str]:
        """Enhanced field extraction with special handling for complex fields."""
        if field_key not in self.patterns:
            return None

        # Special handling for enhanced fields
        if field_key == 'risk_surveyor_report':
            enhanced_result = self.extract_surveyor_report_enhanced(text)
            if enhanced_result:
                return enhanced_result

        if field_key == 'claims_experience':
            enhanced_result = self.extract_claims_experience_enhanced(text)
            if enhanced_result:
                return enhanced_result
            # Fallback to structured claims extraction
            structured_claims = self.extract_structured_claims(text)
            if structured_claims:
                return structured_claims

        if field_key == 'share_offered':
            share_percentage = self.extract_share_percentage(text)
            if share_percentage:
                return share_percentage

        # Standard pattern matching for other fields
        patterns = self.patterns[field_key]
        text_lines = text.split('\n')

        # Context keywords for better targeting
        context_keywords = {
            'risk_surveyor_report': ['survey', 'inspection', 'recommendations', 'technical', 'engineering', 'report'],
            'inward_acceptances': ['inward', 'treaty', 'accepted', 'participating', 'consortium'],
            'claims_experience': ['claims', 'loss', 'experience', 'history', 'ratio', 'frequency'],
            'insured': ['name of insured', 'insured:', 'assured:', 'policyholder:', 'client:'],
            'cedant': ['cedant:', 'ceding company', 'direct insurer', 'original insurer'],
            'broker': ['broker:', 'intermediary:', 'placed by', 'via'],
            'total_sum_insured': ['sum insured', 'tsi', 'insured value', 'coverage amount', 'limit'],
            'share_offered': ['share offered', '% offered', 'seeking', 'participation', 'quota', 'line capacity', 'share'],
            'premium_rates': ['rate', 'pricing', '% rate', 'per mille'],
            'period_of_insurance': ['period', 'effective', 'coverage period', 'policy period', 'duration']
        }

        field_contexts = context_keywords.get(field_key, [])

        # Try contextual extraction first
        for line in text_lines:
            line_lower = line.lower()
            has_context = any(ctx.lower() in line_lower for ctx in field_contexts)

            if has_context:
                for pattern in patterns:
                    try:
                        matches = re.findall(pattern, line, re.IGNORECASE | re.MULTILINE)
                        if matches:
                            for match in matches:
                                cleaned_match = self._validate_and_clean_match(match, field_key)
                                if cleaned_match:
                                    return cleaned_match
                    except re.error as e:
                        logger.warning(f"Regex error: {e}")
                        continue

        # Fallback to full text search
        for pattern in patterns:
            try:
                matches = re.findall(pattern, text, re.IGNORECASE | re.MULTILINE)
                if matches:
                    for match in matches:
                        cleaned_match = self._validate_and_clean_match(match, field_key)
                        if cleaned_match:
                            return cleaned_match
            except re.error as e:
                logger.warning(f"Regex error: {e}")
                continue

        return None

    def _validate_and_clean_match(self, match: Any, field_key: str) -> Optional[str]:
        """Enhanced validation with content filtering."""
        if isinstance(match, tuple):
            match = match[0] if match[0] else (match[1] if len(match) > 1 else '')

        if not match or not str(match).strip():
            return None

        cleaned_match = str(match).strip()

        # Global exclusions (policy/technical language)
        excluded_terms = [
            'accepts no liability', 'damage caused by viruses', 'transmitted via email',
            'warranties excluding', 'usa/canada', 'cross border reinsurers',
            'divisions deputed', 'plant head', 'consulting engineers fees',
            'foreign expert visit', 'limit 5%', 'subject to maximum'
        ]

        if any(term in cleaned_match.lower() for term in excluded_terms):
            return None

        # Field-specific validation
        if field_key in ['insured', 'cedant', 'broker']:
            # Must be a reasonable company name
            if len(cleaned_match) < 5 or len(cleaned_match) > 150:
                return None

            # Should not contain technical terms
            technical_terms = ['clause', 'warranty', 'exclusion', 'limit', 'subject to']
            if any(term in cleaned_match.lower() for term in technical_terms):
                return None

            # Should contain company indicators or be a proper business name
            company_indicators = ['ltd', 'limited', 'inc', 'corp', 'company', 'co', 'pvt', 'private', 'assurance',
                                  'insurance']
            word_count = len(cleaned_match.split())
            has_company_indicator = any(indicator in cleaned_match.lower() for indicator in company_indicators)

            if not has_company_indicator and word_count < 2:
                return None

        elif field_key == 'claims_experience':
            # Combine fragmented claims data into coherent sentence
            if len(cleaned_match) < 15:  # Too short to be meaningful
                return None
            # Clean up fragmented dates and descriptions
            cleaned_match = re.sub(r'\s+', ' ', cleaned_match)
            cleaned_match = re.sub(r';\s*', '; ', cleaned_match)

        return cleaned_match.title() if field_key in ['insured', 'cedant', 'broker'] else cleaned_match


@dataclass
class ExtractedInfo:
    """Holds extracted facultative reinsurance information."""
    insured: Optional[str] = None
    cedant: Optional[str] = None
    broker: Optional[str] = None
    occupation_of_insured: Optional[str] = None
    main_activities: Optional[str] = None
    perils_covered: Optional[str] = None
    geographical_limit: Optional[str] = None
    situation_of_risk: Optional[str] = None
    total_sum_insured: Optional[str] = None
    total_sum_insured_float: Optional[float] = None
    currency: Optional[str] = None
    period_of_insurance: Optional[str] = None
    excess_deductible: str = "TBD"
    retention_of_cedant: str = "TBD"
    possible_maximum_loss: str = "TBD"
    cat_exposure: str = "TBD"
    claims_experience: str = "TBD"
    reinsurance_deductions: str = "TBD"
    share_offered: str = "TBD"
    inward_acceptances: str = "TBD"
    risk_surveyor_report: str = "TBD"
    premium_rates: str = "TBD"
    premium: str = "TBD"
    climate_change_risk_factors: str = "TBD"
    esg_risk_assessment: str = "TBD"

    # System fields
    source_file: Optional[str] = None
    extraction_timestamp: Optional[str] = None
    confidence_score: float = 0.0
    used_ocr: bool = False
    email_metadata: Optional[Dict[str, str]] = None
    processing_method: Optional[str] = None

class FacultativeReinsuranceExtractor:
    """main workflow."""

    def __init__(self, use_cpu: bool = True, lang: str = 'en'):
        """Initialize the enhanced facultative reinsurance extractor."""
        self.email_processor = OptimizedEmailProcessor() if HAS_EMAIL_LIBS else None
        self.doc_processor = SmartDocumentProcessor()
        self.info_extractor = EnhancedInformationExtractor()

        # Supported file extensions with processing priorities
        self.supported_extensions = {
            # High priority - likely to contain structured data
            '.msg': 'email',
            '.pdf': 'document',
            '.docx': 'document',
            # Medium priority
            '.xlsx': 'spreadsheet',
            '.xls': 'spreadsheet',
            '.csv': 'data',
            '.pptx': 'presentation',
            # Low priority - require OCR
            '.jpg': 'image',
            '.jpeg': 'image',
            '.png': 'image',
            '.bmp': 'image',
            '.tiff': 'image'
        }

        logger.info("Enhanced FacultativeReinsuranceExtractor initialized")

    def process_file_smart(self, file_path: str) -> Tuple[Optional[str], Dict[str, Any]]:
        """
        Smart file processing with optimized workflow.

        Returns:
            Tuple of (extracted_text, processing_metadata)
        """
        file_path = Path(file_path)
        extension = file_path.suffix.lower()

        if extension not in self.supported_extensions:
            logger.warning(f"Unsupported file type: {extension}")
            return None, {'error': 'unsupported_file_type'}

        processing_metadata = {
            'file_type': self.supported_extensions[extension],
            'used_ocr': False,
            'processing_method': 'unknown',
            'email_metadata': {}
        }

        try:
            # EMAIL PROCESSING - Priority 1
            if extension == '.msg' and self.email_processor:
                text, email_metadata = self.email_processor.process_msg_with_attachments(
                    str(file_path), self.doc_processor
                )
                processing_metadata['processing_method'] = 'email_extraction'  # Keep original name
                processing_metadata['email_metadata'] = email_metadata
                processing_metadata['file_type'] = 'email'  # Add this line

                # Quick relevance check for emails
                if text and not self.email_processor.quick_check_email_relevance(text, email_metadata):
                    logger.info(f"Email {file_path.name} appears to be non-facultative related")
                    processing_metadata['relevance_score'] = 'low'
                else:
                    processing_metadata['relevance_score'] = 'high'

                return text, processing_metadata

            # DOCUMENT PROCESSING - Priority 2
            elif extension == '.pdf':
                text, used_ocr = self.doc_processor.process_pdf_smart(str(file_path))
                processing_metadata['used_ocr'] = used_ocr
                processing_metadata['processing_method'] = 'pdf_smart'
                return text, processing_metadata

            elif extension == '.docx':
                text = self.doc_processor.process_docx(str(file_path))
                processing_metadata['processing_method'] = 'docx_text'
                return text, processing_metadata

            # SPREADSHEET PROCESSING - Priority 3
            elif extension in ['.xlsx', '.xls']:
                text = self.doc_processor.process_excel_smart(str(file_path))
                processing_metadata['processing_method'] = 'excel_smart'
                return text, processing_metadata

            elif extension == '.csv':
                text = self._process_csv_simple(str(file_path))
                processing_metadata['processing_method'] = 'csv_pandas'
                return text, processing_metadata

            # PRESENTATION PROCESSING - Priority 4
            elif extension == '.pptx':
                text = self._process_pptx_simple(str(file_path))
                processing_metadata['processing_method'] = 'pptx_text'
                return text, processing_metadata

            # IMAGE PROCESSING - Priority 5 (OCR required)
            elif extension in ['.jpg', '.jpeg', '.png', '.bmp', '.tiff']:
                text = self._process_image_ocr(str(file_path))
                processing_metadata['used_ocr'] = True
                processing_metadata['processing_method'] = 'image_ocr'
                return text, processing_metadata

        except Exception as e:
            logger.error(f"Error processing file {file_path}: {e}")
            processing_metadata['error'] = str(e)
            return None, processing_metadata

        return None, processing_metadata

    def _process_csv_simple(self, file_path: str) -> str:
        """Simple CSV processing with memory optimization."""
        try:
            df = pd.read_csv(file_path, nrows=1000)  # Limit rows
            return df.to_string(index=False, na_rep='')
        except Exception as e:
            logger.error(f"CSV processing failed for {file_path}: {e}")
            return ""

    def _process_pptx_simple(self, file_path: str) -> str:
        """Simple PPTX text extraction."""
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text = ""
            for slide_num, slide in enumerate(prs.slides):
                text += f"\n--- Slide {slide_num + 1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
            return text
        except Exception as e:
            logger.error(f"PPTX processing failed for {file_path}: {e}")
            return ""

    def _process_image_ocr(self, file_path: str) -> str:
        """Process images using OCR (only when necessary)."""
        if not PADDLEOCR_AVAILABLE:
            logger.warning(f"OCR not available for image {file_path}")
            return ""

        try:
            ocr = self.doc_processor._get_ocr()
            if ocr is None:
                return ""

            result = ocr.ocr(file_path, cls=True)
            if result and result[0]:
                text = '\n'.join([line[1][0] for line in result[0] if line[1][0].strip()])
                return text
            return ""
        except Exception as e:
            logger.error(f"OCR processing failed for {file_path}: {e}")
            return ""

    def extract_information_enhanced(self, text: str, source_file: str,
                                     processing_metadata: Dict[str, Any]) -> ExtractedInfo:
        """
        Enhanced information extraction with TSI parsing.
        """
        extracted = ExtractedInfo(
            source_file=source_file,
            extraction_timestamp=datetime.now().isoformat(),
            used_ocr=processing_metadata.get('used_ocr', False),
            email_metadata=processing_metadata.get('email_metadata', {}),
            processing_method=processing_metadata.get('processing_method', 'unknown')
        )

        # Field mapping for all extractable fields
        field_mapping = {
            'insured': 'insured',
            'cedant': 'cedant',
            'broker': 'broker',
            'perils_covered': 'perils_covered',
            'geographical_limit': 'geographical_limit',
            'situation_of_risk': 'situation_of_risk',
            'occupation_of_insured': 'occupation_of_insured',
            'main_activities': 'main_activities',
            'total_sum_insured': 'total_sum_insured',
            'currency': 'currency',
            'period_of_insurance': 'period_of_insurance',
            'risk_surveyor_report': 'risk_surveyor_report',
            'inward_acceptances': 'inward_acceptances',
            'share_offered': 'share_offered',
            'claims_experience': 'claims_experience',
        }

        extracted_count = 0
        for pattern_key, attr_name in field_mapping.items():
            value = self.info_extractor.extract_field(text, pattern_key)
            if value and value.strip():
                setattr(extracted, attr_name, value)
                extracted_count += 1

        # Special handling for TSI with enhanced parsing
        if extracted.total_sum_insured:
            parsed_tsi, extracted_currency = self.info_extractor.parse_tsi_and_currency(extracted.total_sum_insured)
            extracted.total_sum_insured_float = parsed_tsi
            if extracted_currency and not extracted.currency:
                extracted.currency = extracted_currency

        # Calculate confidence score
        base_confidence = extracted_count / len(field_mapping)

        # Boost confidence for high-priority email sources
        if processing_metadata.get('file_type') == 'email' and processing_metadata.get('relevance_score') == 'high':
            base_confidence *= 1.2  # 20% boost for relevant emails

        extracted.confidence_score = min(base_confidence, 1.0)  # Cap at 1.0

        return extracted

    def process_folder_optimized(self, folder_path: str) -> List[ExtractedInfo]:
        """
        Optimized folder processing with smart file prioritization.
        """
        folder_path = Path(folder_path)
        if not folder_path.exists():
            logger.error(f"Folder not found: {folder_path}")
            return []

        results = []

        # Get all supported files
        all_files = []
        for file_path in folder_path.rglob('*'):
            if file_path.is_file() and file_path.suffix.lower() in self.supported_extensions:
                all_files.append(file_path)

        # Sort files by priority (emails first, then documents, then images)
        def get_file_priority(file_path):
            ext = file_path.suffix.lower()
            priority_map = {
                '.msg': 1,  # Highest priority
                '.pdf': 2,
                '.docx': 2,
                '.xlsx': 3,
                '.xls': 3,
                '.csv': 3,
                '.pptx': 4,
                '.jpg': 5,  # Lowest priority (OCR required)
                '.jpeg': 5,
                '.png': 5,
                '.bmp': 5,
                '.tiff': 5
            }
            return priority_map.get(ext, 6)

        sorted_files = sorted(all_files, key=get_file_priority)

        logger.info(f"Processing {len(sorted_files)} files in priority order...")

        for file_path in sorted_files:
            logger.info(f"Processing: {file_path}")

            text_content, processing_metadata = self.process_file_smart(str(file_path))

            if text_content and text_content.strip():
                extracted_info = self.extract_information_enhanced(
                    text_content, str(file_path.name), processing_metadata
                )
                results.append(extracted_info)

                logger.info(f"Extracted info from {file_path.name}, "
                            f"confidence: {extracted_info.confidence_score:.2f}, "
                            f"method: {extracted_info.processing_method}")
            else:
                logger.warning(f"No text extracted from {file_path.name}")

        return results

    def save_results_enhanced(self, results: List[ExtractedInfo], output_dir: str) -> None:
        """Save results with separate files per email source."""
        if not results:
            logger.warning("No results to save")
            return

        output_dir = Path(output_dir).resolve()
        try:
            output_dir.mkdir(parents=True, exist_ok=True)
            logger.info(f"Output directory ready: {output_dir}")
        except Exception as e:
            logger.error(f"Directory creation failed: {e}")
            return

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")

        # Group results by email source
        email_groups = {}
        non_email_results = []

        for result in results:
            if result.processing_method == 'email_extraction' and result.email_metadata:
                sender = result.email_metadata.get('sender', 'unknown_sender')
                # Clean sender name for filename
                clean_sender = re.sub(r'[<>:"/\\|?*]', '_', sender.split('@')[0] if '@' in sender else sender)

                if clean_sender not in email_groups:
                    email_groups[clean_sender] = []
                email_groups[clean_sender].append(result)
            else:
                non_email_results.append(result)

        # Save individual email files
        for sender, email_results in email_groups.items():
            try:
                self._save_email_group(email_results, output_dir, sender, timestamp)
            except Exception as e:
                logger.error(f"Failed to save email group {sender}: {e}")

        # Save combined non-email results
        if non_email_results:
            try:
                self._save_non_email_results(non_email_results, output_dir, timestamp)
            except Exception as e:
                logger.error(f"Failed to save non-email results: {e}")

        # Save overall summary
        try:
            self._save_overall_summary(results, output_dir, timestamp, len(email_groups))
        except Exception as e:
            logger.error(f"Failed to save overall summary: {e}")

    def _save_email_group(self, results: List[ExtractedInfo], output_dir: Path, sender: str, timestamp: str):
        """Save individual email group results."""
        data = []
        for result in results:
            result_dict = self._convert_result_to_dict(result)
            data.append(result_dict)

        df = pd.DataFrame(data)

        # Clean sender name further for filename
        safe_sender = sender[:30]  # Limit length

        # Save CSV
        csv_path = output_dir / f"email_{safe_sender}_{timestamp}.csv"
        df.to_csv(csv_path, index=False, encoding='utf-8')

        # Save Excel
        excel_path = output_dir / f"email_{safe_sender}_{timestamp}.xlsx"
        with pd.ExcelWriter(excel_path, engine='openpyxl') as writer:
            df.to_excel(writer, sheet_name='Extraction_Results', index=False)

            # Add email metadata sheet
            if results and results[0].email_metadata:
                email_info = pd.DataFrame([results[0].email_metadata])
                email_info.to_excel(writer, sheet_name='Email_Info', index=False)

        logger.info(f"Email group saved: {safe_sender} ({len(results)} files)")

    def _save_non_email_results(self, results: List[ExtractedInfo], output_dir: Path, timestamp: str):
        """Save non-email results."""
        data = []
        for result in results:
            result_dict = self._convert_result_to_dict(result)
            data.append(result_dict)

        df = pd.DataFrame(data)

        csv_path = output_dir / f"non_email_results_{timestamp}.csv"
        df.to_csv(csv_path, index=False, encoding='utf-8')

        excel_path = output_dir / f"non_email_results_{timestamp}.xlsx"
        df.to_excel(excel_path, index=False)

        logger.info(f"Non-email results saved ({len(results)} files)")

    def _convert_result_to_dict(self, result: ExtractedInfo) -> Dict[str, Any]:
        """Convert ExtractedInfo to dictionary with all fields."""
        result_dict = {
            # Core extracted fields
            'insured': str(result.insured) if result.insured else None,
            'cedant': str(result.cedant) if result.cedant else None,
            'broker': str(result.broker) if result.broker else None,
            'perils_covered': str(result.perils_covered) if result.perils_covered else None,
            'geographical_limit': str(result.geographical_limit) if result.geographical_limit else None,
            'situation_of_risk': str(result.situation_of_risk) if result.situation_of_risk else None,
            'occupation_of_insured': str(result.occupation_of_insured) if result.occupation_of_insured else None,
            'main_activities': str(result.main_activities) if result.main_activities else None,
            'total_sum_insured': str(result.total_sum_insured) if result.total_sum_insured else None,
            'currency': str(result.currency) if result.currency else None,
            'period_of_insurance': str(result.period_of_insurance) if result.period_of_insurance else None,

            # TBD fields
            'excess_deductible': result.excess_deductible,
            'retention_of_cedant': result.retention_of_cedant,
            'possible_maximum_loss': result.possible_maximum_loss,
            'cat_exposure': result.cat_exposure,
            'claims_experience': result.claims_experience,
            'reinsurance_deductions': result.reinsurance_deductions,
            'share_offered': result.share_offered,
            'inward_acceptances': result.inward_acceptances,
            'risk_surveyor_report': result.risk_surveyor_report,
            'premium_rates': result.premium_rates,
            'premium': result.premium,
            'climate_change_risk_factors': result.climate_change_risk_factors,
            'esg_risk_assessment': result.esg_risk_assessment,

            # System fields
            'source_file': str(result.source_file) if result.source_file else None,
            'total_sum_insured_float': float(
                result.total_sum_insured_float) if result.total_sum_insured_float is not None else None,
            'confidence_score': float(result.confidence_score),
            'used_ocr': bool(result.used_ocr),
            'extraction_timestamp': str(result.extraction_timestamp) if result.extraction_timestamp else None,
            'processing_method': str(result.processing_method) if result.processing_method else None,
        }

        # Add email metadata if available
        if result.email_metadata:
            result_dict['email_subject'] = str(result.email_metadata.get('subject', ''))
            result_dict['email_sender'] = str(result.email_metadata.get('sender', ''))
            result_dict['email_date'] = str(result.email_metadata.get('date', ''))

        return result_dict

    def _save_overall_summary(self, results: List[ExtractedInfo], output_dir: Path, timestamp: str,
                              email_groups_count: int):
        """Save overall processing summary."""
        summary = {
            'total_files_processed': len(results),
            'email_groups_created': email_groups_count,
            'processing_timestamp': datetime.now().isoformat(),
            'average_confidence_score': sum(r.confidence_score for r in results) / len(results),
            'high_confidence_count': sum(1 for r in results if r.confidence_score >= 0.5),
            'email_sources_count': sum(1 for r in results if r.processing_method == 'email_extraction'),
            'ocr_used_count': sum(1 for r in results if r.used_ocr),
            'tsi_parsed_count': sum(1 for r in results if r.total_sum_insured_float is not None),
        }

        summary_path = output_dir / f"overall_summary_{timestamp}.json"
        with open(summary_path, 'w', encoding='utf-8') as f:
            json.dump(summary, f, indent=2, ensure_ascii=False, default=str)

        logger.info(f"Overall summary saved: {summary_path}")


# def main():
#     # CONFIGURATION - Edit these paths according to your setup
#     INPUT_FOLDER = "test insuarance docs"
#     OUTPUT_FOLDER = "Friday Morning Sept 26 2025 "
# 
#     # Check if email libraries are available
#     if not HAS_EMAIL_LIBS:
#         logger.warning("Email processing libraries not available. MSG files will be skipped.")
# 
#     if not HAS_DOC_LIBS:
#         logger.error("Document processing libraries not available. Exiting.")
#         return
# 
#     # Initialize enhanced extractor
#     logger.info("Initializing Enhanced Facultative Reinsurance Extractor...")
#     try:
#         extractor = FacultativeReinsuranceExtractor()
#     except Exception as e:
#         logger.error(f"Failed to initialize extractor: {e}")
#         return
# 
#     # Process documents with optimized workflow
#     logger.info(f"Starting optimized extraction from: {INPUT_FOLDER}")
#     start_time = datetime.now()
# 
#     results = extractor.process_folder_optimized(INPUT_FOLDER)
# 
#     end_time = datetime.now()
#     processing_duration = (end_time - start_time).total_seconds()
# 
#     # Save enhanced results
#     if results:
#         # Use the enhanced save method
#         # extractor.save_results_enhanced = lambda r, o: save_results_enhanced(extractor, r, o)
#         extractor.save_results_enhanced(results, OUTPUT_FOLDER)
# 
#         logger.info(f"Extraction complete in {processing_duration:.2f} seconds. Processed {len(results)} files.")
# 
#         # Calculate and print enhanced summary
#         high_confidence_count = sum(1 for r in results if r.confidence_score >= 0.5)
#         avg_confidence = sum(r.confidence_score for r in results) / len(results)
#         email_count = sum(1 for r in results if r.processing_method == 'email_extraction')
#         ocr_count = sum(1 for r in results if r.used_ocr)
#         tsi_parsed_count = sum(1 for r in results if r.total_sum_insured_float is not None)
# 
#         print("\n" + "=" * 60)
#         print("ENHANCED EXTRACTION SUMMARY")
#         print("=" * 60)
#         print(f"Total files processed: {len(results)}")
#         print(f"Processing time: {processing_duration:.2f} seconds")
#         print(f"Average time per file: {processing_duration / len(results):.2f} seconds")
#         print(
#             f"High confidence extractions: {high_confidence_count} ({high_confidence_count / len(results) * 100:.1f}%)")
#         print(f"Email sources processed: {email_count}")
#         print(f"Files requiring OCR: {ocr_count}")
#         print(f"TSI values parsed: {tsi_parsed_count}")
#         print(f"Average confidence score: {avg_confidence:.3f}")
# 
#         # Show TSI summary if available
#         if tsi_parsed_count > 0:
#             tsi_values = [r.total_sum_insured_float for r in results if r.total_sum_insured_float is not None]
#             print(f"\nTSI Analysis:")
#             print(f"  Total TSI values: {len(tsi_values)}")
#             print(f"  Average TSI: {sum(tsi_values) / len(tsi_values):,.0f}")
#             print(f"  Min TSI: {min(tsi_values):,.0f}")
#             print(f"  Max TSI: {max(tsi_values):,.0f}")
# 
#         # Show processing method breakdown
#         method_counts = {}
#         for r in results:
#             method = r.processing_method or 'unknown'
#             method_counts[method] = method_counts.get(method, 0) + 1
# 
#         print(f"\nProcessing Methods:")
#         for method, count in sorted(method_counts.items()):
#             print(f"  {method}: {count}")
# 
#         print(f"\nResults saved to: {OUTPUT_FOLDER}")
#         print("=" * 60)
# 
#     else:
#         logger.warning("No files were successfully processed.")
#         print("\nNo files were processed successfully. Check:")
#         print("1. Input folder path is correct")
#         print("2. Folder contains supported file types (.msg, .pdf, .docx, etc.)")
#         print("3. Files are not corrupted or password-protected")
#         print("4. Required libraries are installed")
# 
# if __name__ == "__main__":
#     main()
