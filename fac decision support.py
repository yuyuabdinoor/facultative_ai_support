import tempfile
from typing import Iterator, Union

import cv2
import extract_msg
import fitz
import numpy as np
import openpyxl
import pandas as pd
from PIL import Image
from docx import Document
from paddleocr import PaddleOCR
from pptx import Presentation
from tqdm import tqdm

from config import OLLAMA_CONFIG, OCR_CONFIG, PROCESSING_CONFIG, CACHE_CONFIG
from utility import *


@dataclass
class PDFPageInfo:
    """Information about a PDF page"""
    page_number: int
    width: int
    height: int
    rotation: int
    text_length: int


class StreamingPDFProcessor:
    """
    Memory-efficient PDF processor that processes pages one at a time.

    Features:
    - Streaming page processing
    - Automatic memory management
    - Adaptive quality control
    - Error recovery per page
    """

    def __init__(
            self,
            initial_dpi: int = 300,
            enable_memory_monitoring: bool = True,
            enable_adaptive_quality: bool = True
    ):
        """
        Initialize PDF processor.

        Args:
            initial_dpi: Initial rendering DPI
            enable_memory_monitoring: Enable memory monitoring and cleanup
            enable_adaptive_quality: Enable adaptive quality control
        """
        self.initial_dpi = initial_dpi
        self.memory_monitor = MemoryMonitor() if enable_memory_monitoring else None
        self.quality_controller = AdaptiveQualityController(initial_dpi) if enable_adaptive_quality else None
        self.logger = get_logger(__name__)

    @track_performance('pdf.open')
    def get_pdf_info(self, pdf_path: Path) -> Dict[str, Any]:
        """
        Get PDF metadata without loading pages.

        Args:
            pdf_path: Path to PDF file

        Returns:
            Dictionary with metadata
        """

        try:
            doc = fitz.open(pdf_path)

            # Create PDFPageInfo objects for each page
            pages_info = []
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_info = PDFPageInfo(
                    page_number=page_num + 1,
                    width=int(page.rect.width),
                    height=int(page.rect.height),
                    rotation=page.rotation,
                    text_length=len(page.get_text())
                )
                pages_info.append(page_info)

            info = {
                'page_count': len(doc),
                'pages': [vars(p) for p in pages_info],  # Convert to dict
                'metadata': doc.metadata,
                'is_encrypted': doc.is_encrypted,
                'file_size_bytes': pdf_path.stat().st_size,
                'file_size_mb': pdf_path.stat().st_size / (1024 * 1024)
            }

            doc.close()
            return info

        except Exception as e:
            self.logger.error(f"Error getting PDF info: {e}")
            raise

    def _get_current_dpi(self) -> int:
        """Get current DPI from quality controller or default"""
        if self.quality_controller:
            return self.quality_controller.current_dpi
        return self.initial_dpi

    def _check_memory(self):
        """Check memory and adjust quality if needed"""
        if not self.memory_monitor:
            return

        if self.memory_monitor.is_memory_critical():
            self.memory_monitor.force_cleanup()

            if self.quality_controller:
                self.quality_controller.reduce_quality()

    @track_performance('pdf.process_page')
    def process_page(
            self,
            doc: fitz.Document,
            page_num: int,
            dpi: Optional[int] = None
    ) -> Tuple[np.ndarray, np.ndarray]:
        """
        Process a single PDF page.

        Args:
            doc: Open PDF document
            page_num: Page number (0-indexed)
            dpi: Optional DPI override

        Returns:
            Tuple of (processed_image, original_image) as numpy arrays
        """
        try:
            # Check memory before processing
            self._check_memory()

            # Get DPI
            if dpi is None:
                dpi = self._get_current_dpi()

            # Load page
            page = doc.load_page(page_num)

            # Create transformation matrix
            mat = fitz.Matrix(dpi / 72, dpi / 72)

            # Render page
            pix = page.get_pixmap(matrix=mat)

            # Convert to numpy array
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                pix.height, pix.width, pix.n
            )

            # Convert RGBA to RGB if needed
            if img.shape[2] == 4:
                img = cv2.cvtColor(img, cv2.COLOR_RGBA2RGB)

            # Make a copy for processed version
            processed = img.copy()
            original = img.copy()

            # Cleanup
            del pix

            record_metric('pdf.page_processed', 1, page_num=page_num, dpi=dpi)

            return processed, original

        except Exception as e:
            self.logger.error(f"Error processing page {page_num}: {e}")
            record_metric('pdf.page_error', 1, page_num=page_num)
            raise

    def iter_pages(
            self,
            pdf_path: Path,
            start_page: int = 0,
            end_page: Optional[int] = None
    ) -> Iterator[Tuple[int, np.ndarray, np.ndarray]]:
        """
        Iterate through PDF pages one at a time.

        Args:
            pdf_path: Path to PDF
            start_page: First page to process (0-indexed)
            end_page: Last page to process (exclusive), None = all pages

        Yields:
            Tuples of (page_num, processed_image, original_image)
        """
        doc = None

        try:
            doc = fitz.open(pdf_path)
            total_pages = len(doc)

            if end_page is None:
                end_page = total_pages
            else:
                end_page = min(end_page, total_pages)

            self.logger.info(f"Processing PDF pages {start_page} to {end_page - 1}")

            for page_num in range(start_page, end_page):
                try:
                    processed, original = self.process_page(doc, page_num)
                    yield page_num, processed, original

                    # Cleanup after each page
                    if self.memory_monitor:
                        self.memory_monitor.check_and_cleanup()

                except Exception as e:
                    self.logger.error(f"Failed to process page {page_num}: {e}")
                    # Continue with next page
                    continue

        finally:
            if doc:
                doc.close()
                gc.collect()  # Force cleanup after closing document

    def process_pdf_chunked(
            self,
            pdf_path: Path,
            chunk_size: int = 10,
            max_pages: Optional[int] = None
    ) -> Iterator[List[Tuple[int, np.ndarray, np.ndarray]]]:
        """
        Process PDF in chunks for better memory management.

        Args:
            pdf_path: Path to PDF
            chunk_size: Number of pages per chunk
            max_pages: Maximum total pages to process

        Yields:
            Lists of (page_num, processed_image, original_image) tuples
        """
        info = self.get_pdf_info(pdf_path)
        total_pages = info['page_count']

        if max_pages:
            total_pages = min(total_pages, max_pages)

        self.logger.info(f"Processing {total_pages} pages in chunks of {chunk_size}")

        for chunk_start in range(0, total_pages, chunk_size):
            chunk_end = min(chunk_start + chunk_size, total_pages)

            chunk = []
            for page_num, processed, original in self.iter_pages(pdf_path, chunk_start, chunk_end):
                chunk.append((page_num, processed, original))

            if chunk:
                yield chunk

            # Force cleanup between chunks
            if self.memory_monitor:
                self.memory_monitor.force_cleanup()

    def extract_page_text(self, pdf_path: Path, page_num: int) -> str:
        """
        Extract text from a single page (lightweight operation).

        Args:
            pdf_path: Path to PDF
            page_num: Page number (0-indexed)

        Returns:
            Extracted text
        """
        try:
            doc = fitz.open(pdf_path)
            page = doc.load_page(page_num)
            text = page.get_text()
            doc.close()
            return text
        except Exception as e:
            self.logger.error(f"Error extracting text from page {page_num}: {e}")
            return ""

    def has_text_layer(self, pdf_path: Path, sample_pages: int = 3) -> bool:
        """
        Check if PDF has a text layer (useful to skip OCR).

        Args:
            pdf_path: Path to PDF
            sample_pages: Number of pages to sample

        Returns:
            True if text layer detected
        """
        try:
            doc = fitz.open(pdf_path)
            total_pages = min(len(doc), sample_pages)

            has_text = False
            for page_num in range(total_pages):
                text = doc.load_page(page_num).get_text().strip()
                if len(text) > 50:  # More than 50 chars suggests real text
                    has_text = True
                    break

            doc.close()
            return has_text

        except Exception as e:
            self.logger.error(f"Error checking text layer: {e}")
            return False


class SmartTextTruncator:
    """
    Intelligently truncate text while preserving important sections.

    Features:
    - Identify and preserve key sections
    - Keep document structure
    - Maintain context around important fields
    """

    IMPORTANT_KEYWORDS = [
        'insured', 'cedant', 'broker', 'sum insured', 'total_sum_insured', 'tsi', 'rate', 'premium rate'
                                                                                          'premium', 'period',
        'coverage', 'retention', 'deductible', 'claims experience',
        'excess', 'share', 'limit', 'country', 'location', 'risk'
    ]

    def __init__(self):
        self.logger = get_logger(__name__)

    def score_section(self, text: str) -> float:
        """
        Score a text section by importance.

        Args:
            text: Text section

        Returns:
            Importance score (0-1)
        """
        text_lower = text.lower()

        # Count important keywords
        keyword_count = sum(1 for kw in self.IMPORTANT_KEYWORDS if kw in text_lower)

        # Presence of numbers (likely financial data)
        has_numbers = bool(re.search(r'\d+[,.]?\d*', text))

        # Contains dates
        has_dates = bool(re.search(r'\d{1,2}[\/\-]\d{1,2}[\/\-]\d{2,4}', text))

        # Calculate score
        score = 0.0
        score += min(keyword_count * 0.2, 0.6)  # Max 0.6 from keywords
        score += 0.2 if has_numbers else 0.0
        score += 0.2 if has_dates else 0.0

        return min(score, 1.0)

    def smart_truncate(
            self,
            text: str,
            max_length: int,
            keep_start_percent: float = 0.3,
            keep_end_percent: float = 0.2
    ) -> str:
        """
        Intelligently truncate text preserving important sections.

        Args:
            text: Full text
            max_length: Maximum length
            keep_start_percent: Percentage of budget for start
            keep_end_percent: Percentage of budget for end

        Returns:
            Truncated text with preserved important sections
        """
        if len(text) <= max_length:
            return text

        # Split into paragraphs
        paragraphs = text.split('\n\n')

        # Score each paragraph
        scored_paragraphs = [
            (i, para, self.score_section(para))
            for i, para in enumerate(paragraphs)
        ]

        # Always keep start and end
        start_budget = int(max_length * keep_start_percent)
        end_budget = int(max_length * keep_end_percent)
        middle_budget = max_length - start_budget - end_budget

        # Build output
        result_parts = []
        current_length = 0

        # Add start
        start_text = text[:start_budget]
        result_parts.append(start_text)
        current_length += len(start_text)

        # Add important middle sections
        middle_paragraphs = sorted(
            scored_paragraphs[1:-1],
            key=lambda x: x[2],
            reverse=True
        )

        included_middle = []
        for idx, para, score in middle_paragraphs:
            if current_length + len(para) < max_length - end_budget:
                included_middle.append((idx, para))
                current_length += len(para)

        # Sort by original order
        included_middle.sort(key=lambda x: x[0])

        if included_middle:
            result_parts.append("\n\n[...MIDDLE SECTIONS...]\n\n")
            result_parts.extend([para for _, para in included_middle])

        # Add end
        result_parts.append("\n\n[...END SECTION...]\n\n")
        end_text = text[-end_budget:]
        result_parts.append(end_text)

        return ''.join(result_parts)


class OfficeDocumentProcessor:
    """Process Office documents (docx, pptx, xlsx, csv) with table extraction"""

    @staticmethod
    def extract_docx(file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extract text and tables from Word with formatting preservation.

        Preserves:
        - Paragraph styles (Heading 1-9, Normal, etc.)
        - Bold, italic, underline
        - Font size
        - Text alignment
        - Table structure with cell styling
        """
        doc: Document = Document(file_path)
        content: Dict[str, Any] = {
            'paragraphs': [],
            'tables': [],
            'metadata': {
                'total_paragraphs': 0,
                'total_tables': 0,
                'table_validation': []
            }
        }

        # Extract paragraphs with formatting
        for para in doc.paragraphs:
            if para.text.strip():
                para_data = {
                    'text': para.text,
                    'style': para.style.name if para.style else 'Normal',
                    'alignment': str(para.alignment) if para.alignment else 'LEFT',
                    'runs': []
                }

                # Extract run-level formatting
                for run in para.runs:
                    if run.text:
                        run_data = {
                            'text': run.text,
                            'bold': run.bold,
                            'italic': run.italic,
                            'underline': run.underline,
                            'font_size': run.font.size.pt if run.font.size else None
                        }
                        para_data['runs'].append(run_data)

                content['paragraphs'].append(para_data)

        # Extract tables with cell formatting
        for table_idx, table in enumerate(doc.tables):
            table_data: Dict[str, Any] = {
                'table_number': table_idx + 1,
                'rows': len(table.rows),
                'columns': len(table.columns),
                'data': [],
                'formatted_data': [],
                'dataframe': None,
                'validation': {
                    'has_header': False,
                    'empty_cells': 0,
                    'merged_cells': 0,
                    'data_quality_score': 0.0
                }
            }

            # Extract with formatting
            total_cells: int = 0
            empty_cells: int = 0

            for row_idx, row in enumerate(table.rows):
                row_data: List[str] = []
                formatted_row: List[Dict] = []

                for cell in row.cells:
                    cell_text: str = cell.text.strip()
                    total_cells += 1
                    if not cell_text:
                        empty_cells += 1

                    row_data.append(cell_text)

                    # Extract cell formatting
                    cell_format = {
                        'text': cell_text,
                        'bold': False,
                        'italic': False,
                        'alignment': 'LEFT',
                        'is_header': False
                    }

                    # Check first paragraph in cell
                    if cell.paragraphs:
                        para = cell.paragraphs[0]
                        if para.runs:
                            # Check if any run is bold (typical for headers)
                            cell_format['bold'] = any(r.bold for r in para.runs if r.bold)
                            cell_format['italic'] = any(r.italic for r in para.runs if r.italic)

                        cell_format['alignment'] = str(para.alignment) if para.alignment else 'LEFT'

                        # Headers typically use Heading styles or are bold
                        if row_idx == 0 or cell_format['bold']:
                            cell_format['is_header'] = True

                    formatted_row.append(cell_format)

                table_data['data'].append(row_data)
                table_data['formatted_data'].append(formatted_row)

            # Validation metrics
            table_data['validation']['empty_cells'] = empty_cells
            table_data['validation']['data_quality_score'] = (
                (total_cells - empty_cells) / total_cells if total_cells > 0 else 0
            )

            # Detect headers based on formatting
            if table_data['formatted_data']:
                first_row = table_data['formatted_data'][0]
                header_score = sum(1 for cell in first_row if cell['is_header']) / len(first_row)
                table_data['validation']['has_header'] = header_score > 0.6

            # Create DataFrame
            try:
                if len(table_data['data']) > 1:
                    df: pd.DataFrame
                    if table_data['validation']['has_header']:
                        df = pd.DataFrame(table_data['data'][1:], columns=table_data['data'][0])
                    else:
                        df = pd.DataFrame(table_data['data'])

                    df = df.dropna(how='all').dropna(axis=1, how='all')

                    if not df.empty:
                        table_data['dataframe'] = df.to_dict('records')
                        table_data['dataframe_shape'] = df.shape
                        table_data['columns'] = df.columns.tolist()
                    else:
                        table_data['dataframe'] = []
                        table_data['dataframe_shape'] = (0, 0)
                        table_data['columns'] = []
                else:
                    table_data['dataframe'] = []
                    table_data['dataframe_shape'] = (0, 0)
                    table_data['columns'] = []

            except Exception as e:
                print(f"Could not convert table {table_idx + 1} to DataFrame: {e}")
                table_data['dataframe'] = []
                table_data['dataframe_shape'] = (0, 0)
                table_data['columns'] = []

            content['tables'].append(table_data)
            content['metadata']['table_validation'].append(table_data['validation'])

        content['metadata']['total_paragraphs'] = len(content['paragraphs'])
        content['metadata']['total_tables'] = len(content['tables'])

        return content

    @staticmethod
    def extract_pptx(file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extract text and tables from PowerPoint presentations

        Args:
            file_path: Path to the .pptx file

        Returns:
            Dictionary containing slides with text, tables, and metadata
        """

        prs: Presentation = Presentation(file_path)
        content: Dict[str, Any] = {
            'slides': [],
            'metadata': {
                'total_slides': len(prs.slides),
                'total_tables': 0,
                'total_text_boxes': 0
            }
        }

        for slide_idx, slide in enumerate(prs.slides):
            slide_content: Dict[str, Any] = {
                'slide_number': slide_idx + 1,
                'title': '',
                'title_formatted': None,
                'text_boxes': [],  # Changed from 'text' list
                'tables': []
            }

            # Extract title with formatting
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                slide_content['title'] = title_text

                # Extract title formatting
                if slide.shapes.title.text_frame:
                    title_formatted = {
                        'text': title_text,
                        'runs': []
                    }

                    for paragraph in slide.shapes.title.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                run_data = {
                                    'text': run.text,
                                    'bold': run.font.bold if run.font.bold is not None else False,
                                    'italic': run.font.italic if run.font.italic is not None else False,
                                    'underline': run.font.underline if run.font.underline is not None else False,
                                    'font_size': run.font.size.pt if run.font.size else None
                                }
                                title_formatted['runs'].append(run_data)

                    slide_content['title_formatted'] = title_formatted

            # Extract text boxes and shapes with formatting
            for shape in slide.shapes:
                # Skip title (already processed)
                if shape == slide.shapes.title:
                    continue

                # Text extraction with formatting
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text_box = {
                        'text': shape.text,
                        'paragraphs': []
                    }

                    for paragraph in shape.text_frame.paragraphs:
                        para_data = {
                            'text': paragraph.text,
                            'alignment': str(paragraph.alignment) if paragraph.alignment else 'LEFT',
                            'level': paragraph.level,
                            'runs': []
                        }

                        for run in paragraph.runs:
                            if run.text:
                                run_data = {
                                    'text': run.text,
                                    'bold': run.font.bold if run.font.bold is not None else False,
                                    'italic': run.font.italic if run.font.italic is not None else False,
                                    'underline': run.font.underline if run.font.underline is not None else False,
                                    'font_size': run.font.size.pt if run.font.size else None
                                }
                                para_data['runs'].append(run_data)

                        text_box['paragraphs'].append(para_data)

                    if text_box['text'].strip():
                        slide_content['text_boxes'].append(text_box)
                        content['metadata']['total_text_boxes'] += 1

                # Table extraction (keep existing logic)
                if shape.has_table:
                    table = shape.table
                    table_data: Dict[str, Any] = {
                        'table_number': len(slide_content['tables']) + 1,
                        'rows': len(table.rows),
                        'columns': len(table.columns),
                        'data': [],
                        'formatted_data': [],
                        'dataframe': None
                    }

                    # Extract table with cell formatting
                    for row_idx, row in enumerate(table.rows):
                        row_data: List[str] = []
                        formatted_row: List[Dict] = []

                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            row_data.append(cell_text)

                            # Extract cell formatting
                            cell_format = {
                                'text': cell_text,
                                'bold': False,
                                'italic': False,
                                'is_header': False
                            }

                            # Check text frame in cell
                            if cell.text_frame and cell.text_frame.paragraphs:
                                for para in cell.text_frame.paragraphs:
                                    if para.runs:
                                        # Check if any run is bold
                                        cell_format['bold'] = any(
                                            r.font.bold for r in para.runs
                                            if r.font.bold is not None and r.font.bold
                                        )
                                        cell_format['italic'] = any(
                                            r.font.italic for r in para.runs
                                            if r.font.italic is not None and r.font.italic
                                        )

                            # First row or bold cells are likely headers
                            if row_idx == 0 or cell_format['bold']:
                                cell_format['is_header'] = True

                            formatted_row.append(cell_format)

                        table_data['data'].append(row_data)
                        table_data['formatted_data'].append(formatted_row)

                    # Convert to DataFrame (existing logic)
                    try:
                        df: pd.DataFrame
                        if len(table_data['data']) > 1:
                            # Check if first row has headers
                            first_row_formatted = table_data['formatted_data'][0]
                            header_score = sum(1 for cell in first_row_formatted if cell['is_header']) / len(
                                first_row_formatted)

                            if header_score > 0.6:
                                df = pd.DataFrame(table_data['data'][1:], columns=table_data['data'][0])
                            else:
                                df = pd.DataFrame(table_data['data'])

                            table_data['dataframe'] = df.to_dict('records')
                            table_data['dataframe_shape'] = df.shape
                        else:
                            df = pd.DataFrame(table_data['data'])
                            table_data['dataframe'] = df.to_dict('records')
                            table_data['dataframe_shape'] = df.shape
                    except Exception as e:
                        print(f"Could not convert table to DataFrame: {e}")

                    slide_content['tables'].append(table_data)
                    content['metadata']['total_tables'] += 1

            content['slides'].append(slide_content)

        return content

    @staticmethod
    def extract_xlsx(file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Excel extraction with sheet analysis

        Args:
            file_path: Path to the .xlsx file

        Returns:
            Dictionary containing sheets with data and metadata
        """
        wb: openpyxl.Workbook = openpyxl.load_workbook(file_path, data_only=True)
        content: Dict[str, Any] = {
            'sheets': [],
            'metadata': {
                'total_sheets': len(wb.sheetnames),
                'sheet_names': wb.sheetnames,
                'sheet_analysis': []
            }
        }

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_data: Dict[str, Any] = {
                'sheet_name': sheet_name,
                'used_range': f"{ws.dimensions}",
                'data': [],
                'dataframe': None,
                'analysis': {
                    'has_data': False,
                    'data_density': 0.0,
                    'likely_headers': False
                }
            }

            # Extract all rows
            all_rows: List[List[Any]] = []
            for row in ws.iter_rows(values_only=True):
                all_rows.append(list(row))

            sheet_data['data'] = all_rows

            # Analyze sheet content
            if all_rows:
                # Check if sheet has meaningful data
                non_empty_cells: int = sum(
                    1 for row in all_rows for cell in row
                    if cell is not None and str(cell).strip()
                )
                total_cells: int = sum(len(row) for row in all_rows)
                sheet_data['analysis']['data_density'] = (
                    non_empty_cells / total_cells if total_cells > 0 else 0
                )
                sheet_data['analysis']['has_data'] = sheet_data['analysis']['data_density'] > 0.1

                try:
                    if len(all_rows) > 1:
                        # Better header detection
                        first_row: List[Any] = all_rows[0]
                        header_indicators: int = sum(
                            1 for cell in first_row
                            if cell and isinstance(cell, str) and
                            not str(cell).replace('.', '').replace(',', '').replace('-', '').isdigit()
                        )

                        df: pd.DataFrame
                        if header_indicators > len(first_row) * 0.5:
                            sheet_data['analysis']['likely_headers'] = True
                            df = pd.DataFrame(all_rows[1:], columns=first_row)
                        else:
                            df = pd.DataFrame(all_rows)
                    else:
                        df = pd.DataFrame(all_rows)

                    # Clean DataFrame
                    df = df.dropna(how='all').dropna(axis=1, how='all')

                    if not df.empty:
                        sheet_data['dataframe'] = df.to_dict('records')
                        sheet_data['dataframe_shape'] = df.shape
                        sheet_data['columns'] = df.columns.tolist()
                    else:
                        sheet_data['dataframe'] = []
                        sheet_data['dataframe_shape'] = (0, 0)
                        sheet_data['columns'] = []

                except Exception as e:
                    print(f"Could not convert sheet '{sheet_name}' to DataFrame: {e}")
                    sheet_data['dataframe'] = []
                    sheet_data['dataframe_shape'] = (0, 0)
                    sheet_data['columns'] = []

            content['sheets'].append(sheet_data)
            content['metadata']['sheet_analysis'].append(sheet_data['analysis'])

        return content

    @staticmethod
    def extract_csv(file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extract data from CSV files with encoding detection

        Args:
            file_path: Path to the .csv file

        Returns:
            Dictionary containing CSV data and metadata
        """
        content: Dict[str, Any] = {
            'filename': Path(file_path).name,
            'data': [],
            'dataframe': None,
            'metadata': {
                'encoding': 'unknown',
                'delimiter': 'unknown',
                'data_quality': 0.0
            }
        }

        # Try multiple encodings
        encodings: List[str] = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']

        for encoding in encodings:
            try:
                df: pd.DataFrame = pd.read_csv(file_path, encoding=encoding, on_bad_lines='skip')
                content['metadata']['encoding'] = encoding

                # Detect delimiter
                with open(file_path, 'r', encoding=encoding) as f:
                    sample: str = f.read(1024)
                    delimiters: List[str] = [',', ';', '\t', '|']
                    delimiter_counts: Dict[str, int] = {
                        delim: sample.count(delim) for delim in delimiters
                    }
                    content['metadata']['delimiter'] = max(delimiter_counts, key=delimiter_counts.get)

                # Calculate data quality
                total_cells: int = df.size
                non_null_cells: int = df.count().sum()
                content['metadata']['data_quality'] = (
                    non_null_cells / total_cells if total_cells > 0 else 0
                )

                content['dataframe'] = df.to_dict('records')
                content['dataframe_shape'] = df.shape
                content['columns'] = df.columns.tolist()
                content['metadata']['rows'] = len(df)
                content['metadata']['columns'] = len(df.columns)

                # Store raw data
                content['data'] = [df.columns.tolist()] + df.values.tolist()
                break

            except Exception as e:
                continue

        return content


class EmailProcessor:
    """Process .msg files with existing attachments in folder structure"""

    def __init__(self, subfolder_path: Union[str, Path]) -> None:
        """
        Initialize EmailProcessor

        Args:
            subfolder_path: Path to the folder containing the .msg file
        """
        self.subfolder_path: Path = Path(subfolder_path)
        self.msg_file: Optional[Path] = None
        self.msg_data: Optional[Dict[str, Any]] = None

    def find_msg_file(self) -> Optional[Path]:
        """
        Find the .msg file in the subfolder

        Returns:
            Path to .msg file if found, None otherwise
        """
        # iterate entries and match suffix case-insensitively
        for p in sorted(self.subfolder_path.iterdir(), key=lambda x: x.name.lower()):
            if p.is_file() and p.suffix.lower() == ".msg":
                self.msg_file = p
                return p
        return None

    def _safe_str(self, value: Any) -> Optional[str]:
        """
        Safely convert any value to string, handling bytes

        Args:
            value: Value to convert

        Returns:
            String representation or None
        """
        if value is None:
            return None
        if isinstance(value, bytes):
            try:
                return value.decode('utf-8', errors='ignore')
            except:
                return str(value)
        return str(value)

    def extract_email_data(self) -> Optional[Dict[str, Any]]:
        """
        Extract structured data from .msg file

        Returns:
            Dictionary containing email metadata and content, or None on error
        """
        if not self.msg_file:
            self.find_msg_file()

        if not self.msg_file:
            return None

        try:
            msg = extract_msg.Message(str(self.msg_file))

            self.msg_data = {
                'metadata': {
                    'subject': self._safe_str(msg.subject) if msg.subject else '',
                    'sender': self._safe_str(msg.sender) if msg.sender else '',
                    'sender_email': self._safe_str(getattr(msg, 'senderEmail', None)),
                    'date': self._safe_str(msg.date) if msg.date else None,
                    'received_time': self._safe_str(getattr(msg, 'receivedTime', None)) if hasattr(msg,
                                                                                                   'receivedTime') else None,
                    'message_id': self._safe_str(msg.messageId) if msg.messageId else '',
                    'importance': self._safe_str(msg.importance) if msg.importance else '',
                },
                'recipients': {
                    'to': self._format_recipients(msg.to),
                    'cc': self._format_recipients(msg.cc),
                    'bcc': self._format_recipients(msg.bcc)
                },
                'body': {
                    'plain_text': self._safe_str(msg.body) if msg.body else '',
                    'html': self._safe_str(msg.htmlBody) if msg.htmlBody else '',
                },
                'attachments_in_msg': self._get_msg_attachments(msg),
                'folder_path': str(self.subfolder_path),
                'msg_filename': self.msg_file.name
            }

            msg.close()
            return self.msg_data

        except Exception as e:
            print(f"Error extracting {self.msg_file}: {e}")
            traceback.print_exc()
            return None

    def _format_recipients(self, recipients: Any) -> List[Dict[str, str]]:
        """
        Format recipient list - handles both string and object formats

        Args:
            recipients: Recipients in various formats

        Returns:
            List of dictionaries with name and email
        """
        if not recipients:
            return []

        formatted: List[Dict[str, str]] = []

        if isinstance(recipients, str):
            return [{'name': self._safe_str(recipients), 'email': self._safe_str(recipients)}]

        if isinstance(recipients, list):
            for r in recipients:
                if isinstance(r, str):
                    formatted.append({'name': self._safe_str(r), 'email': self._safe_str(r)})
                elif hasattr(r, 'name') and hasattr(r, 'email'):
                    formatted.append({
                        'name': self._safe_str(r.name) if r.name else '',
                        'email': self._safe_str(r.email) if r.email else ''
                    })
                else:
                    formatted.append({'name': self._safe_str(r), 'email': self._safe_str(r)})

        return formatted

    def _get_msg_attachments(self, msg: Any) -> List[Dict[str, Any]]:
        """
        Get attachment info from .msg file - WITHOUT binary data

        Args:
            msg: Message object from extract_msg

        Returns:
            List of attachment metadata dictionaries
        """
        attachments: List[Dict[str, Any]] = []
        try:
            for att in msg.attachments:
                filename: str = att.longFilename or att.shortFilename or 'unknown'
                att_info: Dict[str, Any] = {
                    'filename': self._safe_str(filename),
                    'size': len(att.data) if att.data else 0,
                }
                attachments.append(att_info)
        except Exception as e:
            print(f"  Warning: Could not read attachments from msg: {e}")

        return attachments

    def list_folder_files(self) -> List[Dict[str, Any]]:
        """
        List all files in the subfolder (saved attachments)

        Returns:
            List of file metadata dictionaries
        """
        all_files: List[Dict[str, Any]] = []
        for file in self.subfolder_path.iterdir():
            if file.is_file() and not file.name.endswith('.msg'):
                all_files.append({
                    'filename': file.name,
                    'path': str(file),
                    'size': file.stat().st_size
                })
        return all_files

    def get_complete_data(self) -> Optional[Dict[str, Any]]:
        """
        Get email data + list of existing attachment files

        Returns:
            Complete email data dictionary or None on error
        """
        email_data: Optional[Dict[str, Any]] = self.extract_email_data()
        if not email_data:
            return None

        email_data['saved_attachments'] = self.list_folder_files()
        return email_data

    def process_attachments_with_ocr(
            self,
            email_data: Dict[str, Any],
            ocr_engine: 'Support',
            save_visuals: bool = True
    ) -> List[Dict[str, Any]]:
        """
        Process attachments with OCR, office documents, and direct text extraction.

        Args:
            email_data: Dictionary containing email metadata and attachments
            ocr_engine: Instance of Detect_and_Recognize class for OCR operations
            save_visuals: If True, saves annotated images and JSON results to:
                         <email_folder>/text_detected_and_recognized/<TIMESTAMP>/

        Returns:
            List of dictionaries containing extracted text and metadata for each attachment
        """
        attachment_texts: List[Dict[str, Any]] = []
        office_processor = OfficeDocumentProcessor()

        # If saving visuals, create timestamped output root in the email folder
        out_root: Optional[Path] = None
        if save_visuals:
            ts: str = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
            out_root = Path(self.subfolder_path) / "text_detected_and_recognized" / ts
            out_root.mkdir(parents=True, exist_ok=True)
            tqdm.write(f"[OCR] Saving visuals to: {out_root}")

        attachments: List[Dict[str, Any]] = email_data.get('saved_attachments', [])
        resource_limits = ResourceLimits()

        # Main attachment processing loop with progress bar
        for attachment in tqdm(attachments, desc="Processing attachments", unit="file", leave=False):
            file_path: str = attachment['path']
            ext: str = Path(file_path).suffix.lower()
            base_name: str = Path(file_path).stem

            if ext in ['.json']:
                continue

            start_time: float = time.time()

            # OFFICE DOCUMENTS - Process first before OCR
            x = SUPPORTED_EXTENSIONS['office']
            # print(f'Office extension; {ext} index; {x.index(ext)}')
            if ext in x[x.index('.docx')]:
                try:
                    content: Dict[str, Any] = office_processor.extract_docx(file_path)

                    # Extract all text for LLM context with style hints
                    text_parts: List[str] = []

                    # Extract paragraph text with heading markers
                    for para in content['paragraphs']:
                        if isinstance(para, dict):
                            text = para['text']
                            style = para.get('style', 'Normal')

                            # Add heading markers for important sections
                            if 'Heading' in style:
                                text = f"\n## {text} ##\n"

                            text_parts.append(text)
                        else:
                            text_parts.append(str(para))

                    # Extract table text
                    for table in content['tables']:
                        # Check if table has header row
                        has_header = table.get('validation', {}).get('has_header', False)

                        table_text: str = f"\n[TABLE {table['table_number']} - {table['rows']}x{table['columns']}]\n"

                        for row_idx, row in enumerate(table['data']):
                            row_text = " | ".join(str(cell) for cell in row)

                            # Mark header rows
                            if row_idx == 0 and has_header:
                                row_text = f"HEADER: {row_text}"

                            table_text += row_text + "\n"

                        text_parts.append(table_text)

                    full_text: str = "\n".join(text_parts)
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': full_text,
                        'method': 'docx_extraction',
                        'structured_content': content,
                        'time': time.time() - start_time
                    })

                    # Save structured data
                    if save_visuals and out_root is not None:
                        json_path: Path = out_root / f"{base_name}_docx_extracted.json"
                        with open(json_path, 'w', encoding='utf-8') as jf:
                            json.dump(content, jf, indent=2, ensure_ascii=False)

                except Exception as e:
                    tqdm.write(f"[DOCX] Error processing {file_path}: {e}")
                    traceback.print_exc()

            elif ext in x[x.index('.pptx')]:
                try:
                    content: Dict[str, Any] = office_processor.extract_pptx(file_path)

                    # Extract all text for LLM context
                    text_parts: List[str] = []
                    for slide in content['slides']:
                        slide_text: str = f"\n[SLIDE {slide['slide_number']}: {slide['title']}]\n"

                        # Add text boxes
                        for text_box in slide.get('text_boxes', []):
                            slide_text += text_box['text'] + "\n"

                        # Add tables
                        for table in slide['tables']:
                            table_text: str = f"\n[TABLE {table['table_number']} - {table['rows']}x{table['columns']}]\n"
                            for row in table['data']:
                                table_text += " | ".join(str(cell) for cell in row) + "\n"
                            slide_text += table_text

                        text_parts.append(slide_text)

                    full_text: str = "\n".join(text_parts)
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': full_text,
                        'method': 'pptx_extraction',
                        'structured_content': content,
                        'time': time.time() - start_time
                    })

                    # Save structured data with formatting
                    if save_visuals and out_root is not None:
                        json_path: Path = out_root / f"{base_name}_pptx_extracted.json"
                        with open(json_path, 'w', encoding='utf-8') as jf:
                            json.dump(content, jf, indent=2, ensure_ascii=False)

                except Exception as e:
                    tqdm.write(f"[PPTX] Error processing {file_path}: {e}")

            elif ext in x[x.index('.xlsx')]:
                try:
                    content: Dict[str, Any] = office_processor.extract_xlsx(file_path)

                    # Extract all sheets for LLM context
                    text_parts: List[str] = []
                    for sheet in content['sheets']:
                        sheet_text: str = f"\n[SHEET: {sheet['sheet_name']} - {sheet['used_range']}]\n"
                        if sheet['dataframe']:
                            # Create table representation
                            for row in sheet['data'][:100]:  # Limit to first 100 rows
                                sheet_text += " | ".join(str(cell) if cell is not None else "" for cell in row) + "\n"
                        text_parts.append(sheet_text)

                    full_text: str = "\n".join(text_parts)
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': full_text,
                        'method': 'xlsx_extraction',
                        'structured_content': content,
                        'time': time.time() - start_time
                    })

                    # Save structured data
                    if save_visuals and out_root is not None:
                        json_path: Path = out_root / f"{base_name}_xlsx_extracted.json"
                        with open(json_path, 'w', encoding='utf-8') as jf:
                            json.dump(content, jf, indent=2, ensure_ascii=False)

                except Exception as e:
                    tqdm.write(f"[XLSX] Error processing {file_path}: {e}")

            elif ext in x[x.index('.csv')]:
                try:
                    content: Dict[str, Any] = office_processor.extract_csv(file_path)

                    # Create text representation
                    text_parts: List[str] = [
                        f"[CSV FILE - {content['metadata'].get('rows', 0)} rows x {content['metadata'].get('columns', 0)} columns]\n"
                    ]
                    for row in content['data'][:100]:  # Limit to first 100 rows
                        text_parts.append(" | ".join(str(cell) if cell is not None else "" for cell in row))

                    full_text: str = "\n".join(text_parts)
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': full_text,
                        'method': 'csv_extraction',
                        'structured_content': content,
                        'time': time.time() - start_time
                    })

                    # Save structured data
                    if save_visuals and out_root is not None:
                        json_path: Path = out_root / f"{base_name}_csv_extracted.json"
                        with open(json_path, 'w', encoding='utf-8') as jf:
                            json.dump(content, jf, indent=2, ensure_ascii=False)

                except Exception as e:
                    tqdm.write(f"[CSV] Error processing {file_path}: {e}")

            # PDF - OCR Processing with nested progress bar
            elif ext in SUPPORTED_EXTENSIONS["documents"]:
                try:
                    pdf_start = time.time()

                    streaming_processor = StreamingPDFProcessor(
                        initial_dpi=PROCESSING_CONFIG.get('pdf_dpi', 300),
                        enable_memory_monitoring=True,
                        enable_adaptive_quality=True
                    )

                    # Open PDF once to analyze all pages
                    doc = fitz.open(file_path)
                    total_pages = len(doc)
                    tqdm.write(f"  [PDF] {base_name} - {total_pages} pages")

                    # Analyze each page for text content AND images
                    pages_analysis = []
                    for page_num in range(total_pages):
                        page = doc.load_page(page_num)
                        text = page.get_text().strip()

                        # Check for images on page
                        image_list = page.get_images(full=True)
                        has_images = len(image_list) > 0

                        # Determine processing strategy
                        has_text = len(text) > 50

                        pages_analysis.append({
                            'page_num': page_num,
                            'has_text': has_text,
                            'has_images': has_images,
                            'text_length': len(text),
                            'image_count': len(image_list),
                            'strategy': None  # Will determine below
                        })

                    doc.close()

                    # Determine strategy for each page
                    for page_info in pages_analysis:
                        if page_info['has_text'] and not page_info['has_images']:
                            # Pure text page - text extraction only
                            page_info['strategy'] = 'text_only'
                        elif not page_info['has_text'] and page_info['has_images']:
                            # Pure image page (scanned) - OCR only
                            page_info['strategy'] = 'ocr_only'
                        elif page_info['has_text'] and page_info['has_images']:
                            # Mixed content - BOTH text extraction AND OCR
                            page_info['strategy'] = 'hybrid'
                        else:
                            # Empty page
                            page_info['strategy'] = 'skip'

                    # Log strategy summary
                    strategy_counts = {}
                    for p in pages_analysis:
                        strat = p['strategy']
                        strategy_counts[strat] = strategy_counts.get(strat, 0) + 1

                    tqdm.write(f"  [PDF] Strategy: {strategy_counts}")

                    # Process each page according to its strategy
                    for page_info in pages_analysis:
                        page_num = page_info['page_num']
                        strategy = page_info['strategy']

                        if strategy == 'skip':
                            tqdm.write(f"    [PDF] Page {page_num + 1}: Empty, skipping")
                            continue

                        # TEXT EXTRACTION (for 'text_only' and 'hybrid')
                        if strategy in ['text_only', 'hybrid']:
                            page_text = streaming_processor.extract_page_text(file_path, page_num)
                            if page_text.strip():
                                attachment_texts.append({
                                    'file': attachment['filename'],
                                    'page': page_num + 1,
                                    'text': page_text,
                                    'method': 'pdf_text_layer',
                                    'source': 'text_extraction',
                                    'time': time.time() - pdf_start
                                })
                                tqdm.write(
                                    f"    [PDF] Page {page_num + 1}: Text extracted ({page_info['text_length']} chars)")

                        # OCR (for 'ocr_only' and 'hybrid')
                        if strategy in ['ocr_only', 'hybrid']:
                            reason = "scanned page" if strategy == 'ocr_only' else f"has {page_info['image_count']} images"
                            tqdm.write(f"    [PDF] Page {page_num + 1}: Running OCR ({reason})")

                            page_start = time.time()

                            # Use streaming processor to get page as image
                            doc = fitz.open(file_path)
                            try:
                                processed, original = streaming_processor.process_page(doc, page_num)

                                # Save temp image
                                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                                    img_pil = Image.fromarray(original)
                                    img_pil.save(tmp.name)
                                    tmp_path = tmp.name

                                # OCR
                                ocr_start = time.time()
                                results = ocr_engine.ocr_with_detection_and_recognition(tmp_path)
                                ocr_time = time.time() - ocr_start

                                # Save visuals
                                if save_visuals and out_root is not None and results:
                                    try:
                                        if isinstance(results, list) and len(results) > 0:
                                            result_obj = results[0]
                                            page_dir = out_root / f"{base_name}_page{page_num + 1}_OCR"
                                            page_dir.mkdir(exist_ok=True)
                                            result_obj.save_to_img(str(page_dir))

                                            # Save JSON with strategy info
                                            json_path = page_dir / "ocr_result.json"
                                            result_obj.save_to_json(str(json_path))

                                            # Save strategy info separately
                                            strategy_path = page_dir / "page_info.json"
                                            with open(strategy_path, 'w', encoding='utf-8') as f:
                                                json.dump(page_info, f, indent=2)

                                    except Exception as e:
                                        tqdm.write(f"      [PDF] Warning: Could not save page {page_num + 1}: {e}")

                                # Extract text from OCR
                                text_parts = []
                                for res in results:
                                    result_data = res if isinstance(res, dict) else (getattr(res, "json", None) or res)
                                    rec_texts = result_data.get('rec_texts', []) or []
                                    rec_scores = result_data.get('rec_scores', []) or []

                                    if rec_scores and len(rec_scores) == len(rec_texts):
                                        for t, s in zip(rec_texts, rec_scores):
                                            try:
                                                if float(s) >= OCR_CONFIG.get('confidence_threshold', 0.7) and str(
                                                        t).strip():
                                                    text_parts.append(str(t).strip())
                                            except Exception:
                                                continue

                                ocr_text = ' '.join(text_parts).strip()

                                if ocr_text:
                                    attachment_texts.append({
                                        'file': attachment['filename'],
                                        'page': page_num + 1,
                                        'text': ocr_text,
                                        'method': 'ocr',
                                        'source': 'image_ocr' if strategy == 'ocr_only' else 'hybrid_ocr',
                                        'time': time.time() - page_start,
                                        'ocr_time': ocr_time,
                                        'image_count': page_info['image_count']
                                    })
                                    tqdm.write(
                                        f"      [PDF] Page {page_num + 1}: OCR completed ({ocr_time:.2f}s, {len(text_parts)} text blocks)")
                                else:
                                    tqdm.write(f"      [PDF] Page {page_num + 1}: OCR found no text")

                                os.unlink(tmp_path)

                            finally:
                                doc.close()

                            # Memory check after each OCR page
                            if streaming_processor.memory_monitor:
                                streaming_processor.memory_monitor.check_and_cleanup()

                    tqdm.write(f"  [PDF] Total: {time.time() - pdf_start:.2f}s")

                except Exception as e:
                    tqdm.write(f"[OCR] Error processing PDF {file_path}: {e}")
                    traceback.print_exc()

            # IMAGE FILES
            elif ext in SUPPORTED_EXTENSIONS["images"]:
                try:
                    img_start: float = time.time()

                    # OCR with timing
                    ocr_start: float = time.time()
                    results: List[Any] = ocr_engine.ocr_with_detection_and_recognition(file_path)
                    ocr_time: float = time.time() - ocr_start

                    # SAVE FIRST - before any text filtering
                    if save_visuals and out_root is not None and results:
                        save_start: float = time.time()
                        try:
                            # Handle results structure - should be a list
                            result_obj: Any
                            if isinstance(results, list) and len(results) > 0:
                                result_obj = results[0]
                            else:
                                result_obj = results

                            # Create file-specific subdirectory for this image's outputs
                            img_dir: Path = out_root / f"{base_name}_image"
                            img_dir.mkdir(exist_ok=True)

                            # save_to_img expects a DIRECTORY
                            result_obj.save_to_img(str(img_dir))

                            # Save JSON with specific filename
                            json_path: Path = img_dir / "ocr_result.json"
                            result_obj.save_to_json(str(json_path))

                            save_time: float = time.time() - save_start
                            tqdm.write(f"  [IMG] {base_name} - OCR: {ocr_time:.2f}s, Save: {save_time:.2f}s")
                        except Exception as e:
                            tqdm.write(f"  [IMG] Warning: Could not save visuals: {e}")
                            traceback.print_exc()

                    # NOW extract text from result objects
                    text_parts: List[str] = []
                    for res in results:
                        # Access the JSON dict directly from the result object
                        result_data: Dict[str, Any]
                        if isinstance(res, dict):
                            result_data = res
                        else:
                            result_data = getattr(res, "json", None) or res

                        rec_texts: List[str] = result_data.get('rec_texts', []) or []
                        rec_scores: List[float] = result_data.get('rec_scores', []) or []

                        if rec_scores and len(rec_scores) == len(rec_texts):
                            for t, s in zip(rec_texts, rec_scores):
                                try:
                                    if float(s) >= 0.7 and str(t).strip():
                                        text_parts.append(str(t).strip())
                                except Exception:
                                    continue
                        else:
                            for t in rec_texts:
                                if isinstance(t, str) and len(t.strip()) >= 2:
                                    text_parts.append(t.strip())

                    full_text: str = ' '.join(text_parts).strip()

                    # Only add to attachment_texts if there's filtered text
                    # But we already saved visuals above regardless
                    if full_text:
                        attachment_texts.append({
                            'file': attachment['filename'],
                            'text': full_text,
                            'method': 'ocr',
                            'time': time.time() - img_start,
                            'ocr_time': ocr_time
                        })

                except Exception as e:
                    tqdm.write(f"[OCR] Error processing image {file_path}: {e}")
                    traceback.print_exc()

            # TEXT FILES - Direct extraction
            elif ext in SUPPORTED_EXTENSIONS["text"]:
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text: str = f.read()
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': text,
                        'method': 'direct',
                        'time': time.time() - start_time
                    })

                    # optionally save a JSON copy of the raw text to visuals folder
                    if save_visuals and out_root is not None:
                        json_name: str = f"{base_name}_raw_text.json"
                        json_path: Path = out_root / json_name
                        try:
                            with open(json_path, 'w', encoding='utf-8') as jf:
                                json.dump({"file": attachment['filename'], "text": text}, jf, indent=2,
                                          ensure_ascii=False)
                        except Exception as e:
                            tqdm.write(f"[OCR] Warning: failed to save raw text json {json_path}: {e}")
                except Exception as e:
                    tqdm.write(f"[TXT] Error processing {file_path}: {e}")

            # UNKNOWN - Try OCR as fallback
            else:
                try:
                    ocr_start: float = time.time()
                    results: List[Any] = ocr_engine.ocr_with_detection_and_recognition(file_path)
                    ocr_time: float = time.time() - ocr_start

                    # Extract text from results
                    text_parts: List[str] = []
                    for res in results:
                        result_data: Dict[str, Any]
                        if isinstance(res, dict):
                            result_data = res
                        else:
                            result_data = getattr(res, "json", None) or res

                        rec_texts: List[str] = result_data.get('rec_texts', []) or []
                        rec_scores: List[float] = result_data.get('rec_scores', []) or []

                        if rec_scores and len(rec_scores) == len(rec_texts):
                            for t, s in zip(rec_texts, rec_scores):
                                try:
                                    if float(s) >= 0.7 and str(t).strip():
                                        text_parts.append(str(t).strip())
                                except Exception:
                                    continue

                    text: str = ' '.join(text_parts).strip()

                    if text:
                        attachment_texts.append({
                            'file': attachment['filename'],
                            'text': text,
                            'method': 'ocr_fallback',
                            'time': time.time() - start_time,
                            'ocr_time': ocr_time
                        })

                    if save_visuals and out_root is not None and results:
                        try:
                            result_obj: Any = results[0] if isinstance(results, list) and len(results) > 0 else results
                            unknown_dir: Path = out_root / f"{base_name}_unknown"
                            unknown_dir.mkdir(exist_ok=True)
                            result_obj.save_to_img(str(unknown_dir))
                            json_path: Path = unknown_dir / "ocr_result.json"
                            result_obj.save_to_json(str(json_path))
                        except Exception as e:
                            tqdm.write(f"[OCR] Warning: failed to save unknown file type visuals: {e}")

                except Exception as e:
                    tqdm.write(f"[OCR] Skipping unknown file type {file_path}: {e}")

        return attachment_texts


class Support:
    def __init__(self, model_name: Optional[str] = None) -> None:
        self.model_name: str = model_name or OLLAMA_CONFIG["model_name"]
        self.ocr: PaddleOCR = PaddleOCR(
            # doc_orientation_classify_model_name='PP-LCNet_x1_0_doc_ori',
            # doc_orientation_classify_model_dir='PP-LCNet_x1_0_doc_ori',
            # textline_orientation_model_name='PP-LCNet_x1_0_textline_ori',
            # textline_orientation_model_dir='PP-LCNet_x1_0_textline_ori',
            text_detection_model_name=MODEL_PATHS['detection model'],
            text_recognition_model_name=MODEL_PATHS['recognition model'],
            text_detection_model_dir=MODEL_PATHS['detection folder'],
            text_recognition_model_dir=MODEL_PATHS['recognition folder'],
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False,
            device=OCR_CONFIG['device'],
            cpu_threads=OCR_CONFIG['cpu_threads'],
            enable_mkldnn=OCR_CONFIG['enable_mkldnn'],
            text_det_limit_side_len=OCR_CONFIG['det_limit_side_len'],
            text_det_limit_type='max',
            text_det_thresh=OCR_CONFIG['text_det_thresh'],
            text_det_box_thresh=OCR_CONFIG['text_det_box_thresh'],
            text_det_unclip_ratio=1.5,
            text_recognition_batch_size=OCR_CONFIG['text_recognition_batch_size'],
            # ocr_version="PP-OCRv5",
        )
        self._logger = get_logger(__name__)

    def __enter__(self):
        """Context manager entry"""
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        """Context manager exit - cleanup resources"""
        self.cleanup()
        return False

    def cleanup(self):
        """Cleanup OCR resources"""
        try:
            if hasattr(self, 'ocr') and self.ocr is not None:
                # PaddleOCR doesn't have explicit cleanup, but we can clear references
                del self.ocr
                self.ocr = None
                gc.collect()
                self._logger.debug("OCR resources cleaned up")
        except Exception as e:
            self._logger.warning(f"Error during OCR cleanup: {e}")

    def __del__(self):
        """Destructor - ensure cleanup"""
        self.cleanup()

    def ocr_with_detection_and_recognition(
            self,
            input_path: Union[str, np.ndarray, List[np.ndarray]]
    ) -> Any:
        """
        Perform OCR with both detection and recognition

        Args:
            input_path: Path to image, PDF, directory, or numpy array(s)

        Returns:
            PaddleOCR result object(s) with detection boxes and recognized text
            Format: List of page results, each containing:
                    [[bbox, (text, confidence)], ...]
        """
        return self.ocr.predict(input_path)

    def generate_response(
            self,
            prompt: str,
            max_tokens: int = 4000,
            temperature: float = 0.1
    ) -> Dict[str, Any]:
        """
        Generate response from local Ollama model using ollama library

        Args:
            prompt: The input prompt
            max_tokens: Maximum tokens to generate
            temperature: Sampling temperature (0.0 to 1.0)

        Returns:
            Dictionary with response data and metadata
        """
        start_time: float = time.time()

        try:
            with tqdm(total=2, desc="Generating response", leave=False) as pbar:
                pbar.set_description("Sending request to model")
                result = ollama.generate(
                    model=self.model_name,
                    prompt=prompt,
                    options={
                        "temperature": temperature,
                        "num_predict": max_tokens,
                        "top_p": 0.9,
                        "top_k": 40
                    }
                )
                pbar.update(1)

                pbar.set_description("Processing response")
                generation_time: float = time.time() - start_time
                response_text: str = result.get("response", "")
                pbar.update(1)

            return {
                "success": True,
                "response": response_text,
                "model": result.get("model", self.model_name),
                "generation_time": generation_time,
                "tokens_generated": len(response_text.split()),
                "metadata": {
                    "prompt_length": len(prompt),
                    "temperature": temperature,
                    "max_tokens": max_tokens
                }
            }

        except Exception as e:
            return {
                "success": False,
                "error": str(e),
                "generation_time": time.time() - start_time
            }

    def extract_structured_data(self, llm_prompt: str) -> Dict[str, Any]:
        """
        Extract structured reinsurance data using the local LLM

        Args:
            llm_prompt: The formatted prompt with email and attachment data

        Returns:
            Dictionary with extracted data and processing metadata
        """

        # Generate response
        result: Dict[str, Any] = self.generate_response(
            llm_prompt,
            max_tokens=OLLAMA_CONFIG["max_tokens"],
            temperature=OLLAMA_CONFIG["temperature"]
        )

        if not result["success"]:
            return {
                "success": False,
                "error": result["error"],
                "raw_response": None,
                "extracted_data": None
            }

        # Try to parse JSON response
        raw_response: str = result["response"].strip()

        # Clean up response - remove any markdown formatting
        with tqdm(total=3, desc="Processing LLM output", leave=False) as pbar:
            pbar.set_description("Cleaning response")
            if raw_response.startswith("```json"):
                raw_response = raw_response[7:]
            if raw_response.endswith("```"):
                raw_response = raw_response[:-3]
            raw_response = raw_response.strip()
            pbar.update(1)

            try:
                pbar.set_description("Parsing JSON")
                extracted_data: Dict[str, Any] = json.loads(raw_response)
                pbar.update(2)

                return {
                    "success": True,
                    "raw_response": raw_response,
                    "extracted_data": extracted_data,
                    "generation_time": result["generation_time"],
                    "tokens_generated": result["tokens_generated"],
                    "model_used": result["model"]
                }
            except json.JSONDecodeError as e:
                pbar.set_description("JSON parse failed, trying extraction")
                pbar.update(1)

                # Try to extract JSON from the response
                json_match = re.search(r'\{.*\}', raw_response, re.DOTALL)
                if json_match:
                    try:
                        extracted_data = json.loads(json_match.group())
                        pbar.update(1)
                        return {
                            "success": True,
                            "raw_response": raw_response,
                            "extracted_data": extracted_data,
                            "generation_time": result["generation_time"],
                            "tokens_generated": result["tokens_generated"],
                            "model_used": result["model"],
                            "warning": "JSON extracted from response with regex"
                        }
                    except json.JSONDecodeError:
                        pass

                pbar.update(1)
                return {
                    "success": False,
                    "error": f"JSON parsing failed: {e}",
                    "raw_response": raw_response,
                    "extracted_data": None,
                    "generation_time": result["generation_time"]
                }


MASTER_PROMPT: str = """You are a reinsurance data extraction specialist. Extract structured information from the provided email and attachments into valid JSON format.

# CRITICAL OUTPUT REQUIREMENTS
1. Output ONLY valid JSON - no markdown, no explanations, no code fences
2. Use "TBD" for any field not found in the text
3. All fields must be strings EXCEPT total_sum_insured_float (numeric)
4. Preserve original values - do not invent or calculate data not present
5. When multiple values exist for a field, concatenate with " | " separator

# EXTRACTION FIELDS

## Core Entities (Required)
- insured: Original risk owner (factory, airline, etc.)
- cedant: Insurance company ceding the risk
- broker: Intermediary facilitating placement
- policy_reference: Policy/certificate/slip number
- email_subject: Subject line of the email

## Reinsurance Structure
- reinsurance_type: "Facultative" | "Treaty" | "Facultative Obligatory" | "TBD"
- coverage_basis: "Proportional" | "Non-Proportional" | "Quota Share" | "Surplus" | "Excess of Loss" | "Stop Loss" | "TBD"
- layer_structure: For XL treaties (e.g., "USD 5M xs USD 10M" meaning $5M excess of $10M)

## Risk Details
- occupation_of_insured: Business/industry of insured
- main_activities: Core operations
- perils_covered: Specific risks (fire, flood, earthquake, etc.)
- geographical_limit: Coverage territory
- situation_of_risk: Physical location/route/address
- country: Risk country

## Financial Terms (Critical)
- total_sum_insured_and_breakdown: Full insured value with details (STRING)
  * Format: "Building: X | Machinery: Y | Stock: Z | Total: W"
  * Include all components if breakdown provided
- total_sum_insured_float: Numeric value only (FLOAT)
- currency: ISO code or symbol (USD, EUR, KES, GBP, $, €, etc.)
- premium: Actual premium amount (gross unless specified net)
- premium_rates: Rate as % or per mille (‰) - include unit symbol
- period_of_insurance: Coverage dates (from/to format preferred)
- valuation_basis: "Replacement Cost" | "Market Value" | "Agreed Value" | "Actual Cash Value" | "TBD"

## Reinsurance Terms
- retention_of_cedant: Cedant's kept share (% or amount)
- share_offered: Portion offered to this reinsurer (% or amount)
- excess_deductible: Loss threshold before coverage applies
- reinstatements: Number and premium terms (e.g., "2 reinstatements at 100% premium each")
- commission_rate: For proportional treaties (brokerage + overriding commission if split)
- profit_commission: Profit sharing terms if applicable
- aggregate_limit: Annual aggregate limit if specified

## Risk Assessment
- possible_maximum_loss: PML estimate (% of TSI or amount)
- cat_exposure: Catastrophe exposure details
- claims_experience: 3-year claims history with amounts and causes
- risk_surveyor_report: Technical inspection findings, recommendations, risk quality
- loss_ratio: Historical loss ratio if provided (calculated or stated)

## Multi-Party & Additional
- co_reinsurers: Other reinsurers and their shares (format: "Reinsurer A: 30% | Reinsurer B: 20%")
- lead_reinsurer: Lead reinsurer if subscription market
- reinsurance_deductions: Brokerage, taxes, levies (itemized if possible)
- inward_acceptances: Inward reinsurance flow (rare)
- warranties_conditions: Special conditions or warranties required
- territorial_exclusions: Excluded territories if any

## ESG & Emerging Risks
- climate_change_risk_factors: Climate exposure assessment ("High" | "Moderate" | "Low" | description)
- esg_risk_assessment: ESG risk level ("Low" | "Medium" | "High" | description)
- cyber_risk_exposure: If cyber risk mentioned or relevant
- pandemic_exclusions: COVID-19 or pandemic-related terms

# CALCULATION FORMULAS (Apply ONLY when source values are present)

## Premium Rate Calculation
When Premium and TSI are given:
- As percentage: Rate% = (Premium ÷ TSI) × 100
- As per mille: Rate‰ = (Premium ÷ TSI) × 1000
- Include calculated value in premium_rates field with note: "Calculated: X%"

## Premium Calculation
When Rate and TSI are given:
- From %: Premium = TSI × (Rate% ÷ 100)
- From ‰: Premium = TSI × (Rate‰ ÷ 1000)
- Mark as "Calculated: [amount]" if derived

## Loss Ratio Calculation
When claims data provided:
- Formula: Loss Ratio% = [(Paid Claims + Outstanding Claims - Recoveries) ÷ Earned Premium] × 100
- Use 3-year period if available
- Show calculation: "35% (Based on: Paid 300K + Outstanding 100K - Recoveries 50K / Premium 1M)"

## Consistency Checks
- If retention + share_offered = 100%, mark as verified
- If TSI breakdown components don't sum to total, flag: "Note: Components sum to X, Total stated as Y"
- Cross-reference premium calculation with stated rate when both present

# FIELD MAPPING GUIDE
Map these common synonyms to correct fields:

**Insured Variations:**
"Assured" | "Policyholder" | "Named Insured" | "Risk Owner" | "Original Insured" → insured

**Cedant Variations:**
"Ceding Company" | "Ceding Insurer" | "Reinsured" | "Direct Insurer" | "Original Insurer" → cedant

**Broker Variations:**
"Reinsurance Broker" | "Intermediary" | "Placing Broker" | "Insurance Broker" → broker

**Sum Insured Variations:**
"TSI" | "Total Limit" | "Aggregate Limit" | "Sum Insured" | "Limit of Liability" | "Total Insured Value" → total_sum_insured_and_breakdown

**Coverage Variations:**
"Covered Perils" | "Risks" | "Coverage Type" | "Insured Perils" | "Scope of Cover" → perils_covered

**Territory Variations:**
"Territory" | "Territorial Limit" | "Geographic Scope" | "Territorial Scope" → geographical_limit

**Deductible Variations:**
"Excess" | "Deductible" | "Franchise" | "First Loss" → excess_deductible

**Loss Variations:**
"PML" | "Maximum Loss" | "EML" | "Estimated Maximum Loss" | "Maximum Probable Loss" → possible_maximum_loss

**Report Variations:**
"Surveyor Report" | "Risk Report" | "Inspection Report" | "Survey Report" | "Risk Assessment" | "Site Survey" → risk_surveyor_report

# CURRENCY & NUMBER PARSING

## Currency Detection
- Symbols: $ (USD) | € (EUR) | £ (GBP) | ¥ (JPY) | KES/KSh (Kenyan Shilling)
- Handle formats: $1,000,000 | 1,000,000 USD | €1.000.000,00 | 1 000 000 EUR
- Recognize abbreviations: 50M = 50,000,000 | 2.5M = 2,500,000 | 1K = 1,000

## Number Parsing Rules
- Strip thousand separators: 1,000,000 or 1.000.000 or 1 000 000 → 1000000
- Handle decimals: Respect context (US: 1,000.50 | EU: 1.000,50)
- For total_sum_insured_float: Convert to clean float (no commas, spaces)
- Preserve original formatting in string fields

## Multi-Currency Handling
- If multiple currencies mentioned, extract original currency
- Note conversions if provided: "USD 1M (approx KES 130M at rate 130)"
- For TSI float, use primary/largest currency amount

# VALIDATION RULES

## Mandatory Consistency Checks
1. **TSI Validation**: If total_sum_insured_float present, currency must be present
2. **Date Validation**: Accept formats DD/MM/YYYY | YYYY-MM-DD | Month DD, YYYY | DD-Mon-YYYY
3. **Percentage Validation**: Retention + Share Offered should ≤ 100% (flag if exceeds)
4. **Premium-Rate-TSI Triangle**: If 2 of 3 present, verify consistency
5. **Period Validation**: End date should be after start date

## Data Quality Flags
When uncertain or inconsistent data found, append note to field:
- "Value unclear - extracted from: [source location]"
- "Conflict: Email states X, attachment states Y"
- "Assumed based on context"
- "Calculated - verify accuracy"

## Edge Cases
- **Missing breakdown but total given**: Extract total only, note "No breakdown provided"
- **Percentage vs absolute retention**: Specify which (e.g., "30%" vs "KES 50M")
- **Multiple periods**: Use policy period, note renewal if mentioned
- **Ranges**: Keep as range "5%-7.5%" or "Between X and Y"
- **TBD vs 0 vs Null**: Use "TBD" for unknown, "0" only if explicitly stated as zero/nil

# SPECIAL EXTRACTION SCENARIOS

## Email with Multiple Risks
- If email covers multiple risks, extract the primary/first risk
- Note in policy_reference: "Multiple risks - extracted: [Risk Name]"

## Attachments Priority
1. Formal slip/placement document (highest priority)
2. Risk survey report
3. Broker presentation
4. Email body
5. Signatures/chains (lowest priority)

## OCR Artifacts Handling
- Common OCR errors: 0/O, 1/I/l, 5/S, 8/B
- Validate extracted numbers against context (e.g., premium shouldn't exceed TSI)

## Implicit Information
- Country: Derive from situation_of_risk address if not stated
- Reinsurance_type: If 100% quota share or specific risk → likely Facultative
- Coverage_basis: If retention% mentioned → Proportional; if excess mentioned → Non-Proportional


# FINAL CHECKLIST BEFORE OUTPUT
□ All required fields present (even if "TBD")
□ total_sum_insured_float is numeric (no quotes)
□ All other fields are strings
□ Currency extracted when TSI present
□ Dates in consistent format
□ No markdown formatting or code fences
□ JSON is valid and parseable
□ Original phrasing preserved where possible
□ Calculations noted as "Calculated: X" when derived
□ No explanatory text outside JSON

# IMPORTANT REMINDERS
- Missing values = "TBD" (not null, not empty string)
- Do NOT calculate values not present in text UNLESS both source values exist
- Preserve original terminology and phrasing
- Extract verbatim where possible, paraphrase minimally
- When in doubt, use "TBD" rather than guessing
- Return ONLY the JSON object - no preamble, no postscript

---TEXT BLOCK STARTS---
{LLM_INPUT}
---TEXT BLOCK ENDS---
"""


def create_llm_context(
        email_data: Dict[str, Any],
        attachment_texts: List[Dict[str, Any]]
) -> Tuple[Dict[str, Any], str]:
    """
    Build LLM context and fill the MASTER_PROMPT with the combined text.
    - Skips attachments with empty text.
    - Returns (context_dict, llm_prompt_string).

    Args:
        email_data: Dictionary containing email metadata and body
        attachment_texts: List of dictionaries with attachment text and metadata

    Returns:
        Tuple of (context dictionary, formatted LLM prompt string)
    """

    # Initialize truncator
    truncator = SmartTextTruncator()

    # Filter attachments
    kept_attachments = []
    for att in attachment_texts:
        txt = att.get('text', '')
        if isinstance(txt, str) and txt.strip():
            # Smart truncate long attachments
            if len(txt) > PROCESSING_CONFIG.get('max_attachment_length', 15000):
                txt = truncator.smart_truncate(
                    txt,
                    max_length=PROCESSING_CONFIG.get('max_attachment_length', 15000),
                    keep_start_percent=0.4,
                    keep_end_percent=0.2
                )
                att['text'] = txt
                att['truncated'] = True
            kept_attachments.append(att)

    # Build context
    context = {
        'email': {
            'subject': email_data['metadata'].get('subject', ''),
            'from': email_data['metadata'].get('sender', ''),
            'date': email_data['metadata'].get('date', ''),
            'body': email_data['body'].get('plain_text', '') if 'body' in email_data else ''
        },
        'attachments': kept_attachments,
        'total_attachments': len(kept_attachments),
        'processing_times': {
            'email': 0,
            'attachments': sum(a.get('time', 0) for a in kept_attachments)
        }
    }

    # Build combined text
    parts = []
    parts.append(f"Email Subject: {context['email']['subject']}")
    parts.append(f"From: {context['email']['from']}")
    parts.append(f"Date: {context['email']['date']}")

    # Smart truncate email body
    email_body = context['email']['body'].strip()
    if len(email_body) > 5000:
        email_body = truncator.smart_truncate(email_body, max_length=5000)
    parts.append(f"\nEmail Body:\n{email_body}")

    if kept_attachments:
        parts.append("\nAttachments:")
        for i, att in enumerate(kept_attachments, start=1):
            header = f"\n--- Attachment {i}: {att['file']}"
            if att.get('page'):
                header += f" (page {att['page']})"
            if att.get('truncated'):
                header += " [TRUNCATED]"
            header += " ---"
            parts.append(header + "\n" + att['text'])

    combined_text = "\n".join(parts).strip()

    # Extract pattern candidates
    try:
        patterns = ReinsurancePatterns()
        candidates = extract_with_patterns(combined_text, patterns)

        if candidates:
            hint_lines = ["\n[AUTO-EXTRACTED HINTS]:"]
            for k, v in candidates.items():
                val_str = str(v.value if hasattr(v, 'value') else v)[:100]
                hint_lines.append(f"{k}: {val_str}")
            combined_text += "\n" + "\n".join(hint_lines)
    except Exception:
        pass

    # Final truncation using smart truncator
    if len(combined_text) > PROCESSING_CONFIG.get('max_text_length', 80000):
        combined_text = truncator.smart_truncate(
            combined_text,
            max_length=PROCESSING_CONFIG.get('max_text_length', 80000)
        )

    llm_prompt = MASTER_PROMPT.replace("{LLM_INPUT}", combined_text)

    return context, llm_prompt


def validate_processing_results(folder_path: Path) -> Dict[str, Any]:
    """
    Returns:
      {
        'validation': { 'context': True, ... },
        'missing_files': ['master_prompt.json'],
        'invalid_json': ['llm_response.json']
      }
    """
    validation: Dict[str, bool] = {}
    missing: List[str] = []
    invalid: List[str] = []
    expected_files: Dict[str, str] = {
        'context': 'llm_context.json',
        'prompt': 'master_prompt.json',
        'response': 'llm_response.json',
        'error': 'llm_error.json',
        'extraction': 'extracted_reinsurance_data.json'
    }

    for file_type, filename in expected_files.items():
        file_path = folder_path / filename
        if not file_path.exists():
            validation[file_type] = False
            missing.append(filename)
            continue
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                json.load(f)
            validation[file_type] = True
        except json.JSONDecodeError:
            validation[file_type] = False
            invalid.append(filename)

    return {'validation': validation, 'missing_files': missing, 'invalid_json': invalid}


def load_processing_results(folder_path: Path, strict: bool = False) -> Optional[Dict[str, Any]]:
    """
    Load results; if strict=True raise exceptions on missing/invalid files,
    otherwise return a dict with whatever could be read plus a 'diagnostics' key.
    """
    results: Dict[str, Any] = {}
    diag = {'missing': [], 'invalid': []}

    def try_load(name: str, filename: str):
        p = folder_path / filename
        if not p.exists():
            diag['missing'].append(filename)
            return None
        try:
            with open(p, 'r', encoding='utf-8') as f:
                return json.load(f)
        except Exception:
            diag['invalid'].append(filename)
            return None

    results['context'] = try_load('context', 'llm_context.json')
    results['prompt'] = try_load('prompt', 'master_prompt.json')
    # prefer response, fallback to error
    r = try_load('response', 'llm_response.json')
    if r is None:
        r = try_load('error', 'llm_error.json')
    results['llm_result'] = r
    results['extracted_data'] = try_load('extraction', 'extracted_reinsurance_data.json')
    results['diagnostics'] = diag

    if strict and (diag['missing'] or diag['invalid']):
        raise RuntimeError(f"Load failed: {diag}")

    return results


def get_log_folder(root_folder: Path) -> Path:
    """Create timestamped log subfolder within root folder"""
    timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S_%A")  # Includes day name
    log_folder = Path(root_folder) / "logs" / timestamp
    log_folder.mkdir(parents=True, exist_ok=True)
    return log_folder


def validate_extraction_json(data: Dict[str, Any]) -> ReinsuranceExtraction:
    """
    Validate and parse extraction JSON from LLM.

    Args:
        data: Raw dictionary from LLM response

    Returns:
        Validated ReinsuranceExtraction object

    Raises:
        ValidationError: If data doesn't match schema
    """
    return ReinsuranceExtraction(**data)


def validate_config(config_dict: Dict[str, Any]) -> AppConfig:
    """
    Validate application configuration.

    Args:
        config_dict: Configuration dictionary

    Returns:
        Validated AppConfig object

    Raises:
        ValidationError: If config is invalid
    """
    return AppConfig.from_dict(config_dict)


def safe_parse_extraction(
        data: Dict[str, Any],
        strict: bool = False
) -> tuple[Optional[ReinsuranceExtraction], List[str]]:
    """
    Safely parse extraction with error collection.

    Args:
        data: Raw extraction data
        strict: If True, raise on validation errors

    Returns:
        Tuple of (extraction_object or None, list of error messages)
    """
    try:
        extraction = ReinsuranceExtraction(**data)
        return extraction, []
    except Exception as e:
        if strict:
            raise
        return None, [str(e)]


def run_pipeline(root_folder: str, model_name: str = None, skip_processed: bool = True,
                 use_model_chaining: bool = False,
                 validation_model: str = None):
    """
    Main pipeline orchestration with all components integrated
    """
    root_path = Path(root_folder)
    cache_dir = root_path / ".prompt_cache"
    cache_dir.mkdir(exist_ok=True)

    # 1. Setup logging with timestamp
    log_folder = get_log_folder(root_path)
    logger = setup_logging(log_folder, log_level=LogLevel.INFO)
    logger.info("Pipeline starting", extra={'root_folder': root_folder, 'model': model_name})

    # 2. Health checks
    logger.info("Running health checks...")
    health_checker = HealthChecker()
    health_results = health_checker.run_all_checks(
        model_name or OLLAMA_CONFIG['model_name'],
        root_path
    )

    critical_checks = ['ollama', 'filesystem']
    critical_failed = any(
        not health_results[comp].healthy
        for comp in critical_checks
        if comp in health_results
    )

    if critical_failed:
        logger.error("Critical health checks failed - cannot proceed")
        return

    # Warn about memory but continue
    if not health_results.get('memory').healthy:
        logger.warning("Memory usage high - continuing with caution")
        logger.warning("Consider closing other applications")
    else:
        logger.info("All health checks passed")

    # 3 Continue with pipeline...
    logger.info("Initializing components...")

    with Support(model_name) as ds:
        rate_limiter = RateLimiter(max_operations=10, window_seconds=60)
        resource_limits = ResourceLimits(
            max_file_size_mb=PROCESSING_CONFIG.get('max_file_size_mb', 50),
            max_pdf_pages=PROCESSING_CONFIG.get('max_pdf_pages', 100)
        )
        input_validator = InputValidator()
        patterns = ReinsurancePatterns()

        # 4. Memory monitoring
        memory_monitor = MemoryMonitor()

        # 5. Process folders
        subfolders = [f for f in root_path.iterdir() if f.is_dir() and f.name != "logs"]
        logger.info(f"Found {len(subfolders)} folders to process")

        for subfolder in tqdm(sorted(subfolders), desc="Processing emails", unit="email"):
            process_single_folder(
                subfolder=subfolder,
                ds=ds,
                rate_limiter=rate_limiter,
                resource_limits=resource_limits,
                input_validator=input_validator,
                patterns=patterns,
                memory_monitor=memory_monitor,
                logger=logger,
                skip_processed=skip_processed
                # cache_dir=cache_dir,
                # use_model_chaining=use_model_chaining,
                # validation_model=validation_model
            )

        logger.info("Pipeline complete")

        # Export metrics
        metrics_path = log_folder / "metrics.json"
        export_metrics(metrics_path)
        logger.info(f"Metrics exported to {metrics_path}")

    # ds.cleanup() is automatically called here by context manager
    logger.info("Resources cleaned up")


@track_performance('folder.process')
def process_single_folder(
        subfolder: Path,
        ds: Support,
        rate_limiter: RateLimiter,
        resource_limits: ResourceLimits,
        input_validator: InputValidator,
        patterns: ReinsurancePatterns,
        memory_monitor: MemoryMonitor,
        logger: logging.Logger,
        skip_processed: bool = True
):
    """Process a single email folder with intelligent caching"""
    folder_logger = create_logger_with_context(__name__, folder=subfolder.name)

    try:
        folder_logger.info(f"Processing folder: {subfolder.name}")

        # Initialize cache
        cache = ExtractionCache(subfolder)
        model_name = ds.model_name
        timestamp = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")

        # Check if LLM extraction already exists for this model
        has_llm_cache, cached_response_path = cache.has_llm_cache(model_name)

        if has_llm_cache and skip_processed:
            folder_logger.info(f"Found cached LLM extraction for model {model_name}")

            # Load and return cached result
            with open(cached_response_path, 'r', encoding='utf-8') as f:
                cached_data = json.load(f)

            folder_logger.info(f"Using cached extraction from {cached_response_path.name}")
            return

        # Rate limiting
        if not rate_limiter.acquire():
            wait_time = rate_limiter.wait_time()
            folder_logger.warning(f"Rate limited, waiting {wait_time:.1f}s")
            time.sleep(wait_time)

        # Memory check
        memory_monitor.check_and_cleanup()

        # Extract email
        processor = EmailProcessor(subfolder)
        email_data = processor.get_complete_data()

        if not email_data:
            folder_logger.error("No email data found")
            return

        folder_logger.info(f"Email: {email_data['metadata']['subject']}")

        # Process attachments with caching
        attachment_texts = process_attachments_with_cache(
            processor=processor,
            email_data=email_data,
            ds=ds,
            cache=cache,
            timestamp=timestamp,
            logger=folder_logger
        )

        folder_logger.info(f"Processed {len(attachment_texts)} attachments")

        # Create LLM context
        context, llm_prompt = create_llm_context(email_data, attachment_texts)

        # Generate context hash for cache validation
        context_hash = hashlib.sha256(json.dumps(context, sort_keys=True).encode()).hexdigest()
        prompt_hash = hashlib.sha256(llm_prompt.encode()).hexdigest()

        # Create model-specific directory
        model_dir = subfolder / model_name
        model_dir.mkdir(exist_ok=True)

        # Extract candidates using patterns
        combined_text = llm_prompt
        candidates = extract_with_patterns(combined_text, patterns)

        if candidates:
            folder_logger.debug("Pattern extraction candidates", extra={'candidates': candidates})

        # Run LLM extraction
        folder_logger.info("Running LLM extraction...")
        llm_result = ds.extract_structured_data(llm_prompt)

        # Validate extraction
        if llm_result.get("success") and llm_result.get("extracted_data"):
            try:
                validated_extraction = validate_extraction_json(llm_result["extracted_data"])

                completeness = validated_extraction.get_completeness_score()
                missing_critical = validated_extraction.get_missing_critical_fields()

                folder_logger.info(
                    f"Extraction validated - Completeness: {completeness:.1f}%",
                    extra={
                        'completeness': completeness,
                        'missing_critical': missing_critical
                    }
                )

                llm_result["extracted_data"] = validated_extraction.dict()
                llm_result["completeness_score"] = completeness
                llm_result["missing_critical_fields"] = missing_critical

            except Exception as e:
                folder_logger.error("Extraction validation failed")
                log_exception(folder_logger, e, {'folder': subfolder.name})

        # Save results with timestamp
        save_processing_results_with_cache(
            model_dir=model_dir,
            context=context,
            llm_prompt=llm_prompt,
            llm_result=llm_result,
            timestamp=timestamp
        )

        # Register in cache
        cache.register_llm_cache(
            model_name=model_name,
            timestamp=timestamp,
            context_hash=context_hash,
            prompt_hash=prompt_hash
        )

        folder_logger.info(f"Results saved to {model_dir.name}/{timestamp}")

    except Exception as e:
        folder_logger.error(f"Error processing folder: {e}")
        log_exception(folder_logger, e, {'folder': subfolder.name})
        traceback.print_exc()


def process_attachments_with_cache(
        processor: EmailProcessor,
        email_data: Dict[str, Any],
        ds: Support,
        cache: ExtractionCache,
        timestamp: str,
        logger: logging.Logger
) -> List[Dict[str, Any]]:
    """
    Process attachments with intelligent caching.
    Skip OCR/extraction if already cached.
    """
    attachment_texts: List[Dict[str, Any]] = []
    office_processor = OfficeDocumentProcessor()
    attachments: List[Dict[str, Any]] = email_data.get('saved_attachments', [])

    # Create output root for new extractions
    out_root = Path(processor.subfolder_path) / "text_detected_and_recognized" / timestamp

    for attachment in tqdm(attachments, desc="Processing attachments", unit="file", leave=False):
        file_path: str = attachment['path']
        file_path_obj = Path(file_path)
        ext: str = file_path_obj.suffix.lower()
        base_name: str = file_path_obj.stem

        if ext in ['.json']:
            continue

        # Check cache first
        if cache.has_ocr_cache(file_path_obj):
            logger.info(f"Using cached extraction for {attachment['filename']}")
            cached_data = cache.get_cached_ocr_data(file_path_obj)

            if cached_data:
                # Convert cached data back to attachment text format
                attachment_texts.append({
                    'file': attachment['filename'],
                    'text': cached_data.get('text', ''),
                    'method': cached_data.get('method', 'cached'),
                    'cached': True,
                    'time': 0.0
                })
                continue

        # No cache - process normally
        start_time: float = time.time()

        # Process based on file type (existing logic)
        # ... [Keep existing processing logic] ...

        attachment_texts = processor.process_attachments_with_ocr(
            email_data,
            ds,
            save_visuals=PROCESSING_CONFIG['save_visuals']
        )

        # After successful processing, register in cache
        if out_root:
            cache.register_ocr_cache(
                file_path=file_path_obj,
                method=attachment_texts[-1].get('method', 'unknown') if attachment_texts else 'unknown',
                timestamp=timestamp,
                output_dir=out_root
            )

    return attachment_texts


def save_processing_results_with_cache(
        model_dir: Path,
        context: Dict[str, Any],
        llm_prompt: str,
        llm_result: Dict[str, Any],
        timestamp: str
) -> None:
    """
    Save processing results with timestamp in model-specific directory.
    """
    ts_iso = datetime.now().isoformat()

    with tqdm(total=4, desc="Saving results", unit="file", leave=False) as pbar:

        # Save LLM context (shared across runs)
        pbar.set_description("Saving LLM context")
        context_path: Path = model_dir / 'llm_context.json'
        context_data: Dict[str, Any] = {
            "context": context,
            "timestamp": ts_iso,
            "metadata": {
                "email_subject": context.get('email', {}).get('subject', ''),
                "total_attachments": context.get('total_attachments', 0),
                "processing_times": context.get('processing_times', {})
            }
        }
        with open(context_path, 'w', encoding='utf-8') as f:
            json.dump(context_data, f, indent=2, ensure_ascii=False)
        pbar.update(1)

        # Save master prompt (shared)
        pbar.set_description("Saving master prompt")
        prompt_path: Path = model_dir / 'master_prompt.json'
        prompt_data: Dict[str, Any] = {
            "master_prompt": llm_prompt,
            "timestamp": ts_iso,
            "prompt_length": len(llm_prompt),
        }
        with open(prompt_path, 'w', encoding='utf-8') as f:
            json.dump(prompt_data, f, indent=2, ensure_ascii=False)
        pbar.update(1)

        # Save timestamped LLM response
        if llm_result.get("success", False):
            pbar.set_description("Saving LLM response")
            result_path: Path = model_dir / f'llm_response_{timestamp}.json'

            metadata: Dict[str, Any] = llm_result.get("metadata", {})
            result_data: Dict[str, Any] = {
                "success": True,
                "extracted_data": llm_result.get("extracted_data", {}),
                "raw_response": llm_result.get("raw_response", ""),
                "timing": {
                    "generation_time_seconds": llm_result.get("generation_time", 0.0),
                    "tokens_generated": llm_result.get("tokens_generated", 0),
                    "timestamp": ts_iso
                },
                "model_info": {
                    "model_used": llm_result.get("model_used", "unknown"),
                    "temperature": metadata.get("temperature", 0.1),
                    "max_tokens": metadata.get("max_tokens", 3000)
                }
            }

            with open(result_path, 'w', encoding='utf-8') as f:
                json.dump(result_data, f, indent=2, ensure_ascii=False)
            pbar.update(1)

            # Save clean extraction
            pbar.set_description("Saving clean extraction")
            clean_path: Path = model_dir / f'extracted_reinsurance_data_{timestamp}.json'
            with open(clean_path, 'w', encoding='utf-8') as f:
                json.dump(llm_result.get("extracted_data", {}), f, indent=2, ensure_ascii=False)
            pbar.update(1)
        else:
            error_path: Path = model_dir / f'llm_error_{timestamp}.json'
            error_data: Dict[str, Any] = {
                "success": False,
                "error": llm_result.get("error", "Unknown error"),
                "timestamp": ts_iso
            }
            with open(error_path, 'w', encoding='utf-8') as f:
                json.dump(error_data, f, indent=2, ensure_ascii=False)
            pbar.update(2)


class ModelChain:
    """
    Multi-stage LLM processing with model chaining.

    Strategy:
    1. Small/fast model for initial extraction
    2. Larger model for validation and refinement
    3. Optional: Specialized model for specific fields
    """

    def __init__(
            self,
            primary_model: str,
            validation_model: Optional[str] = None,
            use_validation: bool = True
    ):
        """
        Initialize model chain.

        Args:
            primary_model: Main extraction model (fast)
            validation_model: Validation model (accurate)
            use_validation: Whether to use validation stage
        """
        self.primary_model = primary_model
        self.validation_model = validation_model or primary_model
        self.use_validation = use_validation
        self.logger = get_logger(__name__)

    def extract_with_chain(
            self,
            prompt: str,
            context: Dict[str, Any],
            support_instance: Support
    ) -> Dict[str, Any]:
        """
        Extract using model chain.

        Returns:
            Combined result with both stages
        """
        results = {
            'primary': None,
            'validation': None,
            'final': None,
            'confidence': 0.0,
            'stages_used': []
        }

        # Stage 1: Primary extraction
        self.logger.info(f"Stage 1: Primary extraction with {self.primary_model}")

        primary_result = support_instance.generate_response(
            prompt=prompt,
            max_tokens=OLLAMA_CONFIG['max_tokens'],
            temperature=0.1
        )

        results['primary'] = primary_result
        results['stages_used'].append('primary')

        if not primary_result.get('success'):
            self.logger.error("Primary extraction failed")
            return results

        # Try to parse primary result
        try:
            primary_data = self._parse_json_response(primary_result['response'])
            results['primary']['parsed_data'] = primary_data
        except Exception as e:
            self.logger.error(f"Primary parsing failed: {e}")
            return results

        # Stage 2: Validation (optional)
        if self.use_validation and self.validation_model != self.primary_model:
            self.logger.info(f"Stage 2: Validation with {self.validation_model}")

            validation_prompt = self._create_validation_prompt(
                original_extraction=primary_data,
                context=context
            )

            # Switch model temporarily
            original_model = support_instance.model_name
            support_instance.model_name = self.validation_model

            validation_result = support_instance.generate_response(
                prompt=validation_prompt,
                max_tokens=OLLAMA_CONFIG['max_tokens'],
                temperature=0.05  # Lower temp for validation
            )

            support_instance.model_name = original_model

            results['validation'] = validation_result
            results['stages_used'].append('validation')

            if validation_result.get('success'):
                try:
                    validation_data = self._parse_json_response(
                        validation_result['response']
                    )
                    results['validation']['parsed_data'] = validation_data

                    # Merge results (validation takes precedence for conflicts)
                    results['final'] = self._merge_extractions(
                        primary_data,
                        validation_data
                    )
                    results['confidence'] = 0.9

                except Exception as e:
                    self.logger.warning(f"Validation parsing failed: {e}")
                    results['final'] = primary_data
                    results['confidence'] = 0.7
            else:
                results['final'] = primary_data
                results['confidence'] = 0.7
        else:
            results['final'] = primary_data
            results['confidence'] = 0.8

        return results

    def _parse_json_response(self, response: str) -> Dict[str, Any]:
        """Parse JSON from LLM response"""
        cleaned = response.strip()
        if cleaned.startswith("```json"):
            cleaned = cleaned[7:]
        if cleaned.endswith("```"):
            cleaned = cleaned[:-3]
        cleaned = cleaned.strip()

        return json.loads(cleaned)

    def _create_validation_prompt(
            self,
            original_extraction: Dict[str, Any],
            context: Dict[str, Any]
    ) -> str:
        """Create validation prompt"""
        return f"""Review and validate the following reinsurance data extraction.
            Check for:
            1. Accuracy of extracted values
            2. Consistency between related fields
            3. Completeness of critical fields
            4. Proper formatting

            Original extraction:
            {json.dumps(original_extraction, indent=2)}

            Source context (first 2000 chars):
            {str(context)[:2000]}

            Output: Return the VALIDATED extraction as JSON. Fix any errors, but keep correct values unchanged. Use "TBD" for missing fields.
        """

    def _merge_extractions(
            self,
            primary: Dict[str, Any],
            validation: Dict[str, Any]
    ) -> Dict[str, Any]:
        """
        Merge two extractions, preferring validation for non-TBD values.
        """
        merged = primary.copy()

        for key, val_value in validation.items():
            # Skip if validation says TBD but primary has value
            if val_value == "TBD" and merged.get(key) != "TBD":
                continue

            # Prefer validation's value
            if val_value != "TBD":
                merged[key] = val_value

        return merged


def extract_with_caching_and_chaining(
        support: Support,
        prompt: str,
        context: Dict[str, Any],
        cache_dir: Path,
        use_cache: bool = True,
        use_chain: bool = False,
        validation_model: Optional[str] = None
) -> Dict[str, Any]:
    """
    Extract with optional caching and chaining.

    Args:
        support: Support instance
        prompt: Extraction prompt
        context: Email context
        cache_dir: Cache directory
        use_cache: Enable prompt caching
        use_chain: Enable model chaining
        validation_model: Optional validation model

    Returns:
        Extraction result
    """
    # Check cache first
    if use_cache:
        cache = PromptCache(cache_dir)
        cached = cache.get_cached_response(prompt, support.model_name)
        if cached:
            return cached['response']

    # Extract
    if use_chain and validation_model:
        chain = ModelChain(
            primary_model=support.model_name,
            validation_model=validation_model,
            use_validation=True
        )
        result = chain.extract_with_chain(prompt, context, support)
        extraction = result['final']
    else:
        result = support.extract_structured_data(prompt)
        extraction = result

    # Cache result
    if use_cache and result.get('success'):
        cache.cache_response(prompt, extraction, support.model_name)

    return extraction


root_folder = "thursday evening test data version two"
model_name = OLLAMA_CONFIG['model_name']

run_pipeline(
    root_folder=root_folder,
    model_name=model_name,
    skip_processed=PROCESSING_CONFIG['skip_processed'],
    validation_model=CACHE_CONFIG['validation_model']
)
