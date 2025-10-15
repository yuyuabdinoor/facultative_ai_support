import tempfile
from typing import Iterator, Union

import cv2
import extract_msg
import fitz
import numpy as np
import openpyxl
import pandas as pd
from PIL import Image
from docx import Document
from paddleocr import PaddleOCR
from pptx import Presentation
from tqdm import tqdm

from config import OCR_CONFIG, PROCESSING_CONFIG
from utility import *


class DetectAndRecognize:
    def __init__(self) -> None:
        self.ocr: PaddleOCR = PaddleOCR(
            # doc_orientation_classify_model_name='PP-LCNet_x1_0_doc_ori',
            # doc_orientation_classify_model_dir='PP-LCNet_x1_0_doc_ori',
            # textline_orientation_model_name='PP-LCNet_x1_0_textline_ori',
            # textline_orientation_model_dir='PP-LCNet_x1_0_textline_ori',
            text_detection_model_name=MODEL_PATHS['detection model'],
            text_recognition_model_name=MODEL_PATHS['recognition model'],
            text_detection_model_dir=MODEL_PATHS['detection folder'],
            text_recognition_model_dir=MODEL_PATHS['recognition folder'],
            use_doc_orientation_classify=False,
            use_doc_unwarping=False,
            use_textline_orientation=False,
            device=OCR_CONFIG['device'],
            cpu_threads=OCR_CONFIG['cpu_threads'],
            enable_mkldnn=OCR_CONFIG['enable_mkldnn'],
            text_det_limit_side_len=OCR_CONFIG['det_limit_side_len'],
            text_det_limit_type='max',
            text_det_thresh=OCR_CONFIG['text_det_thresh'],
            text_det_box_thresh=OCR_CONFIG['text_det_box_thresh'],
            text_det_unclip_ratio=1.5,
            text_recognition_batch_size=OCR_CONFIG['text_recognition_batch_size'],
        )
        self._logger = get_logger(__name__)

    def __enter__(self):
        """Context manager entry"""
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        """Context manager exit - cleanup resources"""
        self.cleanup()
        return False

    def cleanup(self):
        """Cleanup OCR resources"""
        try:
            if hasattr(self, 'ocr') and self.ocr is not None:
                # Clearing references
                del self.ocr
                self.ocr = None
                gc.collect()
                self._logger.debug("OCR resources cleaned up")
        except Exception as e:
            self._logger.warning(f"Error during OCR cleanup: {e}")

    def __del__(self):
        """Destructor - ensure cleanup"""
        self.cleanup()

    def ocr_with_detection_and_recognition(
            self,
            input_path: Union[str, np.ndarray, List[np.ndarray]]
    ) -> Any:
        """
        Perform OCR with both detection and recognition

        Args:
            input_path: Path to image, PDF, directory, or numpy array(s)

        Returns:
            PaddleOCR result object(s) with detection boxes and recognized text
            Format: List of page results, each containing:
                    [[bbox, (text, confidence)], ...]
        """
        return self.ocr.predict(input_path)


@dataclass
class PDFPageInfo:
    """Information about a PDF page"""
    page_number: int
    width: int
    height: int
    rotation: int
    text_length: int


class StreamingPDFProcessor:
    """
    Memory-efficient PDF processor that processes pages one at a time.

    Features:
    - Streaming page processing
    - Automatic memory management
    - Adaptive quality control
    - Error recovery per page
    """

    def __init__(
            self,
            initial_dpi: int = 300,
            enable_memory_monitoring: bool = True,
            enable_adaptive_quality: bool = True
    ):
        """
        Initialize PDF processor.

        Args:
            initial_dpi: Initial rendering DPI
            enable_memory_monitoring: Enable memory monitoring and cleanup
            enable_adaptive_quality: Enable adaptive quality control
        """
        self.initial_dpi = initial_dpi
        self.memory_monitor = MemoryMonitor() if enable_memory_monitoring else None
        self.quality_controller = AdaptiveQualityController(initial_dpi) if enable_adaptive_quality else None
        self.logger = get_logger(__name__)

    @track_performance('pdf.open')
    def get_pdf_info(self, pdf_path: Path) -> Dict[str, Any]:
        """
        Get PDF metadata without loading pages.

        Args:
            pdf_path: Path to PDF file

        Returns:
            Dictionary with metadata
        """

        try:
            doc = fitz.open(pdf_path)

            # Create PDFPageInfo objects for each page
            pages_info = []
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_info = PDFPageInfo(
                    page_number=page_num + 1,
                    width=int(page.rect.width),
                    height=int(page.rect.height),
                    rotation=page.rotation,
                    text_length=len(page.get_text())
                )
                pages_info.append(page_info)

            info = {
                'page_count': len(doc),
                'pages': [vars(p) for p in pages_info],  # Convert to dict
                'metadata': doc.metadata,
                'is_encrypted': doc.is_encrypted,
                'file_size_bytes': pdf_path.stat().st_size,
                'file_size_mb': pdf_path.stat().st_size / (1024 * 1024)
            }

            doc.close()
            return info

        except Exception as e:
            self.logger.error(f"Error getting PDF info: {e}")
            raise

    def _get_current_dpi(self) -> int:
        """Get current DPI from quality controller or default"""
        if self.quality_controller:
            return self.quality_controller.current_dpi
        return self.initial_dpi

    def _check_memory(self):
        """Check memory and adjust quality if needed"""
        if not self.memory_monitor:
            return

        if self.memory_monitor.is_memory_critical():
            self.memory_monitor.force_cleanup()

            if self.quality_controller:
                self.quality_controller.reduce_quality()

    @track_performance('pdf.process_page')
    def process_page(
            self,
            doc: fitz.Document,
            page_num: int,
            dpi: Optional[int] = None
    ) -> Tuple[np.ndarray, np.ndarray]:
        """
        Process a single PDF page.

        Args:
            doc: Open PDF document
            page_num: Page number (0-indexed)
            dpi: Optional DPI override

        Returns:
            Tuple of (processed_image, original_image) as numpy arrays
        """
        try:
            # Check memory before processing
            self._check_memory()

            # Get DPI
            if dpi is None:
                dpi = self._get_current_dpi()

            # Load page
            page = doc.load_page(page_num)

            # Create transformation matrix
            mat = fitz.Matrix(dpi / 72, dpi / 72)

            # Render page
            pix = page.get_pixmap(matrix=mat)

            # Convert to numpy array
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(
                pix.height, pix.width, pix.n
            )

            # Convert RGBA to RGB if needed
            if img.shape[2] == 4:
                img = cv2.cvtColor(img, cv2.COLOR_RGBA2RGB)

            # Make a copy for processed version
            processed = img.copy()
            original = img.copy()

            # Cleanup
            del pix

            record_metric('pdf.page_processed', 1, page_num=page_num, dpi=dpi)

            return processed, original

        except Exception as e:
            self.logger.error(f"Error processing page {page_num}: {e}")
            record_metric('pdf.page_error', 1, page_num=page_num)
            raise

    def iter_pages(
            self,
            pdf_path: Path,
            start_page: int = 0,
            end_page: Optional[int] = None
    ) -> Iterator[Tuple[int, np.ndarray, np.ndarray]]:
        """
        Iterate through PDF pages one at a time.

        Args:
            pdf_path: Path to PDF
            start_page: First page to process (0-indexed)
            end_page: Last page to process (exclusive), None = all pages

        Yields:
            Tuples of (page_num, processed_image, original_image)
        """
        doc = None

        try:
            doc = fitz.open(pdf_path)
            total_pages = len(doc)

            if end_page is None:
                end_page = total_pages
            else:
                end_page = min(end_page, total_pages)

            self.logger.info(f"Processing PDF pages {start_page} to {end_page - 1}")

            for page_num in range(start_page, end_page):
                try:
                    processed, original = self.process_page(doc, page_num)
                    yield page_num, processed, original

                    # Cleanup after each page
                    if self.memory_monitor:
                        self.memory_monitor.check_and_cleanup()

                except Exception as e:
                    self.logger.error(f"Failed to process page {page_num}: {e}")
                    # Continue with next page
                    continue

        finally:
            if doc:
                doc.close()
                gc.collect()  # Force cleanup after closing document

    def process_pdf_chunked(
            self,
            pdf_path: Path,
            chunk_size: int = 10,
            max_pages: Optional[int] = None
    ) -> Iterator[List[Tuple[int, np.ndarray, np.ndarray]]]:
        """
        Process PDF in chunks for better memory management.

        Args:
            pdf_path: Path to PDF
            chunk_size: Number of pages per chunk
            max_pages: Maximum total pages to process

        Yields:
            Lists of (page_num, processed_image, original_image) tuples
        """
        info = self.get_pdf_info(pdf_path)
        total_pages = info['page_count']

        if max_pages:
            total_pages = min(total_pages, max_pages)

        self.logger.info(f"Processing {total_pages} pages in chunks of {chunk_size}")

        for chunk_start in range(0, total_pages, chunk_size):
            chunk_end = min(chunk_start + chunk_size, total_pages)

            chunk = []
            for page_num, processed, original in self.iter_pages(pdf_path, chunk_start, chunk_end):
                chunk.append((page_num, processed, original))

            if chunk:
                yield chunk

            # Force cleanup between chunks
            if self.memory_monitor:
                self.memory_monitor.force_cleanup()

    def extract_page_text(self, pdf_path: Path, page_num: int) -> str:
        """
        Extract text from a single page (lightweight operation).

        Args:
            pdf_path: Path to PDF
            page_num: Page number (0-indexed)

        Returns:
            Extracted text
        """
        try:
            doc = fitz.open(pdf_path)
            page = doc.load_page(page_num)
            text = page.get_text()
            doc.close()
            return text
        except Exception as e:
            self.logger.error(f"Error extracting text from page {page_num}: {e}")
            return ""

    def has_text_layer(self, pdf_path: Path, sample_pages: int = 3) -> bool:
        """
        Check if PDF has a text layer (useful to skip OCR).

        Args:
            pdf_path: Path to PDF
            sample_pages: Number of pages to sample

        Returns:
            True if text layer detected
        """
        try:
            doc = fitz.open(pdf_path)
            total_pages = min(len(doc), sample_pages)

            has_text = False
            for page_num in range(total_pages):
                text = doc.load_page(page_num).get_text().strip()
                if len(text) > 50:  # More than 50 chars suggests real text
                    has_text = True
                    break

            doc.close()
            return has_text

        except Exception as e:
            self.logger.error(f"Error checking text layer: {e}")
            return False


class OfficeDocumentProcessor:
    """Process Office documents (docx, pptx, xlsx, csv) with table extraction"""

    def __init__(self, logger=None):
        self.logger = logger or get_logger(__name__)

    def _extract_images_from_shape(self, shape, file_stem: str, output_dir: Path, shape_idx: int) -> List[
        Dict[str, Any]]:
        """
        Extract images from Office document shapes (Word/PowerPoint).

        Args:
            shape: Shape object from docx or pptx
            file_stem: Base filename without extension
            output_dir: Directory to save extracted images
            shape_idx: Index of the shape for unique naming

        Returns:
            List of image metadata dicts
        """
        images = []

        try:
            # Check if shape contains an image
            if hasattr(shape, 'image'):
                image_bytes = shape.image.blob
                image_ext = shape.image.ext  # e.g., 'png', 'jpeg'

                # Create images subdirectory
                images_dir = output_dir / "extracted_images"
                images_dir.mkdir(exist_ok=True)

                # Save image
                image_filename = f"{file_stem}_shape{shape_idx}_image.{image_ext}"
                image_path = images_dir / image_filename

                with open(image_path, 'wb') as f:
                    f.write(image_bytes)

                images.append({
                    'filename': image_filename,
                    'path': str(image_path),
                    'size_bytes': len(image_bytes),
                    'format': image_ext,
                    'shape_index': shape_idx
                })

                self.logger.debug(f"Extracted image: {image_filename}")

        except AttributeError:
            # Shape doesn't contain an image
            pass
        except Exception as e:
            self.logger.warning(f"Failed to extract image from shape {shape_idx}: {e}")

        return images

    def extract_docx(self, file_path: Union[str, Path], output_dir: Optional[Path] = None) -> Dict[str, Any]:
        """
        Extract text and tables from Word with formatting preservation.

        Preserves:
        - Paragraph styles (Heading 1-9, Normal, etc.)
        - Bold, italic, underline
        - Font size
        - Text alignment
        - Table structure with cell styling
        """
        doc: Document = Document(file_path)
        content: Dict[str, Any] = {
            'paragraphs': [],
            'tables': [],
            'images': [],
            'metadata': {
                'total_paragraphs': 0,
                'total_tables': 0,
                'table_validation': [],
                'total_images': 0
            }
        }

        # Extract paragraphs with formatting
        for para in doc.paragraphs:
            if para.text.strip():
                para_data = {
                    'text': para.text,
                    'style': para.style.name if para.style else 'Normal',
                    'alignment': str(para.alignment) if para.alignment else 'LEFT',
                    'runs': []
                }

                # Extract run-level formatting
                for run in para.runs:
                    if run.text:
                        run_data = {
                            'text': run.text,
                            'bold': run.bold,
                            'italic': run.italic,
                            'underline': run.underline,
                            'font_size': run.font.size.pt if run.font.size else None
                        }
                        para_data['runs'].append(run_data)

                content['paragraphs'].append(para_data)

        # Extract tables with cell formatting
        for table_idx, table in enumerate(doc.tables):
            table_data: Dict[str, Any] = {
                'table_number': table_idx + 1,
                'rows': len(table.rows),
                'columns': len(table.columns),
                'data': [],
                'formatted_data': [],
                'dataframe': None,
                'validation': {
                    'has_header': False,
                    'empty_cells': 0,
                    'merged_cells': 0,
                    'data_quality_score': 0.0
                }
            }

            # Extract with formatting
            total_cells: int = 0
            empty_cells: int = 0

            for row_idx, row in enumerate(table.rows):
                row_data: List[str] = []
                formatted_row: List[Dict] = []

                for cell in row.cells:
                    cell_text: str = cell.text.strip()
                    total_cells += 1
                    if not cell_text:
                        empty_cells += 1

                    row_data.append(cell_text)

                    # Extract cell formatting
                    cell_format = {
                        'text': cell_text,
                        'bold': False,
                        'italic': False,
                        'alignment': 'LEFT',
                        'is_header': False
                    }

                    # Check first paragraph in cell
                    if cell.paragraphs:
                        para = cell.paragraphs[0]
                        if para.runs:
                            # Check if any run is bold (typical for headers)
                            cell_format['bold'] = any(r.bold for r in para.runs if r.bold)
                            cell_format['italic'] = any(r.italic for r in para.runs if r.italic)

                        cell_format['alignment'] = str(para.alignment) if para.alignment else 'LEFT'

                        # Headers typically use Heading styles or are bold
                        if row_idx == 0 or cell_format['bold']:
                            cell_format['is_header'] = True

                    formatted_row.append(cell_format)

                table_data['data'].append(row_data)
                table_data['formatted_data'].append(formatted_row)

            # Validation metrics
            table_data['validation']['empty_cells'] = empty_cells
            table_data['validation']['data_quality_score'] = (
                (total_cells - empty_cells) / total_cells if total_cells > 0 else 0
            )

            # Detect headers based on formatting
            if table_data['formatted_data']:
                first_row = table_data['formatted_data'][0]
                header_score = sum(1 for cell in first_row if cell['is_header']) / len(first_row)
                table_data['validation']['has_header'] = header_score > 0.6

            # Create DataFrame
            try:
                if len(table_data['data']) > 1:
                    df: pd.DataFrame
                    if table_data['validation']['has_header']:
                        df = pd.DataFrame(table_data['data'][1:], columns=table_data['data'][0])
                    else:
                        df = pd.DataFrame(table_data['data'])

                    df = df.dropna(how='all').dropna(axis=1, how='all')

                    if not df.empty:
                        table_data['dataframe'] = df.to_dict('records')
                        table_data['dataframe_shape'] = df.shape
                        table_data['columns'] = df.columns.tolist()
                    else:
                        table_data['dataframe'] = []
                        table_data['dataframe_shape'] = (0, 0)
                        table_data['columns'] = []
                else:
                    table_data['dataframe'] = []
                    table_data['dataframe_shape'] = (0, 0)
                    table_data['columns'] = []

            except Exception as e:
                self.logger.warning(f"Could not convert table {table_idx + 1} to DataFrame: {e}",
                                    exc_info=False, extra={'custom_fields': {'table_idx': table_idx}})
                table_data['dataframe'] = []
                table_data['dataframe_shape'] = (0, 0)
                table_data['columns'] = []

            content['tables'].append(table_data)
            content['metadata']['table_validation'].append(table_data['validation'])

        if output_dir:
            file_stem = Path(file_path).stem
            shape_idx = 0

            # Extract from document body
            for rel in doc.part.rels.values():
                if "image" in rel.target_ref:
                    try:
                        image_bytes = rel.target_part.blob
                        image_ext = rel.target_part.partname.split('.')[-1]

                        images_dir = output_dir / "extracted_images"
                        images_dir.mkdir(exist_ok=True)

                        image_filename = f"{file_stem}_image{shape_idx}.{image_ext}"
                        image_path = images_dir / image_filename

                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)

                        content['images'].append({
                            'filename': image_filename,
                            'path': str(image_path),
                            'size_bytes': len(image_bytes),
                            'format': image_ext,
                            'index': shape_idx
                        })

                        shape_idx += 1

                    except Exception as e:
                        self.logger.warning(f"Failed to extract image {shape_idx}: {e}")

        content['metadata']['total_paragraphs'] = len(content['paragraphs'])
        content['metadata']['total_tables'] = len(content['tables'])
        content['metadata']['total_images'] = len(content['images'])

        return content

    def extract_pptx(self, file_path: Union[str, Path], output_dir: Optional[Path] = None) -> Dict[str, Any]:
        """
        Extract text and tables from PowerPoint presentations

        Args:
            file_path: Path to the .pptx file

        Returns:
            Dictionary containing slides with text, tables, and metadata
        """

        prs: Presentation = Presentation(file_path)
        content: Dict[str, Any] = {
            'slides': [],
            'metadata': {
                'total_slides': len(prs.slides),
                'total_tables': 0,
                'total_text_boxes': 0,
                'total_images': 0
            }
        }

        file_stem = Path(file_path).stem if output_dir else None

        for slide_idx, slide in enumerate(prs.slides):
            slide_content: Dict[str, Any] = {
                'slide_number': slide_idx + 1,
                'title': '',
                'title_formatted': None,
                'text_boxes': [],
                'tables': [],
                'images': []
            }

            # Extract title with formatting
            if slide.shapes.title:
                title_text = slide.shapes.title.text
                slide_content['title'] = title_text

                # Extract title formatting
                if slide.shapes.title.text_frame:
                    title_formatted = {
                        'text': title_text,
                        'runs': []
                    }

                    for paragraph in slide.shapes.title.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if run.text:
                                run_data = {
                                    'text': run.text,
                                    'bold': run.font.bold if run.font.bold is not None else False,
                                    'italic': run.font.italic if run.font.italic is not None else False,
                                    'underline': run.font.underline if run.font.underline is not None else False,
                                    'font_size': run.font.size.pt if run.font.size else None
                                }
                                title_formatted['runs'].append(run_data)

                    slide_content['title_formatted'] = title_formatted

            # Extract text boxes and shapes with formatting
            for shape_idx, shape in enumerate(slide.shapes):
                # Skip title (already processed)
                if shape == slide.shapes.title:
                    continue

                if output_dir and hasattr(shape, 'image'):
                    try:
                        image_bytes = shape.image.blob
                        image_ext = shape.image.ext

                        images_dir = output_dir / "extracted_images"
                        images_dir.mkdir(exist_ok=True)

                        image_filename = f"{file_stem}_slide{slide_idx + 1}_shape{shape_idx}.{image_ext}"
                        image_path = images_dir / image_filename

                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)

                        slide_content['images'].append({
                            'filename': image_filename,
                            'path': str(image_path),
                            'size_bytes': len(image_bytes),
                            'format': image_ext,
                            'slide': slide_idx + 1,
                            'shape_index': shape_idx
                        })

                        content['metadata']['total_images'] += 1

                    except Exception as e:
                        self.logger.warning(
                            f"Failed to extract image from slide {slide_idx + 1}, shape {shape_idx}: {e}")

                # Text extraction with formatting
                if hasattr(shape, "text_frame") and shape.text_frame:
                    text_box = {
                        'text': shape.text,
                        'paragraphs': []
                    }

                    for paragraph in shape.text_frame.paragraphs:
                        para_data = {
                            'text': paragraph.text,
                            'alignment': str(paragraph.alignment) if paragraph.alignment else 'LEFT',
                            'level': paragraph.level,
                            'runs': []
                        }

                        for run in paragraph.runs:
                            if run.text:
                                run_data = {
                                    'text': run.text,
                                    'bold': run.font.bold if run.font.bold is not None else False,
                                    'italic': run.font.italic if run.font.italic is not None else False,
                                    'underline': run.font.underline if run.font.underline is not None else False,
                                    'font_size': run.font.size.pt if run.font.size else None
                                }
                                para_data['runs'].append(run_data)

                        text_box['paragraphs'].append(para_data)

                    if text_box['text'].strip():
                        slide_content['text_boxes'].append(text_box)
                        content['metadata']['total_text_boxes'] += 1

                # Table extraction (keep existing logic)
                if shape.has_table:
                    table = shape.table
                    table_data: Dict[str, Any] = {
                        'table_number': len(slide_content['tables']) + 1,
                        'rows': len(table.rows),
                        'columns': len(table.columns),
                        'data': [],
                        'formatted_data': [],
                        'dataframe': None
                    }

                    # Extract table with cell formatting
                    for row_idx, row in enumerate(table.rows):
                        row_data: List[str] = []
                        formatted_row: List[Dict] = []

                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            row_data.append(cell_text)

                            # Extract cell formatting
                            cell_format = {
                                'text': cell_text,
                                'bold': False,
                                'italic': False,
                                'is_header': False
                            }

                            # Check text frame in cell
                            if cell.text_frame and cell.text_frame.paragraphs:
                                for para in cell.text_frame.paragraphs:
                                    if para.runs:
                                        # Check if any run is bold
                                        cell_format['bold'] = any(
                                            r.font.bold for r in para.runs
                                            if r.font.bold is not None and r.font.bold
                                        )
                                        cell_format['italic'] = any(
                                            r.font.italic for r in para.runs
                                            if r.font.italic is not None and r.font.italic
                                        )

                            # First row or bold cells are likely headers
                            if row_idx == 0 or cell_format['bold']:
                                cell_format['is_header'] = True

                            formatted_row.append(cell_format)

                        table_data['data'].append(row_data)
                        table_data['formatted_data'].append(formatted_row)

                    # Convert to DataFrame (existing logic)
                    try:
                        df: pd.DataFrame
                        if len(table_data['data']) > 1:
                            # Check if first row has headers
                            first_row_formatted = table_data['formatted_data'][0]
                            header_score = sum(1 for cell in first_row_formatted if cell['is_header']) / len(
                                first_row_formatted)

                            if header_score > 0.6:
                                df = pd.DataFrame(table_data['data'][1:], columns=table_data['data'][0])
                            else:
                                df = pd.DataFrame(table_data['data'])

                            table_data['dataframe'] = df.to_dict('records')
                            table_data['dataframe_shape'] = df.shape
                        else:
                            df = pd.DataFrame(table_data['data'])
                            table_data['dataframe'] = df.to_dict('records')
                            table_data['dataframe_shape'] = df.shape
                    except Exception as e:
                        self.logger.warning(f"Could not convert PPTX table to DataFrame: {e}",
                                            exc_info=False, extra={'custom_fields': {'shape_type': 'table'}})

                    slide_content['tables'].append(table_data)
                    content['metadata']['total_tables'] += 1

                if output_dir and hasattr(shape, 'image'):
                    try:
                        image_bytes = shape.image.blob
                        image_ext = shape.image.ext

                        images_dir = output_dir / "extracted_images"
                        images_dir.mkdir(exist_ok=True)

                        image_filename = f"{file_stem}_slide{slide_idx + 1}_shape{shape_idx}.{image_ext}"
                        image_path = images_dir / image_filename

                        with open(image_path, 'wb') as f:
                            f.write(image_bytes)

                        slide_content['images'].append({
                            'filename': image_filename,
                            'path': str(image_path),
                            'size_bytes': len(image_bytes),
                            'format': image_ext,
                            'slide': slide_idx + 1,
                            'shape_index': shape_idx
                        })

                        content['metadata']['total_images'] += 1

                    except Exception as e:
                        self.logger.warning(
                            f"Failed to extract image from slide {slide_idx + 1}, shape {shape_idx}: {e}")

            content['slides'].append(slide_content)

        return content

    def extract_xlsx(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Excel extraction with sheet analysis

        Args:
            file_path: Path to the .xlsx file

        Returns:
            Dictionary containing sheets with data and metadata
        """
        wb: openpyxl.Workbook = openpyxl.load_workbook(file_path, data_only=True)
        content: Dict[str, Any] = {
            'sheets': [],
            'metadata': {
                'total_sheets': len(wb.sheetnames),
                'sheet_names': wb.sheetnames,
                'sheet_analysis': []
            }
        }

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_data: Dict[str, Any] = {
                'sheet_name': sheet_name,
                'used_range': f"{ws.dimensions}",
                'data': [],
                'dataframe': None,
                'analysis': {
                    'has_data': False,
                    'data_density': 0.0,
                    'likely_headers': False
                }
            }

            # Extract all rows
            all_rows: List[List[Any]] = []
            for row in ws.iter_rows(values_only=True):
                all_rows.append(list(row))

            sheet_data['data'] = all_rows

            # Analyze sheet content
            if all_rows:
                # Check if sheet has meaningful data
                non_empty_cells: int = sum(
                    1 for row in all_rows for cell in row
                    if cell is not None and str(cell).strip()
                )
                total_cells: int = sum(len(row) for row in all_rows)
                sheet_data['analysis']['data_density'] = (
                    non_empty_cells / total_cells if total_cells > 0 else 0
                )
                sheet_data['analysis']['has_data'] = sheet_data['analysis']['data_density'] > 0.1

                try:
                    if len(all_rows) > 1:
                        # Better header detection
                        first_row: List[Any] = all_rows[0]
                        header_indicators: int = sum(
                            1 for cell in first_row
                            if cell and isinstance(cell, str) and
                            not str(cell).replace('.', '').replace(',', '').replace('-', '').isdigit()
                        )

                        df: pd.DataFrame
                        if header_indicators > len(first_row) * 0.5:
                            sheet_data['analysis']['likely_headers'] = True
                            df = pd.DataFrame(all_rows[1:], columns=first_row)
                        else:
                            df = pd.DataFrame(all_rows)
                    else:
                        df = pd.DataFrame(all_rows)

                    # Clean DataFrame
                    df = df.dropna(how='all').dropna(axis=1, how='all')

                    if not df.empty:
                        sheet_data['dataframe'] = df.to_dict('records')
                        sheet_data['dataframe_shape'] = df.shape
                        sheet_data['columns'] = df.columns.tolist()
                    else:
                        sheet_data['dataframe'] = []
                        sheet_data['dataframe_shape'] = (0, 0)
                        sheet_data['columns'] = []

                except Exception as e:
                    self.logger.warning(f"Could not convert Excel sheet to DataFrame: {e}",
                                        exc_info=False, extra={'custom_fields': {'sheet_name': sheet_name}})
                    sheet_data['dataframe'] = []
                    sheet_data['dataframe_shape'] = (0, 0)
                    sheet_data['columns'] = []

            content['sheets'].append(sheet_data)
            content['metadata']['sheet_analysis'].append(sheet_data['analysis'])

        return content

    def extract_csv(self, file_path: Union[str, Path]) -> Dict[str, Any]:
        """
        Extract data from CSV files with encoding detection

        Args:
            file_path: Path to the .csv file

        Returns:
            Dictionary containing CSV data and metadata
        """

        content: Dict[str, Any] = {
            'filename': Path(file_path).name,
            'data': [],
            'dataframe': None,
            'metadata': {
                'encoding': 'unknown',
                'delimiter': 'unknown',
                'data_quality': 0.0
            }
        }

        encodings: List[str] = ['utf-8', 'latin-1', 'cp1252', 'iso-8859-1']
        last_error = None

        for encoding in encodings:
            try:
                df: pd.DataFrame = pd.read_csv(file_path, encoding=encoding, on_bad_lines='skip')
                content['metadata']['encoding'] = encoding

                # Detect delimiter
                with open(file_path, 'r', encoding=encoding) as f:
                    sample: str = f.read(1024)
                    delimiters: List[str] = [',', ';', '\t', '|']
                    delimiter_counts: Dict[str, int] = {
                        delim: sample.count(delim) for delim in delimiters
                    }
                    content['metadata']['delimiter'] = max(delimiter_counts, key=delimiter_counts.get)

                # Calculate data quality
                total_cells: int = df.size
                non_null_cells: int = df.count().sum()
                content['metadata']['data_quality'] = (
                    non_null_cells / total_cells if total_cells > 0 else 0
                )

                content['dataframe'] = df.to_dict('records')
                content['dataframe_shape'] = df.shape
                content['columns'] = df.columns.tolist()
                content['metadata']['rows'] = len(df)
                content['metadata']['columns'] = len(df.columns)
                content['data'] = [df.columns.tolist()] + df.values.tolist()

                self.logger.info(f"CSV parsed successfully with {encoding} encoding")
                break

            except Exception as e:
                last_error = e
                self.logger.debug(f"Failed to parse CSV with {encoding}: {e}")
                continue

        if content['metadata']['encoding'] == 'unknown' and last_error:
            self.logger.error(f"Could not parse CSV with any encoding: {last_error}")
            content['metadata']['parse_error'] = str(last_error)

        return content


class EmailProcessor:
    """Process .msg files with existing attachments in folder structure"""

    def __init__(self, subfolder_path: Union[str, Path]) -> None:
        """
        Initialize EmailProcessor

        Args:
            subfolder_path: Path to the folder containing the .msg file
        """
        self.subfolder_path: Path = Path(subfolder_path)
        self.msg_file: Optional[Path] = None
        self.msg_data: Optional[Dict[str, Any]] = None
        self.logger = get_logger(__name__)

    def find_msg_file(self) -> Optional[Path]:
        """
        Find the .msg file in the subfolder

        Returns:
            Path to .msg file if found, None otherwise
        """
        for p in sorted(self.subfolder_path.iterdir(), key=lambda x: x.name.lower()):
            if p.is_file() and p.suffix.lower() == ".msg":
                self.msg_file = p
                return p
        return None

    def _safe_str(self, value: Any) -> Optional[str]:
        """
        Safely convert any value to string, handling bytes

        Args:
            value: Value to convert

        Returns:
            String representation or None
        """
        if value is None:
            return None
        if isinstance(value, bytes):
            try:
                return value.decode('utf-8', errors='ignore')
            except:
                return str(value)
        return str(value)

    def extract_email_data(self) -> Optional[Dict[str, Any]]:
        """
        Extract structured data from .msg file with detailed error handling.

        Returns:
            Dictionary containing email metadata and content, or None on error
        """

        if not self.msg_file:
            self.find_msg_file()

        if not self.msg_file:
            self.logger.debug("No .msg file found in folder")
            return None

        try:
            msg = extract_msg.Message(str(self.msg_file))
            self.logger.debug(f"Message opened successfully: {self.msg_file.name}")
        except FileNotFoundError as e:
            self.logger.error(f"MSG file not found: {self.msg_file}", exc_info=False)
            return None

        except PermissionError as e:
            self.logger.error(f"Permission denied reading MSG: {self.msg_file}", exc_info=False)
            return None

        except Exception as e:
            self.logger.error(f"Failed to extract message: {self.msg_file}", exc_info=True)
            return None

        try:
            self.msg_data = {
                'metadata': {
                    'subject': self._safe_str(msg.subject) if msg.subject else '',
                    'sender': self._safe_str(msg.sender) if msg.sender else '',
                    'sender_email': self._safe_str(getattr(msg, 'senderEmail', None)),
                    'date': self._safe_str(msg.date) if msg.date else None,
                    'received_time': self._safe_str(getattr(msg, 'receivedTime', None))
                    if hasattr(msg, 'receivedTime') else None,
                    'message_id': self._safe_str(msg.messageId) if msg.messageId else '',
                    'importance': self._safe_str(msg.importance) if msg.importance else '',
                },
                'recipients': {
                    'to': self._format_recipients(msg.to),
                    'cc': self._format_recipients(msg.cc),
                    'bcc': self._format_recipients(msg.bcc)
                },
                'body': {
                    'plain_text': self._safe_str(msg.body) if msg.body else '',
                    'html': self._safe_str(msg.htmlBody) if msg.htmlBody else '',
                },
                'attachments_in_msg': self._get_msg_attachments(msg),
                'folder_path': str(self.subfolder_path),
                'msg_filename': self.msg_file.name
            }

            self.logger.info(f"Email extracted: {self.msg_data['metadata']['subject'][:80]}")
            msg.close()
            return self.msg_data

        except AttributeError as e:
            self.logger.error("Missing expected attribute in message object", exc_info=True)
            try:
                msg.close()
            except:
                pass
            return None

        except Exception as e:
            self.logger.error("Error extracting email data fields", exc_info=True)
            try:
                msg.close()
            except:
                pass
            return None

    def _format_recipients(self, recipients: Any) -> List[Dict[str, str]]:
        """
        Format recipient list - handles both string and object formats

        Args:
            recipients: Recipients in various formats

        Returns:
            List of dictionaries with name and email
        """
        if not recipients:
            return []

        formatted: List[Dict[str, str]] = []

        if isinstance(recipients, str):
            return [{'name': self._safe_str(recipients), 'email': self._safe_str(recipients)}]

        if isinstance(recipients, list):
            for r in recipients:
                if isinstance(r, str):
                    formatted.append({'name': self._safe_str(r), 'email': self._safe_str(r)})
                elif hasattr(r, 'name') and hasattr(r, 'email'):
                    formatted.append({
                        'name': self._safe_str(r.name) if r.name else '',
                        'email': self._safe_str(r.email) if r.email else ''
                    })
                else:
                    formatted.append({'name': self._safe_str(r), 'email': self._safe_str(r)})

        return formatted

    def _get_msg_attachments(self, msg: Any) -> List[Dict[str, Any]]:
        """
        Get attachment info from .msg file - WITHOUT binary data

        Args:
            msg: Message object from extract_msg

        Returns:
            List of attachment metadata dictionaries
        """
        attachments: List[Dict[str, Any]] = []
        try:
            for att in msg.attachments:
                filename: str = att.longFilename or att.shortFilename or 'unknown'
                att_info: Dict[str, Any] = {
                    'filename': self._safe_str(filename),
                    'size': len(att.data) if att.data else 0,
                }
                attachments.append(att_info)
        except Exception as e:
            self.logger.warning(f"Warning: Could not read attachments from msg: {e}")

        return attachments

    def list_folder_files(self) -> List[Dict[str, Any]]:
        """
        List all files in the subfolder (saved attachments)

        Returns:
            List of file metadata dictionaries
        """
        all_files: List[Dict[str, Any]] = []
        for file in self.subfolder_path.iterdir():
            if file.is_file() and not file.name.endswith('.msg'):
                all_files.append({
                    'filename': file.name,
                    'path': str(file),
                    'size': file.stat().st_size
                })
        return all_files

    def get_complete_data(self) -> Optional[Dict[str, Any]]:
        """
        Get email data + list of existing attachment files

        Returns:
            Complete email data dictionary or None on error
        """
        email_data: Optional[Dict[str, Any]] = self.extract_email_data()
        if not email_data:
            return None

        email_data['saved_attachments'] = self.list_folder_files()
        return email_data

    def process_attachments_with_ocr(
            self,
            email_data: Dict[str, Any],
            ocr_engine: DetectAndRecognize,
            save_visuals: bool = True
    ) -> list[dict[str, Any]] | None:
        """
        Process attachments with OCR, office documents, and direct text extraction.

        Args:
            email_data: Dictionary containing email metadata and attachments
            ocr_engine: Instance of Detect_and_Recognize class for OCR operations
            save_visuals: If True, saves annotated images and JSON results to:
                         <email_folder>/text_detected_and_recognized/<TIMESTAMP>/

        Returns:
            List of dictionaries containing extracted text and metadata for each attachment
        """
        attachment_texts: List[Dict[str, Any]] = []
        office_processor = OfficeDocumentProcessor()

        # If saving visuals, create timestamped output root in the email folder
        out_root: Optional[Path] = None
        if save_visuals:
            ts: str = datetime.now().strftime("%Y-%m-%d_%H-%M-%S")
            out_root = Path(self.subfolder_path) / "text_detected_and_recognized" / ts
            out_root.mkdir(parents=True, exist_ok=True)
            tqdm.write(f"[OCR] Saving visuals to: {out_root}")
            self.logger.info(f"[OCR] Saving visuals to: {out_root}")

        attachments: List[Dict[str, Any]] = email_data.get('saved_attachments', [])
        resource_limits = ResourceLimits()

        # Main attachment processing loop with progress bar
        for attachment in tqdm(attachments, desc="Processing attachments", unit="file", leave=False):
            file_path: Path = Path(attachment['path'])
            valid, errors = resource_limits.validate_file_size(file_path)
            if not valid:
                self.logger.warning(f"Skipping file due to resource limits: {errors}")
                continue
            ext: str = file_path.suffix.lower()
            base_name: str = file_path.stem

            if ext in ['.json']:
                continue

            start_time: float = time.time()

            # OFFICE DOCUMENTS - Process first before OCR
            x = SUPPORTED_EXTENSIONS['office']
            # print(f'Office extension; {ext} index; {x.index(ext)}')
            file_output_dir = out_root / base_name
            file_output_dir.mkdir(exist_ok=True)
            if ext in x[x.index('.docx')]:
                try:

                    content: Dict[str, Any] = office_processor.extract_docx(file_path, output_dir=file_output_dir)

                    # Extract all text for LLM context with style hints
                    text_parts: List[str] = []

                    # Extract paragraph text with heading markers
                    for para in content['paragraphs']:
                        if isinstance(para, dict):
                            text = para['text']
                            style = para.get('style', 'Normal')

                            # Add heading markers for important sections
                            if 'Heading' in style:
                                text = f"\n## {text} ##\n"

                            text_parts.append(text)
                        else:
                            text_parts.append(str(para))

                    # Extract table text
                    for table in content['tables']:
                        # Check if table has header row
                        has_header = table.get('validation', {}).get('has_header', False)

                        table_text: str = f"\n[TABLE {table['table_number']} - {table['rows']}x{table['columns']}]\n"

                        for row_idx, row in enumerate(table['data']):
                            row_text = " | ".join(str(cell) for cell in row)

                            # Mark header rows
                            if row_idx == 0 and has_header:
                                row_text = f"HEADER: {row_text}"

                            table_text += row_text + "\n"

                        text_parts.append(table_text)

                    extracted_images = content.get('images', [])
                    if extracted_images:
                        self.logger.info(f"Found {len(extracted_images)} images in {file_path.name}, running OCR...")

                        for img_info in extracted_images:
                            img_path = Path(img_info['path'])

                            try:
                                # Run OCR on extracted image
                                img_ocr_results = ocr_engine.ocr_with_detection_and_recognition(str(img_path))

                                # Extract text from OCR results
                                img_text_parts = []
                                for res in img_ocr_results:
                                    result_data = res if isinstance(res, dict) else (getattr(res, "json", None) or res)
                                    rec_texts = result_data.get('rec_texts', []) or []
                                    rec_scores = result_data.get('rec_scores', []) or []

                                    if rec_scores and len(rec_scores) == len(rec_texts):
                                        for t, s in zip(rec_texts, rec_scores):
                                            try:
                                                if float(s) >= OCR_CONFIG.get('confidence_threshold', 0.7) and str(
                                                        t).strip():
                                                    img_text_parts.append(str(t).strip())
                                            except Exception:
                                                continue

                                img_ocr_text = ' '.join(img_text_parts)

                                if img_ocr_text:
                                    text_parts.append(f"\n[IMAGE OCR - {img_info['filename']}]\n{img_ocr_text}\n")
                                    self.logger.debug(
                                        f"OCR extracted {len(img_text_parts)} text blocks from {img_info['filename']}")

                            except Exception as e:
                                self.logger.warning(f"OCR failed on extracted image {img_info['filename']}: {e}")

                    full_text: str = "\n".join(text_parts)
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': full_text,
                        'method': 'docx_extraction_with_image_ocr',  # Updated method name
                        'structured_content': content,
                        'images_processed': len(extracted_images),
                        'time': time.time() - start_time
                    })

                    # Save structured data
                    if save_visuals and out_root is not None:
                        json_path: Path = out_root / f"{base_name}_docx_extracted.json"
                        with open(json_path, 'w', encoding='utf-8') as jf:
                            json.dump(content, jf, indent=2, ensure_ascii=False)

                except Exception as e:
                    tqdm.write(f"[DOCX] Error processing {file_path}: {e}")
                    self.logger.error(f"[DOCX] Error processing {file_path}: {e}", exc_info=True)
                    # traceback.print_exc()

            elif ext in x[x.index('.pptx')]:
                try:
                    content: Dict[str, Any] = office_processor.extract_pptx(file_path, output_dir=file_output_dir)

                    # Extract all text for LLM context
                    text_parts: List[str] = []
                    for slide in content['slides']:
                        slide_text: str = f"\n[SLIDE {slide['slide_number']}: {slide['title']}]\n"

                        # Add text boxes
                        for text_box in slide.get('text_boxes', []):
                            slide_text += text_box['text'] + "\n"

                        # Add tables
                        for table in slide['tables']:
                            table_text: str = f"\n[TABLE {table['table_number']} - {table['rows']}x{table['columns']}]\n"
                            for row in table['data']:
                                table_text += " | ".join(str(cell) for cell in row) + "\n"
                            slide_text += table_text

                        text_parts.append(slide_text)

                        slide_images = slide.get('images', [])
                        if slide_images:
                            for img_info in slide_images:
                                img_path = Path(img_info['path'])

                                try:
                                    img_ocr_results = ocr_engine.ocr_with_detection_and_recognition(str(img_path))
                                    img_text_parts = []
                                    for res in img_ocr_results:
                                        result_data = res if isinstance(res, dict) else (
                                                getattr(res, "json", None) or res)
                                        rec_texts = result_data.get('rec_texts', []) or []
                                        rec_scores = result_data.get('rec_scores', []) or []

                                        if rec_scores and len(rec_scores) == len(rec_texts):
                                            for t, s in zip(rec_texts, rec_scores):
                                                try:
                                                    if float(s) >= OCR_CONFIG.get('confidence_threshold', 0.7) and str(
                                                            t).strip():
                                                        img_text_parts.append(str(t).strip())
                                                except Exception:
                                                    continue

                                    img_ocr_text = ' '.join(img_text_parts)
                                    if img_ocr_text:
                                        text_parts.append(f"\n[IMAGE OCR - {img_info['filename']}]\n{img_ocr_text}\n")
                                        self.logger.debug(
                                            f"OCR extracted {len(img_text_parts)} text blocks from {img_info['filename']}")

                                except Exception as e:
                                    self.logger.warning(f"OCR failed on slide image: {e}")

                    full_text: str = "\n".join(text_parts)
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': full_text,
                        'method': 'pptx_extraction_with_image_ocr',
                        'structured_content': content,
                        'images_processed': sum(len(s.get('images', [])) for s in content['slides']),
                        'time': time.time() - start_time
                    })

                    # Save structured data with formatting
                    if save_visuals and out_root is not None:
                        json_path: Path = out_root / f"{base_name}_pptx_extracted.json"
                        with open(json_path, 'w', encoding='utf-8') as jf:
                            json.dump(content, jf, indent=2, ensure_ascii=False)

                except Exception as e:
                    tqdm.write(f"[PPTX] Error processing {file_path}: {e}")
                    self.logger.error(f"[PPTX] Error processing {file_path}: {e}", exc_info=True)

            elif ext in x[x.index('.xlsx')]:
                try:
                    content: Dict[str, Any] = office_processor.extract_xlsx(file_path)

                    # Extract all sheets for LLM context
                    text_parts: List[str] = []
                    for sheet in content['sheets']:
                        sheet_text: str = f"\n[SHEET: {sheet['sheet_name']} - {sheet['used_range']}]\n"
                        if sheet['dataframe']:
                            # Create table representation
                            for row in sheet['data'][:100]:  # Limit to first 100 rows
                                sheet_text += " | ".join(str(cell) if cell is not None else "" for cell in row) + "\n"
                        text_parts.append(sheet_text)

                    full_text: str = "\n".join(text_parts)
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': full_text,
                        'method': 'xlsx_extraction',
                        'structured_content': content,
                        'time': time.time() - start_time
                    })

                    # Save structured data
                    if save_visuals and out_root is not None:
                        json_path: Path = out_root / f"{base_name}_xlsx_extracted.json"
                        with open(json_path, 'w', encoding='utf-8') as jf:
                            json.dump(content, jf, indent=2, ensure_ascii=False)

                except Exception as e:
                    tqdm.write(f"[XLSX] Error processing {file_path}: {e}")
                    self.logger.error(f"[XLSX] Error processing {file_path}: {e}", exc_info=True)

            elif ext in x[x.index('.csv')]:
                try:
                    content: Dict[str, Any] = office_processor.extract_csv(file_path)

                    # Create text representation
                    text_parts: List[str] = [
                        f"[CSV FILE - {content['metadata'].get('rows', 0)} rows x {content['metadata'].get('columns', 0)} columns]\n"
                    ]
                    for row in content['data'][:100]:  # Limit to first 100 rows
                        text_parts.append(" | ".join(str(cell) if cell is not None else "" for cell in row))

                    full_text: str = "\n".join(text_parts)
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': full_text,
                        'method': 'csv_extraction',
                        'structured_content': content,
                        'time': time.time() - start_time
                    })

                    # Save structured data
                    if save_visuals and out_root is not None:
                        json_path: Path = out_root / f"{base_name}_csv_extracted.json"
                        with open(json_path, 'w', encoding='utf-8') as jf:
                            json.dump(content, jf, indent=2, ensure_ascii=False)

                except Exception as e:
                    tqdm.write(f"[CSV] Error processing {file_path}: {e}")
                    self.logger.error(f"[CSV] Error processing {file_path}: {e}", exc_info=True)

            # PDF - OCR Processing
            elif ext in SUPPORTED_EXTENSIONS["documents"]:
                try:
                    pdf_start = time.time()

                    streaming_processor = StreamingPDFProcessor(
                        initial_dpi=PROCESSING_CONFIG.get('pdf_dpi', 300),
                        enable_memory_monitoring=True,
                        enable_adaptive_quality=True
                    )

                    # Open PDF once to analyze all pages
                    doc = fitz.open(file_path)
                    total_pages = len(doc)
                    tqdm.write(f"  [PDF] {base_name} - {total_pages} pages")

                    # Analyze each page for text content AND images
                    pages_analysis = []
                    for page_num in range(total_pages):
                        page = doc.load_page(page_num)
                        text = page.get_text().strip()

                        # Check for images on page
                        image_list = page.get_images(full=True)
                        has_images = len(image_list) > 0

                        # Determine processing strategy
                        has_text = len(text) > 50

                        pages_analysis.append({
                            'page_num': page_num,
                            'has_text': has_text,
                            'has_images': has_images,
                            'text_length': len(text),
                            'image_count': len(image_list),
                            'strategy': None  # Will determine below
                        })

                    doc.close()

                    # Determine strategy for each page
                    for page_info in pages_analysis:
                        if page_info['has_text'] and not page_info['has_images']:
                            # Pure text page - text extraction only
                            page_info['strategy'] = 'text_only'
                        elif not page_info['has_text'] and page_info['has_images']:
                            # Pure image page (scanned) - OCR only
                            page_info['strategy'] = 'ocr_only'
                        elif page_info['has_text'] and page_info['has_images']:
                            # Mixed content - BOTH text extraction AND OCR
                            page_info['strategy'] = 'hybrid'
                        else:
                            # Empty page
                            page_info['strategy'] = 'skip'

                    # Log strategy summary
                    strategy_counts = {}
                    for p in pages_analysis:
                        strat = p['strategy']
                        strategy_counts[strat] = strategy_counts.get(strat, 0) + 1

                    tqdm.write(f"  [PDF] Strategy: {strategy_counts}")

                    # Process each page according to its strategy
                    for page_info in pages_analysis:
                        page_num = page_info['page_num']
                        strategy = page_info['strategy']

                        if strategy == 'skip':
                            tqdm.write(f"    [PDF] Page {page_num + 1}: Empty, skipping")
                            self.logger.info(f'[PDF] Page {page_num + 1}: Empty, skipping')
                            continue

                        # TEXT EXTRACTION (for 'text_only' and 'hybrid')
                        if strategy in ['text_only', 'hybrid']:
                            page_text = streaming_processor.extract_page_text(file_path, page_num)
                            if page_text.strip():
                                attachment_texts.append({
                                    'file': attachment['filename'],
                                    'page': page_num + 1,
                                    'text': page_text,
                                    'method': 'pdf_text_layer',
                                    'source': 'text_extraction',
                                    'time': time.time() - pdf_start
                                })
                                tqdm.write(
                                    f"    [PDF] Page {page_num + 1}: Text extracted ({page_info['text_length']} chars)")
                                self.logger.info(
                                    f"[PDF] Page {page_num + 1}: Text extracted ({page_info['text_length']} chars)")

                        # OCR (for 'ocr_only' and 'hybrid')
                        if strategy in ['ocr_only', 'hybrid']:
                            reason = "scanned page" if strategy == 'ocr_only' else f"has {page_info['image_count']} images"
                            tqdm.write(f"    [PDF] Page {page_num + 1}: Running OCR ({reason})")
                            self.logger.info(f'[PDF] Page {page_num + 1}: Running OCR ({reason})')

                            page_start = time.time()

                            # Use streaming processor to get page as image
                            doc = fitz.open(file_path)
                            try:
                                processed, original = streaming_processor.process_page(doc, page_num)

                                # Save temp image
                                with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp:
                                    img_pil = Image.fromarray(original)
                                    img_pil.save(tmp.name)
                                    tmp_path = tmp.name

                                # OCR
                                ocr_start = time.time()
                                results = ocr_engine.ocr_with_detection_and_recognition(tmp_path)
                                ocr_time = time.time() - ocr_start

                                # Save visuals
                                if save_visuals and out_root is not None and results:
                                    try:
                                        if isinstance(results, list) and len(results) > 0:
                                            result_obj = results[0]
                                            page_dir = out_root / f"{base_name}_page{page_num + 1}_OCR"
                                            page_dir.mkdir(exist_ok=True)
                                            result_obj.save_to_img(str(page_dir))

                                            # Save JSON with strategy info
                                            json_path = page_dir / "ocr_result.json"
                                            result_obj.save_to_json(str(json_path))

                                            # Save strategy info separately
                                            strategy_path = page_dir / "page_info.json"
                                            with open(strategy_path, 'w', encoding='utf-8') as f:
                                                json.dump(page_info, f, indent=2)

                                    except Exception as e:
                                        tqdm.write(f"      [PDF] Warning: Could not save page {page_num + 1}: {e}")
                                        self.logger.error(f"[PDF] Warning: Could not save page {page_num + 1}: {e}",
                                                          exc_info=True)

                                # Extract text from OCR
                                text_parts = []
                                for res in results:
                                    result_data = res if isinstance(res, dict) else (getattr(res, "json", None) or res)
                                    rec_texts = result_data.get('rec_texts', []) or []
                                    rec_scores = result_data.get('rec_scores', []) or []

                                    if rec_scores and len(rec_scores) == len(rec_texts):
                                        for t, s in zip(rec_texts, rec_scores):
                                            try:
                                                if float(s) >= OCR_CONFIG.get('confidence_threshold', 0.7) and str(
                                                        t).strip():
                                                    text_parts.append(str(t).strip())
                                            except Exception:
                                                continue

                                ocr_text = ' '.join(text_parts).strip()

                                if ocr_text:
                                    attachment_texts.append({
                                        'file': attachment['filename'],
                                        'page': page_num + 1,
                                        'text': ocr_text,
                                        'method': 'ocr',
                                        'source': 'image_ocr' if strategy == 'ocr_only' else 'hybrid_ocr',
                                        'time': time.time() - page_start,
                                        'ocr_time': ocr_time,
                                        'image_count': page_info['image_count']
                                    })
                                    tqdm.write(
                                        f"      [PDF] Page {page_num + 1}: OCR completed ({ocr_time:.2f}s, {len(text_parts)} text blocks)")
                                    self.logger.info(
                                        f'[PDF] Page {page_num + 1}: OCR completed ({ocr_time:.2f}s, {len(text_parts)} text blocks)')
                                else:
                                    tqdm.write(f"      [PDF] Page {page_num + 1}: OCR found no text")
                                    self.logger.info(f'[PDF] Page {page_num + 1}: OCR found no text')

                                os.unlink(tmp_path)

                            finally:
                                doc.close()

                            # Memory check after each OCR page
                            if streaming_processor.memory_monitor:
                                streaming_processor.memory_monitor.check_and_cleanup()

                    tqdm.write(f"  [PDF] Total: {time.time() - pdf_start:.2f}s")
                    self.logger.info(f'[PDF] Total: {time.time() - pdf_start:.2f}s')

                except Exception as e:
                    # tqdm.write(f"[OCR] Error processing PDF {file_path}: {e}")
                    logger = create_logger_with_context(__name__, file=attachment['filename'])
                    logger.error(f"Error processing PDF: {e}", exc_info=True)

            # IMAGE FILES
            elif ext in SUPPORTED_EXTENSIONS["images"]:
                try:
                    img_start: float = time.time()

                    # OCR with timing
                    ocr_start: float = time.time()
                    results: List[Any] = ocr_engine.ocr_with_detection_and_recognition(file_path)
                    ocr_time: float = time.time() - ocr_start

                    # SAVE FIRST - before any text filtering
                    if save_visuals and out_root is not None and results:
                        save_start: float = time.time()
                        try:
                            # Handle results structure - should be a list
                            result_obj: Any
                            if isinstance(results, list) and len(results) > 0:
                                result_obj = results[0]
                            else:
                                result_obj = results

                            # Create file-specific subdirectory for this image's outputs
                            img_dir: Path = out_root / f"{base_name}_image"
                            img_dir.mkdir(exist_ok=True)

                            # save_to_img expects a DIRECTORY
                            result_obj.save_to_img(str(img_dir))

                            # Save JSON with specific filename
                            json_path: Path = img_dir / "ocr_result.json"
                            result_obj.save_to_json(str(json_path))

                            save_time: float = time.time() - save_start
                            tqdm.write(f"  [IMG] {base_name} - OCR: {ocr_time:.2f}s, Save: {save_time:.2f}s")
                        except Exception as e:
                            tqdm.write(f"  [IMG] Warning: Could not save visuals: {e}")
                            self.logger.error(f'[IMG] Warning: Could not save visuals: {e}', exc_info=True)
                            # traceback.print_exc()

                    # NOW extract text from result objects
                    text_parts: List[str] = []
                    for res in results:
                        # Access the JSON dict directly from the result object
                        result_data: Dict[str, Any]
                        if isinstance(res, dict):
                            result_data = res
                        else:
                            result_data = getattr(res, "json", None) or res

                        rec_texts: List[str] = result_data.get('rec_texts', []) or []
                        rec_scores: List[float] = result_data.get('rec_scores', []) or []

                        if rec_scores and len(rec_scores) == len(rec_texts):
                            for t, s in zip(rec_texts, rec_scores):
                                try:
                                    if float(s) >= 0.7 and str(t).strip():
                                        text_parts.append(str(t).strip())
                                except Exception:
                                    continue
                        else:
                            for t in rec_texts:
                                if isinstance(t, str) and len(t.strip()) >= 2:
                                    text_parts.append(t.strip())

                    full_text: str = ' '.join(text_parts).strip()

                    # Only add to attachment_texts if there's filtered text
                    # But we already saved visuals above regardless
                    if full_text:
                        attachment_texts.append({
                            'file': attachment['filename'],
                            'text': full_text,
                            'method': 'ocr',
                            'time': time.time() - img_start,
                            'ocr_time': ocr_time
                        })

                except Exception as e:
                    tqdm.write(f"[OCR] Error processing image {file_path}: {e}")
                    # traceback.print_exc()
                    self.logger.error(f'[OCR] Error processing image {file_path}: {e}', exc_info=True)

            # TEXT FILES - Direct extraction
            elif ext in SUPPORTED_EXTENSIONS["text"]:
                try:
                    with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text: str = f.read()
                    attachment_texts.append({
                        'file': attachment['filename'],
                        'text': text,
                        'method': 'direct',
                        'time': time.time() - start_time
                    })

                    # optionally save a JSON copy of the raw text to visuals folder
                    if save_visuals and out_root is not None:
                        json_name: str = f"{base_name}_raw_text.json"
                        json_path: Path = out_root / json_name
                        try:
                            with open(json_path, 'w', encoding='utf-8') as jf:
                                json.dump({"file": attachment['filename'], "text": text}, jf, indent=2,
                                          ensure_ascii=False)
                        except Exception as e:
                            tqdm.write(f"[OCR] Warning: failed to save raw text json {json_path}: {e}")
                            self.logger.error(f'[OCR] Warning: failed to save raw text json {json_path}: {e}',
                                              exc_info=True)

                except Exception as e:
                    tqdm.write(f"[TXT] Error processing {file_path}: {e}")
                    self.logger.error(f'[TXT] Error processing {file_path}: {e}', exc_info=True)


            # UNKNOWN - Try OCR as fallback
            else:
                try:
                    ocr_start: float = time.time()
                    results: List[Any] = ocr_engine.ocr_with_detection_and_recognition(file_path)
                    ocr_time: float = time.time() - ocr_start

                    # Extract text from results
                    text_parts: List[str] = []
                    for res in results:
                        result_data: Dict[str, Any]
                        if isinstance(res, dict):
                            result_data = res
                        else:
                            result_data = getattr(res, "json", None) or res

                        rec_texts: List[str] = result_data.get('rec_texts', []) or []
                        rec_scores: List[float] = result_data.get('rec_scores', []) or []

                        if rec_scores and len(rec_scores) == len(rec_texts):
                            for t, s in zip(rec_texts, rec_scores):
                                try:
                                    if float(s) >= 0.7 and str(t).strip():
                                        text_parts.append(str(t).strip())
                                except Exception:
                                    continue

                    text: str = ' '.join(text_parts).strip()

                    if text:
                        attachment_texts.append({
                            'file': attachment['filename'],
                            'text': text,
                            'method': 'ocr_fallback',
                            'time': time.time() - start_time,
                            'ocr_time': ocr_time
                        })

                    if save_visuals and out_root is not None and results:
                        try:
                            result_obj: Any = results[0] if isinstance(results, list) and len(results) > 0 else results
                            unknown_dir: Path = out_root / f"{base_name}_unknown"
                            unknown_dir.mkdir(exist_ok=True)
                            result_obj.save_to_img(str(unknown_dir))
                            json_path: Path = unknown_dir / "ocr_result.json"
                            result_obj.save_to_json(str(json_path))
                        except Exception as e:
                            tqdm.write(f"[OCR] Warning: failed to save unknown file type visuals: {e}")
                            self.logger.error(f'[OCR] Warning: failed to save unknown file type visuals: {e}',
                                              exc_info=True)

                except Exception as e:
                    tqdm.write(f"[OCR] Skipping unknown file type {file_path}: {e}")
                    self.logger.error(f'[OCR] Skipping unknown file type {file_path}: {e}', exc_info=True)

        return attachment_texts


class SmartTextTruncator:
    """
    Intelligently truncate text while preserving important sections.

    Features:
    - Identify and preserve key sections
    - Keep document structure
    - Maintain context around important fields
    """

    IMPORTANT_KEYWORDS = [
        'insured', 'cedant', 'broker', 'sum insured', 'total_sum_insured', 'tsi', 'rate', 'premium rate',
        'premium', 'period',
        'coverage', 'retention', 'deductible', 'claims experience',
        'excess', 'share', 'limit', 'country', 'location', 'risk'
    ]

    def __init__(self):
        self.logger = get_logger(__name__)

    def score_section(self, text: str) -> float:
        """
        Score a text section by importance.

        Args:
            text: Text section

        Returns:
            Importance score (0-1)
        """
        text_lower = text.lower()

        # Count important keywords
        keyword_count = sum(1 for kw in self.IMPORTANT_KEYWORDS if kw in text_lower)

        # Presence of numbers (likely financial data)
        has_numbers = bool(re.search(r'\d+[,.]?\d*', text))

        # Contains dates
        has_dates = bool(re.search(r'\d{1,2}[\/\-]\d{1,2}[\/\-]\d{2,4}', text))

        # Calculate score
        score = 0.0
        score += min(keyword_count * 0.2, 0.6)  # Max 0.6 from keywords
        score += 0.2 if has_numbers else 0.0
        score += 0.2 if has_dates else 0.0

        return min(score, 1.0)

    def smart_truncate(
            self,
            text: str,
            max_length: int,
            keep_start_percent: float = 0.3,
            keep_end_percent: float = 0.2
    ) -> str:
        """
        Intelligently truncate text preserving important sections.

        Args:
            text: Full text
            max_length: Maximum length
            keep_start_percent: Percentage of budget for start
            keep_end_percent: Percentage of budget for end

        Returns:
            Truncated text with preserved important sections
        """
        if len(text) <= max_length:
            return text

        # Split into paragraphs
        paragraphs = text.split('\n\n')

        # Score each paragraph
        scored_paragraphs = [
            (i, para, self.score_section(para))
            for i, para in enumerate(paragraphs)
        ]

        # Always keep start and end
        start_budget = int(max_length * keep_start_percent)
        end_budget = int(max_length * keep_end_percent)
        middle_budget = max_length - start_budget - end_budget

        # Build output
        result_parts = []
        current_length = 0

        # Add start
        start_text = text[:start_budget]
        result_parts.append(start_text)
        current_length += len(start_text)

        # Add important middle sections
        middle_paragraphs = sorted(
            scored_paragraphs[1:-1],
            key=lambda x: x[2],
            reverse=True
        )

        included_middle = []
        for idx, para, score in middle_paragraphs:
            if current_length + len(para) < max_length - end_budget:
                included_middle.append((idx, para))
                current_length += len(para)

        # Sort by original order
        included_middle.sort(key=lambda x: x[0])

        if included_middle:
            result_parts.append("\n\n[...MIDDLE SECTIONS...]\n\n")
            result_parts.extend([para for _, para in included_middle])

        # Add end
        result_parts.append("\n\n[...END SECTION...]\n\n")
        end_text = text[-end_budget:]
        result_parts.append(end_text)

        return ''.join(result_parts)


def create_llm_context(
        email_data: Dict[str, Any],
        attachment_texts: List[Dict[str, Any]]
) -> Tuple[Dict[str, Any], str]:
    """
    Build LLM context and fill the USER_PROMPT with the combined text.
    - Skips attachments with empty text.
    - Returns (context_dict, llm_prompt_string).

    Args:
        email_data: Dictionary containing email metadata and body
        attachment_texts: List of dictionaries with attachment text and metadata

    Returns:
        Tuple of (context dictionary, formatted LLM prompt string)
    """

    # Initialize truncator
    truncator = SmartTextTruncator()
    logger = get_logger(__name__)

    # Filter attachments
    kept_attachments = []
    for att in attachment_texts:
        txt = att.get('text', '')
        if isinstance(txt, str) and txt.strip():
            # Smart truncate long attachments
            if len(txt) > PROCESSING_CONFIG.get('max_attachment_length', 15000):
                txt = truncator.smart_truncate(
                    txt,
                    max_length=PROCESSING_CONFIG.get('max_attachment_length', 15000),
                    keep_start_percent=0.4,
                    keep_end_percent=0.2
                )
                att['text'] = txt
                att['truncated'] = True
            kept_attachments.append(att)

    # Build context
    context = {
        'email': {
            'subject': email_data['metadata'].get('subject', ''),
            'from': email_data['metadata'].get('sender', ''),
            'date': email_data['metadata'].get('date', ''),
            'body': email_data['body'].get('plain_text', '') if 'body' in email_data else ''
        },
        'attachments': kept_attachments,
        'total_attachments': len(kept_attachments),
        'processing_times': {
            'email': 0,
            'attachments': sum(a.get('time', 0) for a in kept_attachments)
        }
    }

    # Build combined text
    parts = []
    parts.append(f"Email Subject: {context['email']['subject']}")
    parts.append(f"From: {context['email']['from']}")
    parts.append(f"Date: {context['email']['date']}")

    # Smart truncate email body
    email_body = context['email']['body'].strip()
    if len(email_body) > 5000:
        email_body = truncator.smart_truncate(email_body, max_length=5000)
    parts.append(f"\nEmail Body:\n{email_body}")

    if kept_attachments:
        parts.append("\nAttachments:")
        for i, att in enumerate(kept_attachments, start=1):
            header = f"\n--- Attachment {i}: {att['file']}"
            if att.get('page'):
                header += f" (page {att['page']})"
            if att.get('truncated'):
                header += " [TRUNCATED]"
            header += " ---"
            parts.append(header + "\n" + att['text'])

    combined_text = "\n".join(parts).strip()

    # Extract pattern candidates
    try:
        patterns = ReinsurancePatterns()
        candidates = extract_with_patterns(combined_text, patterns)

        if candidates:
            logger.debug(f"Pattern extraction found {len(candidates)} candidates")
            hint_lines = ["\n[AUTO-EXTRACTED HINTS]:"]
            for k, v in candidates.items():
                val_str = str(v.value if hasattr(v, 'value') else v)[:100]
                hint_lines.append(f"{k}: {val_str}")
            combined_text += "\n" + "\n".join(hint_lines)
        else:
            logger.debug("Pattern extraction found no candidates")

    except Exception as e:
        logger.warning(f"Pattern extraction failed, continuing without hints: {e}",
                       exc_info=False)

    # Final truncation using smart truncator
    if len(combined_text) > PROCESSING_CONFIG.get('max_text_length', 80000):
        combined_text = truncator.smart_truncate(
            combined_text,
            max_length=PROCESSING_CONFIG.get('max_text_length', 80000)
        )

    user_prompt: str = combined_text

    return context, user_prompt
